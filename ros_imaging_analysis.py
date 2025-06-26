### ROS- Combined Plates
## takes as input a folder that includes an additional folder inside it with txt files
## there should be numeric features in these files & index columns
## must index columns are - cell_type, cell_id, compound, concentration
## follow the steps until a final normalized data frame ready for statistics will be generated.
## all output files will be saved in the main folder including all steps in the analysis, in addition a dashboard with the main results will be saved.
## all text files will be merged and analyzed.

#required libraries :
import os
import logging
import pandas as pd
import numpy as np
import os
import glob
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import itertools
import statsmodels.api as sm
from scipy.stats import boxcox, zscore
from sklearn.neighbors import LocalOutlierFactor
from sklearn.ensemble import RandomForestRegressor
import kaleido
from openpyxl import Workbook
from openpyxl.utils.dataframe import dataframe_to_rows
from scipy.stats import zscore
from sklearn.impute import SimpleImputer
import scipy
from sklearn.decomposition import PCA
import matplotlib.patches as mpatches
from colorama import Fore, Style, init
init(autoreset=True)
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
def process_folders(main_folder):
    try:
        #check main folder
        if not os.path.isdir(main_folder):
            print(f"Error: The folder '{main_folder}' does not exist.")
            return
        #run through subfolders in main folder
        for subfolder_name in os.listdir(main_folder):
            subfolder_path = os.path.join(main_folder, subfolder_name)
            if os.path.isdir(subfolder_path):
                print(f"\nProcessing subfolder: {subfolder_name}")
                #user to choose value for the 'PC' column- plate condition- must.
                pc_value = input("Enter the value for the 'PC' column(e.g- plate condition ), must add (numerical or string): ").strip()
                if not pc_value:
                    print("Error: 'PC' column value cannot be empty. skipping this subfolder.")
                    continue
                #run through text files in subfolder
                for txt_file in glob.glob(os.path.join(subfolder_path, '*.txt')):
                    print(f"Reading file: {txt_file}")
                    try:
                        #Load file and process
                        df = pd.read_csv(txt_file, skiprows=8, engine='python', sep='\t')
                    except Exception as e:
                        print(f"Error reading file '{txt_file}': {e}")
                        continue
                    #adding 'PC' column with the value
                    df['PC'] = pc_value
                    #extract plate name from file name
                    plate_name = os.path.splitext(os.path.basename(txt_file))[0]
                    print(f"Processing plate: {plate_name}")
                    #output folder
                    output_folder = os.path.join(main_folder, "results", f"{subfolder_name}_csv_files")
                    os.makedirs(output_folder, exist_ok=True)
                    #save CSV
                    csv_filename = f"{plate_name}_csv.csv"
                    #prompt user for flagged wells to remove
                    flagged_wells_input = input(
                        f"Enter flagged wells(bad wells) to remove for plate '{plate_name}' (e.g., row_column, ...) or 'none' if none: ").strip()
                    #process flagged wells
                    flagged_wells = []
                    if flagged_wells_input.lower() != 'none':
                        flagged_wells = [fw.strip() for fw in flagged_wells_input.split(',') if fw.strip()]
                    # remove flagged wells
                    for flagged_well in flagged_wells:
                        try:
                            row, column = flagged_well.split('_')
                            row = int(row)
                            column = int(column)
                            df = df.loc[~((df['Row'] == row) & (df['Column'] == column))]
                            print(f"Flagged well removed: Row {row}, Column {column}")
                        except Exception as e:
                            print(f"Error processing flagged well '{flagged_well}': {e}")
                            continue
            #editing column names
                    try:
                        df.columns = [col.split("[")[0].strip() if "[" in col else col for col in df.columns]
                        print("Column names cleaned.")
                    except Exception as e:
                        logging.error(f"Error cleaning column names: {e}")
                        return
                    #save modified DF
                    csv_path = os.path.join(output_folder, csv_filename)
                    try:
                        df.to_csv(csv_path, index=False)
                        print(f"File '{csv_filename}' saved in folder '{output_folder}'")
                    except Exception as e:
                        print(f"Error saving file '{csv_filename}': {e}")
                        continue
        #done processing all subfolders
        print("\nAll subfolders processed successfully.")
        print(f"review the files saved in {main_folder}/results.")

    except Exception as e:
        print(f"Unexpected error occurred: {e}")
##################
def prompt_for_review():
    input("review the csv files. when ready, press enter to continue...")
######################
def process_data_combine(folder_path):
    try:
        #output result folder
        output_folder = os.path.join(main_folder, "ROS_result_combined")
        os.makedirs(output_folder, exist_ok=True)
        result_combine_dir = os.path.join(main_folder, "ROS_result_combined/combined_csv_file")
        if not os.path.exists(result_combine_dir):
            os.makedirs(result_combine_dir, exist_ok=True)
        #subdirectories
        subdirectories = [subdir for subdir in os.listdir(folder_path) if os.path.isdir(os.path.join(folder_path, subdir))]
        combined_dfs = []
        #run on each subdirectory
        for subdir in subdirectories:
            sub_dir_path = os.path.join(folder_path, subdir)
            print("Processing files in:", sub_dir_path)
            dfs = []
            for file in os.listdir(sub_dir_path):
                if file.endswith('.csv'):
                    file_path = os.path.join(sub_dir_path, file)
                    df = pd.read_csv(file_path, encoding='unicode_escape')
                    df['plate'] = os.path.splitext(file)[0]  #add plate column- important
                    dfs.append(df)
        print("plate column added to data frame")
        combined_df = pd.concat(dfs, ignore_index=True)
        combined_dfs.append(combined_df)
        combined_df = pd.concat(combined_dfs, ignore_index=True)
        print("Data frames combined!")
        combined_df['well'] = combined_df['Row'].astype('str') + combined_df['Column'].astype('str')
        combined_df['well'] = combined_df['well'].astype('str')
        print("well column added to data frame")
        combined_df.to_csv(os.path.join(result_combine_dir, "combined_df_data.csv"), index=False)
        print(f"review the files saved in {main_folder}/ROS_result_combined/combined_csv_file.")
    except Exception as e:
        logging.error(f"failed to combine plates...{e}")
####################
def process_plate_data_on_folders(main_folder, unwanted_columns, important_features, index_columns):
    try:
    ## open main folder for dash
        main_results_dir = os.path.join(main_folder, "main_results_folder")
        os.makedirs(main_results_dir, exist_ok=True)
        unwanted_columns = list(set(unwanted_columns))
        important_features = list(set(important_features))
        index_columns = list(set(index_columns))
        #running on each item in the main input folder
        for item_name in os.listdir(main_folder):
            item_path = os.path.join(main_folder, item_name)
            if os.path.isdir(item_path):
                #running on each csv in each subfolder
                for file_name in os.listdir(item_path):
                    if file_name.endswith(".csv"):
                        file_path = os.path.join(item_path, file_name)
                        print(f"Processing file: {file_path}")
                        #load each plate
                        plate_df = pd.read_csv(file_path)
                        #results folder with the name of the csv file
                        results_folder = os.path.join(item_path, f"results_{file_name}_res")
                        os.makedirs(results_folder, exist_ok=True)
                        #drop unwanted columns & start proc
                        try:
                            plate_df = plate_df.drop(columns=unwanted_columns)
                            print(f"Dropped unwanted columns: {unwanted_columns}")
                        except KeyError as e:
                            logging.warning(f"Error dropping unwanted columns: {e}")
                        #checking important features
                        print("These are the important columns chosen for the analysis:")
                        print(important_features)
                        missing_features = [feature for feature in important_features if feature not in plate_df.columns]
                        if missing_features:
                            logging.warning(f"Missing important features: {missing_features}")
                            important_features = [f for f in important_features if f in plate_df.columns]
                        else:
                            print("all selected features found!")
                        if not important_features:
                            logging.error("No valid important features available. Skipping.")
                            continue
                        try:
                            print("These are the index columns chosen for the analysis:")
                            print(index_columns)
                            index_columns.append('well')
                            index_columns.append('plate')
                            print("'well' & 'plate' columns added to index_columns.")
                        except Exception as e:
                            logging.error(f"Error adding 'well' and 'plate' columns to index columns: {e}")
                        #generate combinations of index_columns for new indexes
                        try:
                            new_index_columns = []
                            for r in range(1, len(index_columns) + 1):
                                for comb in itertools.combinations(index_columns, r):
                                    new_index_columns.append(comb)
                            print(f"Generated {len(new_index_columns)} new index column combinations.")
                        except Exception as e:
                            logging.error(f"Error generating new index column combinations: {e}")
                        #create new index columns based on combinations
                        try:
                            new_cols = {}
                            for cols in new_index_columns:
                                col_name = '_'.join(cols)
                                if col_name not in plate_df.columns:
                                    new_cols[col_name] = plate_df[list(cols)].apply(lambda x: '_'.join(map(str, x)), axis=1)
                                else:
                                    print(f"{col_name} already exists and will not be added")
                            plate_df = pd.concat([plate_df, pd.DataFrame(new_cols)], axis=1)
                            print("Created new index columns based on combinations.")
                        except Exception as e:
                            logging.error(f"Error creating new index columns: {e}")
                        ## keep only new_index and imp features
                        try:
                            print("subsetting data to focus on important features & index columns...")
                            columns1 = plate_df[index_columns]
                            columns2 = plate_df[important_features]
                            plate_df_int = pd.concat([columns1, columns2, pd.DataFrame(new_cols)], axis=1)
                        except Exception as e:
                            logging.error(f"Failed to subset{e}")
                            plate_df_int = plate_df.copy()
                        try:
                            # remove rows and columns with ≥90% zero or NA values
                            filtered_df = plate_df_int.drop(columns=index_columns, errors='ignore')
                            filtered_df = filtered_df.dropna(thresh=0.1 * len(filtered_df), axis=0)
                            filtered_df = filtered_df.replace(0, np.nan)
                            filtered_df = filtered_df.dropna(thresh=0.1 * len(filtered_df.columns), axis=1)
                            final_df = pd.concat([columns1, filtered_df], axis=1)
                            print("Dropped rows and columns with ≥90% zero or NA values.")
                            plate_df_2 = final_df.copy()
                        except Exception as e:
                            logging.error(f"Error dropping rows/columns with high NA/zero values: {e}")
                        ## add cell count feat
                        try:
                            print("usual cell_count_feat names: all_cells - Number of Objects/All_cells - Number of Objects/selected_cells - Number of Objects")
                            cell_count_feat = input("Enter cell_count_feat(check in csv if not sure): ")
                            plate_df_2[cell_count_feat] = plate_df[cell_count_feat]
                            plate_df_2[cell_count_feat] = pd.to_numeric(plate_df_2[cell_count_feat])
                        except Exception as e:
                            logging.error(f"failed to locate cell count feat {e}")
                        ##user to remove rows based on any avaialble columns
                        try:
                            print("Section to remove rows based on chosen column, choose column following by valuesto remove, otherwise, type no")
                            while True:
                                print("Columns in the dataset:")
                                for idx, column in enumerate(plate_df_2.columns):
                                    print(f"{idx}: {column}")
                                column_to_filter = input("Enter the column name to filter by (or 'no' to skip): ").strip()
                                if column_to_filter.lower() == 'no':
                                    break
                                if column_to_filter not in plate_df_2.columns:
                                    print("Invalid column name. Please try again.")
                                    continue
                                print(f"Values in the selected column '{column_to_filter}':")
                                print(plate_df_2[column_to_filter].unique())
                                values_to_remove = input("Enter the values to remove (comma-separated): ").split(',')
                                values_to_remove = [val.strip() for val in values_to_remove]
                                plate_df_2 = plate_df_2[~plate_df_2[column_to_filter].isin(values_to_remove)]
                                print(f"Rows with values {values_to_remove} in column '{column_to_filter}' have been removed.")
                        except Exception as e:
                            logging.error(f"Error removing custom column values: {e}")
                            return
                        #remove additional columns if needed
                        try:
                            print("Section to remove columns from the data, choose column to remove, otherwise, type none")
                            for idx, column in enumerate(plate_df_2.columns):
                                print(f"{idx + 1}: {column}")
                            features_to_remove = input("Enter comma-separated column names to remove (or 'none'): ").strip()
                            if features_to_remove.lower() != 'none':
                                features_to_remove = [col.strip() for col in features_to_remove.split(',')]
                                plate_df_2.drop(columns=features_to_remove, inplace=True)
                                print(f"Features removed: {features_to_remove}")
                        except KeyError as e:
                            logging.warning(f"Some features to remove were not found: {e}")
                        except Exception as e:
                            logging.error(f"Error removing features: {e}")
                        #IDs to remove ?
                        try:
                            print("Section to remove cell id values from the data, choose id to remove, otherwise, type none")
                            print(plate_df_2["cell_id"].unique())
                            remove_id = input("Enter ID value to remove (or 'none'): ").strip()
                            if remove_id.lower() != 'none':
                                plate_df_2 = plate_df_2[~plate_df_2['cell_id'].astype(str).str.contains(remove_id)]
                                print(f"Rows with ID '{remove_id}' removed.")
                        except KeyError:
                            logging.warning("Column 'ID' not found. Skipping ID removal.")
                        except Exception as e:
                            logging.error(f"Error removing rows based on ID: {e}")
                        #save first df
                        plate_df_2.to_csv(os.path.join(results_folder, "df_proc_1.csv"), index=False)
                        print("first processed data frame saved to result directory")
                        #directory for main results each subfolder
                        general_results_dir = os.path.join(results_folder, "general_results")
                        os.makedirs(general_results_dir, exist_ok=True)
                        #index column for statistical summary
                        try:
                            print(new_index_columns)
                            index_col = input("Enter the index column for statistical summary: seperated by '_' ").strip()
                            summary_df = plate_df_2.groupby(index_col).describe()
                            print(f"Generated statistical summary grouped by '{index_col}'.")
                        except KeyError:
                            logging.error(f"Index column '{index_col}' not found in DataFrame.")
                            summary_df = pd.DataFrame()
                        except Exception as e:
                            logging.error(f"Error generating statistical summary: {e}")
                            summary_df = pd.DataFrame()
                        #save statistical summary
                        try:
                            summary_excel_path = os.path.join(general_results_dir, "init_statistical_summary.xlsx")
                            summary_df.to_excel(summary_excel_path)
                            print(f"Statistical summary saved to '{summary_excel_path}'")
                        except Exception as e:
                            logging.error(f"Error saving statistical summary to Excel: {e}")
                        #directory for cell count results
                        print("cell count process initiated")
                        cell_count_results_dir = os.path.join(results_folder, "cell_count_results")
                        os.makedirs(cell_count_results_dir, exist_ok=True)
                        #histogram for cell count
                        try:
                            if cell_count_feat in plate_df_2.columns:
                                fig = px.histogram(plate_df_2, x=cell_count_feat, title="Histogram of Cell Count")
                                histogram_file_path = os.path.join(cell_count_results_dir, "histogram.html")
                                fig.write_html(histogram_file_path)
                                print(f"Cell count histogram saved to {histogram_file_path}")
                            else:
                                logging.warning(f"{cell_count_feat} column not found. histogram not generated.")
                        except Exception as e:
                            logging.error(f"Error generating histogram: {e}")
                        #statistics for cell count
                        try:
                            if cell_count_feat in plate_df_2.columns:
                                stats = plate_df_2[cell_count_feat].describe()
                                stats_file_path = os.path.join(cell_count_results_dir, "cell_count-statistics.txt")
                                with open(stats_file_path, "w") as f:
                                    f.write(stats.to_string())
                                print(f"Cell count statistics saved to {stats_file_path}")
                            else:
                                logging.warning(f"{cell_count_feat} column not found. statistics not generated.")
                        except Exception as e:
                            logging.error(f"Error creating cell count statistics: {e}")
                        #input for minimum and maximum cell count values
                        try:
                            min_input = input("Enter the minimum value for cell count (leave blank to skip): ").strip()
                            max_input = input("Enter the maximum value for cell count (leave blank to skip): ").strip()
                            min_value = float(min_input) if min_input else None
                            max_value = float(max_input) if max_input else None
                            if min_value is not None and max_value is not None:
                                plate_df_2 = plate_df_2[(plate_df_2[cell_count_feat] >= min_value) & (
                                        plate_df_2[cell_count_feat] <= max_value)]
                                print(f"Filtered data with cell count between {min_value} and {max_value}")
                            elif min_value is not None:
                                plate_df_2 = plate_df_2[plate_df_2[cell_count_feat] >= min_value]
                                print(f"Filtered data with cell count >= {min_value}")
                            elif max_value is not None:
                                plate_df_2 = plate_df_2[plate_df_2[cell_count_feat] <= max_value]
                                print(f"Filtered data with cell count <= {max_value}")
                            else:
                                print("No filtering applied to cell count.")
                        except ValueError:
                            logging.error("Invalid input for minimum or maximum cell count. Skipping filtering.")
                        except Exception as e:
                            logging.error(f"Error filtering cell count: {e}")
                        plate_df_2.to_csv(os.path.join(results_folder, "df_proc_2.csv"))
                        print("second processed data frame saved to result directory")
                        print(new_index_columns)
                        #directory for box plots results
                        box_plot_dir = os.path.join(results_folder, "box_plots")
                        os.makedirs(box_plot_dir, exist_ok=True)
                        #prompting index column for box plots
                        index_col = input("Enter the index column for box plot vis: ")
                        bar_df = plate_df_2.copy()
                        # box plots for selected index- important features only
                        try:
                            print("box plotting selected features!")
                            for col in bar_df:
                                if col in important_features:
                                    bp = px.box(bar_df, x=index_col, y=bar_df[col], points="all", color=index_col, notched=True)
                                    bp.update_traces(quartilemethod="inclusive")  # or "inclusive", or "linear" by default
                                    bp.update_layout(
                                        font_family="Arial",
                                        font_color="Black",
                                        font_size=20,
                                        font=dict(
                                            family="Arial",
                                            size=20,
                                            color="Black"
                                        )
                                    )
                                    print("one_down")
                                    bp.write_html(os.path.join(box_plot_dir, f'{col}_bar_plot.html'))
                                    #bp.write_image(os.path.join(box_plot_dir, f"box_plot{col}.pdf"), engine="kaleido") ## pdf vs
                            print("box plots done!")
                        except Exception as e:
                            logging.error(f"failed box plotting{e}")
                        #directory for summary files
                        summary_files_dir = os.path.join(results_folder, "summary_files")
                        os.makedirs(summary_files_dir, exist_ok=True)
                        #creating excel  with sheets with split data and summary stats
                        try:
                            print("creating excel files and statistical summary & ppt for combination of groups!")
                            group_columns = ['cell_type', 'compound', 'concentration']
                            excel_file_path = os.path.join(summary_files_dir, f'summary_file.xlsx')
                            with pd.ExcelWriter(excel_file_path, engine='openpyxl') as writer:
                                for name, group in plate_df_2.groupby(group_columns):
                                    try:
                                        sheet_name = f"{name[0]}_{name[1]}_{name[2]}"
                                        group.to_excel(writer, sheet_name=sheet_name, index=False)
                                        ####summary statistics to another sheet
                                        summary_df = group.describe(include='all').T[
                                            ['mean', 'std', 'min', '25%', '50%', '75%', 'max']]
                                        summary_df.to_excel(writer, sheet_name=f"{sheet_name}_Summary")
                                    except ValueError:
                                        pass
                                if not writer.sheets:
                                    pd.DataFrame({'Dummy': [1]}).to_excel(writer, sheet_name='DummySheet', index=False)
                            #ppt with slides for each unique combination
                            print("ppt file generated!")
                            ppt_file_path = os.path.join(summary_files_dir, f'slides.pptx')
                            prs = Presentation()
                            for name, _ in plate_df_2.groupby(group_columns):
                                slide = prs.slides.add_slide(prs.slide_layouts[5])
                                title = slide.shapes.title
                                title.text = f"Combination: {name[0]}, {name[1]}, {name[2]}"
                            prs.save(ppt_file_path)
                        except Exception as e:
                            logging.error(f"failed creating summary{e}")
                        #bar plot for Number of Cells
                        try:
                            print("working on final cell count vis!")
                            group_columns = ['cell_type', 'compound', 'concentration', 'cell_id']
                            plt.figure(figsize=(20, 18))
                            for name, group in plate_df_2.groupby(group_columns):
                                try:
                                    plt.bar(str(name), group[cell_count_feat].mean(),
                                            yerr=group[cell_count_feat].std(), capsize=5,
                                            label=str(name))
                                except ValueError:
                                    pass
                            plt.title(f'Bar Plot for Number of Cells - {file_name}')
                            plt.xlabel('cell_type, compound, concentration, cell_id')
                            plt.ylabel('Number of Cells')
                            plt.xticks(rotation=45, ha='right')
                            plt.legend()
                            plt.savefig(os.path.join(general_results_dir, f'number_of_cells_bar_plot.pdf'))
                            plt.close()
                        except Exception as e:
                            logging.error(f"failed bar plotting cell number{e}")
                        ## remove cell feat
                        try:
                            plate_df_2 = plate_df_2.drop(columns=[cell_count_feat])
                        except Exception as e:
                            logging.error(f"failed to remove {cell_count_feat}...{e}")
                        #heatmap for selected features (mean values & normalized)
                        try:
                            heat_df = plate_df_2.copy()
                            index_col = input("Enter the index column for heatmap: ")
                            print("heatmap prep!")
                            heat_df = heat_df.set_index(index_col, drop=True)
                            numeric_features = heat_df[important_features].select_dtypes(include=np.number)
                            numeric_features_normalized = (numeric_features - numeric_features.mean()) / numeric_features.std()
                            numeric_features_normalized = numeric_features_normalized.groupby(index_col).mean()
                            lan = plt.rcParams['font.family'] = ['Arial']
                            plt.figure(figsize=(20, 12))
                            font = {'family': 'Arial',
                                    'size': 22}
                            plt.rc('font', **font)
                            plt.yticks(family='Arial', rotation="horizontal", size=10)
                            plt.xticks(family='Arial', rotation="vertical", size=8)
                            plt.margins(0.1)
                            plt.subplots_adjust(bottom=0.3)
                            sns.heatmap(numeric_features_normalized, annot=True, cmap='coolwarm')
                            plt.title(f'Heatmap for Selected Numeric Features - {file_name}')
                            plt.tight_layout()
                            plt.savefig(os.path.join(general_results_dir, f'heatmap.pdf'))
                            plt.close()
                            print("heatmap done!")
                        except Exception as e:
                            logging.error(f"Heatmap failed: {e}")
                        #PCA based on all numeric features without ID
                        try:
                            print("creating pca for group with no ID and with id")
                            pca = PCA(n_components=2)
                            # Impute missing values in numeric features
                            numeric_features = plate_df_2.select_dtypes(include=np.number)
                            imputer = SimpleImputer(strategy='mean')
                            numeric_features_imputed = imputer.fit_transform(numeric_features)
                            #interactive pcas
                            pca_result_without_id = pca.fit_transform(numeric_features_imputed)
                            label_without_id = plate_df_2['cell_type'].astype(str) + '_' + \
                                            plate_df_2['compound'].astype(str) + '_' + \
                                            plate_df_2['concentration'].astype(str)
                            color_labels_wo_id = pd.factorize(label_without_id)[0]
                            pca_df1 = pd.DataFrame({
                                'PC1': pca_result_without_id[:, 0],
                                'PC2': pca_result_without_id[:, 1],
                                'Label': label_without_id,
                                'Color': color_labels_wo_id,
                                'cell_id': plate_df_2['cell_id']
                            })
                            fig = px.scatter(pca_df1, x='PC1', y='PC2', color='Label',
                                             title=f'PCA Plot for All Numeric Features without CELL ID - {file_name}', hover_data=['cell_id'],
                                             color_continuous_scale='viridis')
                            fig.update_layout(legend_title='Cell Type, Compound, Concentration')
                            fig.add_annotation(text=f'Explained Variance Ratio: {pca.explained_variance_ratio_}', xref='paper',
                                               yref='paper',
                                               x=0.98, y=0.02, showarrow=False, font=dict(size=10))
                            fig.write_html(os.path.join(general_results_dir, f"pca_plot_wo_id_{file_name[:-4]}.html"))
                            #fig.write_image(os.path.join(general_results_dir, f'pca_plot_without_id2.pdf'), engine="kaleido")
                            # PCA based on all numeric features with CELL ID-2 interactive
                            pca_result_with_id = pca.fit_transform(numeric_features_imputed)
                            # A label for each unique combination of Cell Type, Compound, Concentration, and CELL ID
                            label_with_id = plate_df_2['cell_type'].astype(str) + '_' + \
                                               plate_df_2['compound'].astype(str) + '_' + \
                                               plate_df_2['concentration'].astype(str) + '_' + \
                                               plate_df_2['cell_id'].astype(str)
                            color_labels_with_id = pd.factorize(label_with_id)[0]
                            # Create DataFrame for Plotly
                            pca_df2 = pd.DataFrame({
                                'PC1': pca_result_with_id[:, 0],
                                'PC2': pca_result_with_id[:, 1],
                                'Label': label_with_id,
                                'Color': color_labels_with_id,
                                'cell_id': plate_df_2['cell_id']
                            })
                            fig = px.scatter(pca_df2, x='PC1', y='PC2', color='Label',
                                             title=f'PCA Plot for All Numeric Features with CELL ID - {file_name}', hover_data=['cell_id'],
                                             color_continuous_scale='viridis')
                            fig.update_layout(legend_title='Cell Type, Compound, Concentration, ID')
                            fig.add_annotation(text=f'Explained Variance Ratio: {pca.explained_variance_ratio_}', xref='paper',
                                               yref='paper',
                                               x=0.98, y=0.02, showarrow=False, font=dict(size=10))
                            fig.write_html(os.path.join(general_results_dir, f"pca_plot_with_id_{file_name[:-4]}.html"))
                        except Exception as e:
                            logging.error(f"PCA failed: {e}")
                        print("pca done!")
                        #directory for outlier results
                        try:
                            outlier_files_dir = os.path.join(results_folder, "outlier_files")
                            os.makedirs(outlier_files_dir, exist_ok=True)
                            for col in plate_df_2.columns:
                                if col in index_columns:
                                    plate_df_2[col] = plate_df_2[col].astype(str)
                            perform_outlier_detection = input(
                                "Do you want to perform outlier detection? (yes/no): ").strip().lower()
                            if perform_outlier_detection == 'yes':
                                print("Applying basic Outlier methods...")
                                # prompting index column for outliers detecetion
                                print("Columns available for outlier detection:")
                                for idx, column in enumerate(plate_df_2.columns):
                                    print(f"{idx}: {column}")
                                index_column = input("Choose one column as index for outlier detection: ")
                                out_dat_init = plate_df_2.copy()
                                # final filtered DataFrame
                                init_file_path = os.path.join(outlier_files_dir,
                                                               'out_dat_init.csv')
                                out_dat_init.to_csv(init_file_path, index=False)
                                #outlier detection functions options
                                ## perc
                                def compute_percentiles(df, output_file):
                                    up_bounds = []
                                    low_bounds = []
                                    above_count = []
                                    below_count = []
                                    numeric_cols = df.select_dtypes(include='number').columns
                                    for col in numeric_cols:
                                        percentile_low = df[col].quantile(0.01)
                                        percentile_high = df[col].quantile(0.99)
                                        up_bound = percentile_high
                                        low_bound = percentile_low
                                        above = df[df[col] > percentile_high]
                                        below = df[df[col] < percentile_low]
                                        above_count.append(len(above))
                                        below_count.append(len(below))
                                        df = df[(df[col] >= percentile_low) & (df[col] <= percentile_high)]
                                        up_bounds.append(up_bound)
                                        low_bounds.append(low_bound)
                                    output = pd.DataFrame(
                                        {'Column': numeric_cols, 'Up Bound': up_bounds, 'Low Bound': low_bounds,
                                         'Above Count': above_count,
                                         'Below Count': below_count})
                                    output.to_csv(output_file, index=False)
                                    return df
                                ## iqr
                                def compute_iqr(df, output_file):
                                    up_bounds = []
                                    low_bounds = []
                                    above_count = []
                                    below_count = []
                                    numeric_cols = df.select_dtypes(include='number').columns
                                    for col in numeric_cols:
                                        percentile25th = df[col].quantile(0.25)
                                        percentile75th = df[col].quantile(0.75)
                                        iqr = percentile75th - percentile25th
                                        up_bound = percentile75th + 1.5 * iqr
                                        low_bound = percentile25th - 1.5 * iqr
                                        above = df[df[col] > up_bound]
                                        below = df[df[col] < low_bound]
                                        above_count.append(len(above))
                                        below_count.append(len(below))
                                        df = df[(df[col] >= low_bound) & (df[col] <= up_bound)]
                                        up_bounds.append(up_bound)
                                        low_bounds.append(low_bound)
                                    output = pd.DataFrame(
                                        {'Column': numeric_cols, 'Up Bound': up_bounds, 'Low Bound': low_bounds,
                                         'Above Count': above_count,
                                         'Below Count': below_count})
                                    output.to_csv(output_file, index=False)
                                    return df
                                ## zscore
                                def compute_scores(df, output_file):
                                    up_scores = []
                                    low_scores = []
                                    above_count = []
                                    below_count = []
                                    numeric_cols = df.select_dtypes(include='number').columns
                                    for col in numeric_cols:
                                        up_score = df[col].mean() + 3 * df[col].std()
                                        low_score = df[col].mean() - 3 * df[col].std()
                                        above = df[df[col] > up_score]
                                        below = df[df[col] < low_score]
                                        above_count.append(len(above))
                                        below_count.append(len(below))
                                        df = df[(df[col] >= low_score) & (df[col] <= up_score)]
                                        up_scores.append(up_score)
                                        low_scores.append(low_score)
                                    output = pd.DataFrame(
                                        {'Column': numeric_cols, 'Up Score': up_scores, 'Low Score': low_scores,
                                         'Above Count': above_count,
                                         'Below Count': below_count})
                                    output.to_csv(output_file, index=False)
                                    return df
                                ## LOF
                                def LOF_outlier(df, output_file1, output_file2):
                                    # outlier detection using LOF
                                    types = df[index_column].unique()
                                    print(types)
                                    mask = []
                                    features = df.select_dtypes(include='number').columns
                                    detector_list = [
                                        ("Local Outlier Factor 30", LocalOutlierFactor(n_neighbors=8))
                                    ]
                                    for name, algorithm in detector_list:
                                        errors = np.full(len(df), fill_value=np.nan)
                                        outliers = np.full(len(df), fill_value=np.nan)
                                        for type in types:
                                            x = df.loc[:, features].values
                                            F = x.sum(1)
                                            mask = np.zeros(x.shape[0])
                                            mask[np.isfinite(F)] = 1
                                            mask_type = mask * np.array(df[index_column] == type)
                                            Curr_df = df.loc[mask_type == 1, features]
                                            x = Curr_df.values
                                            if name == "Local Outlier Factor 30":
                                                algorithm.fit(x)
                                                errors[mask_type == 1] = algorithm.negative_outlier_factor_
                                                outliers[mask_type == 1] = algorithm.fit_predict(x)
                                        df[name] = errors
                                        df[f'{name}_outliers'] = outliers
                                        df.set_index(name, inplace=True,
                                                     append=True, drop=False)
                                        df.to_csv(os.path.join(output_file1))
                                        # exclude rows that were defined as outliers
                                        for col in df.columns:
                                            if col.endswith("_outliers"):
                                                df = df[df[col] != -1]
                                    df.to_csv(os.path.join(output_file2))
                                    print("outliers removed!")
                                    return df
                                #run the outlier detection functions
                                df_percentiles = compute_percentiles(out_dat_init,
                                                                     os.path.join(outlier_files_dir,
                                                                                  'percentiles_outliers.csv'))
                                perc_file_path = os.path.join(outlier_files_dir,
                                                               'perc_dat.csv')
                                df_percentiles.to_csv(perc_file_path, index=False)
                                df_iqr = compute_iqr(out_dat_init, os.path.join(outlier_files_dir, 'iqr_outliers.csv'))
                                iqr_file_path = os.path.join(outlier_files_dir,
                                                               'iqr_dat.csv')
                                df_iqr.to_csv(iqr_file_path, index=False)
                                df_scores = compute_scores(out_dat_init,
                                                           os.path.join(outlier_files_dir, 'scores_outliers.csv'))
                                zscore_file_path = os.path.join(outlier_files_dir,
                                                               'zscore_dat.csv')
                                df_scores.to_csv(zscore_file_path, index=False)
                                df_lof = LOF_outlier(out_dat_init,
                                                     os.path.join(outlier_files_dir, 'processed_data_outliers_1.csv'),
                                                     os.path.join(outlier_files_dir, 'processed_data_outliers_2.csv'))
                                print("Please review outlier detection methods in outlier results folder and choose one method to continue with:")
                                method_choice = input("Enter 1 for percentiles, 2 for IQR, 3 for scores, 4 for lof: ")
                                if method_choice == '1':
                                    out_dat = df_percentiles
                                elif method_choice == '2':
                                    out_dat = df_iqr
                                elif method_choice == '3':
                                    out_dat = df_scores
                                elif method_choice == '4':
                                    out_dat = df_lof
                                    out_dat = out_dat.drop(columns=["Local Outlier Factor 30", "Local Outlier Factor 30_outliers", "Local Outlier Factor 30"])
                                else:
                                    print("Invalid choice. Defaulting to original DataFrame.")
                                    out_dat = out_dat_init
                                    out_dat = out_dat.drop(
                                        columns=["Local Outlier Factor 30", "Local Outlier Factor 30_outliers",
                                                 "Local Outlier Factor 30"])
                            else:
                                print("No outlier detection ran.")
                                out_dat = plate_df_2.copy()
                            #final filtered DataFrame
                            final_file_path = os.path.join(results_folder, 'final_filtered_data(post outlier detection).csv')
                            out_dat.to_csv(final_file_path, index=False)
                            print("Final filtered DataFrame saved to", final_file_path)
                        except Exception as e:
                            logging.error(f"Failed to perform outlier detection{e}")
                        #prompting index column for final proc- "PC" is must!
                        for idx, column in enumerate(out_dat.columns):
                            print(f"{idx}: {column}")
                        #columns to set as indexes
                        try:
                            print(Fore.RED + "==== Please make sure you choose all needed index columns - PC is must! ====")
                            columns_to_set_index = input("Enter ALL needed columns to set as indexes -PC IS MUST! (comma-separated): ").split(',')
                            columns_to_set_index = [col.strip() for col in columns_to_set_index if col.strip()]
                            print(f"Columns to set as index: {columns_to_set_index}")
                            if columns_to_set_index:
                                missing_index_cols = [col for col in columns_to_set_index if col not in out_dat.columns]
                                if missing_index_cols:
                                    logging.warning(f"Index columns not found: {missing_index_cols}")
                                out_dat = out_dat.set_index([col for col in columns_to_set_index if col in out_dat.columns], drop=False)
                                print(f"Set columns as index: {columns_to_set_index}")
                        except Exception as e:
                            logging.error(f"Error setting index columns: {e}")
                        out_dat = out_dat.select_dtypes(include=np.number)
                        #target variable selection for feature qc
                        try:
                            #user to select the correct index title as the target variable
                            if out_dat.index.names:
                                print("Select the correct index title to be the target variable for feature quality control:")
                                for idx, index_title in enumerate(out_dat.index.names):
                                    print(f"{idx + 1}: {index_title}")
                                while True:
                                    try:
                                        selected_index = int(
                                            input("Enter the number corresponding to the correct index title: ")) - 1
                                        if selected_index < 0 or selected_index >= len(out_dat.index.names):
                                            raise ValueError("Selection out of range.")
                                        break
                                    except ValueError as e:
                                        logging.error(f"Invalid selection. Please try again. Error: {e}")
                                #selected index title as the target variable
                                target_variable = out_dat.index.names[selected_index]
                                print(f"Target variable selected: {target_variable}")
                            else:
                                raise ValueError("No index names found in the data.")
                            #factor the target variable
                            target_variable_labels, _ = pd.factorize(out_dat.index.get_level_values(target_variable))
                            print(f"Target variable labels: {target_variable_labels}")
                            #directory for feature qc results
                            feature_qc_dir = os.path.join(results_folder, "feature_qc")
                            os.makedirs(feature_qc_dir, exist_ok=True)
                            #statistical summary for features
                            statistical_summary = out_dat.describe()
                            statistical_summary.to_csv(os.path.join(feature_qc_dir, "statistical_summary.csv"))
                            print("Statistical summary saved to 'statistical_summary.csv'.")
                            #feature stability scores -CV - for normalizied features
                            normalized_data = (out_dat - out_dat.min()) / (out_dat.max() - out_dat.min())
                            if normalized_data.isna().any().any():
                                logging.warning("Normalization produced NaN values. Check your data for zero variance features.")
                            stability_scores = normalized_data.var() / normalized_data.mean()
                            print("Stability scores calculated.")
                            #correlation
                            correlation_matrix = out_dat.corr(method='spearman')
                            correlation_matrix.to_csv(os.path.join(feature_qc_dir, "correlation_matrix.csv"))
                            print("Correlation matrix saved to 'correlation_matrix.csv'.")
                            #plot the correlation matrix
                            lan = plt.rcParams['font.family'] = ['Arial']
                            plt.figure(figsize=(32, 18))
                            font = {'family': 'Arial',
                                    'size': 10}
                            plt.rc('font', **font)
                            plt.yticks(family='Arial', rotation="horizontal", size=10)
                            plt.xticks(family='Arial', rotation="vertical", size=8)
                            plt.margins(0.1)
                            plt.subplots_adjust(bottom=0.3)
                            sns.heatmap(correlation_matrix, annot=True, cmap='coolwarm', fmt=".2f")
                            plt.title("Correlation Matrix")
                            plt.tight_layout()
                            correlation_plot_path = os.path.join(feature_qc_dir, "correlation_matrix_plot.png")
                            plt.savefig(correlation_plot_path)
                            plt.close()
                            print(f"Correlation matrix plot saved to: {correlation_plot_path}")
                            #CV plot
                            if stability_scores is not None and not stability_scores.isna().all():
                                plt.figure(figsize=(32, 16))
                                plt.bar(stability_scores.index, stability_scores.values, color='skyblue')
                                plt.title("Feature Stability Scores (Coefficient of Variation)")
                                plt.xlabel("Feature")
                                plt.ylabel("Stability Score")
                                plt.xticks(rotation=45, ha='right')
                                plt.margins(0.1)
                                plt.tight_layout()
                                cv_plot_path = os.path.join(feature_qc_dir, "stability_scores_plot.png")
                                plt.savefig(cv_plot_path)
                                plt.close()
                                print(f"CV plot saved to: {cv_plot_path}")
                            else:
                                logging.warning("Stability scores not available or invalid. Skipping CV plot.")
                            #random forest score computation per features
                            rf_scores = {}
                            try:
                                for feature in out_dat.columns:
                                    model = RandomForestRegressor()
                                    feature_values = out_dat[feature].values.reshape(-1, 1)
                                    model.fit(feature_values, target_variable_labels)
                                    rf_scores[feature] = model.score(feature_values, target_variable_labels)
                                    logging.debug(f"Random Forest score for feature '{feature}': {rf_scores[feature]}")
                            except Exception as e:
                                logging.error(f"Error occurred during Random Forest model training: {e}")
                                rf_scores[feature] = np.nan
                            #RF scores to csv
                            rf_scores_df = pd.DataFrame(rf_scores, index=['RF_Score'])
                            rf_scores_df.to_csv(os.path.join(feature_qc_dir, "random_forest_scores.csv"))
                            print("Random Forest scores saved to 'random_forest_scores.csv'.")
                            #plot RF scores
                            plt.figure(figsize=(32, 16))
                            plt.bar(rf_scores.keys(), rf_scores.values(), color='lightgreen')
                            plt.title("Random Forest Scores for each Feature")
                            plt.xlabel("Feature")
                            plt.ylabel("Random Forest Score")
                            plt.xticks(rotation=45, ha='right')
                            plt.margins(0.1)
                            plt.tight_layout()
                            rf_plot_path = os.path.join(feature_qc_dir, "random_forest_scores_plot.png")
                            plt.savefig(rf_plot_path)
                            plt.close()
                            print(f"Random Forest scores plot saved to: {rf_plot_path}")
                            print(f"Feature QC finsihed and saved in{feature_qc_dir}")
                            #want to change feature names?
                            rename_columns = input("Do you want to rename any columns (recommended for long feature names) ? (yes/no): ").strip().lower()
                            if rename_columns == 'yes':
                                print("Current Columns:")
                                for idx, column in enumerate(out_dat.columns):
                                    logging.info(f"{idx}: {column}")
                                renaming_input = input(
                                    "Enter the column name and new name separated by ':', each pair separated by commas-no space! (e.g., col1:new_col1,col2:new_col2): ").strip()
                                try:
                                    rename_dict = dict(pair.split(':') for pair in renaming_input.split(',') if ':' in pair)
                                    if rename_dict:
                                        out_dat.rename(columns=rename_dict, inplace=True)
                                        print("Names changed successfully.")
                                    else:
                                        logging.warning("No valid renaming pairs provided.")
                                except Exception as e:
                                    logging.error(f"Error renaming columns: {e}")
                            out_dat.to_csv(os.path.join(results_folder, "data_with_new_names.csv"))
                            print(
                                f"Data with new names saved to: {os.path.join(results_folder, 'data_with_new_names.csv')}")
                        except Exception as e:
                            logging.error(f"An error occurred during renaming and feature qc process: {e}")
                            return
                        print("ready_for_norm!")
                        #imputation and mormalization segment
                        try:
                            print("imputing data for clean normalization")
                            #initialize simpleimputer for missing values (NaN) - median!
                            imputer = SimpleImputer(strategy='median')
                            #impute missing values (NaN)
                            df_imputed = pd.DataFrame(imputer.fit_transform(out_dat), columns=out_dat.columns)
                            print("Imputation of missing values completed.")
                            #replace zero values with NaN to use the same imputer
                            df_imputed.replace(0, np.nan, inplace=True)
                            logging.debug("Replaced zero values with NaN for further imputation.")
                            #impute zero values (now NaNs) using the same strategy
                            df_final_imputed = pd.DataFrame(imputer.fit_transform(df_imputed), columns=out_dat.columns)
                            df_final_imputed.index = out_dat.index
                            print("Imputation of zero values completed.")
                            #save imputed data
                            final_imputed_path = os.path.join(results_folder, "df_final_imputed.csv")
                            df_final_imputed.to_csv(final_imputed_path)
                            print(f"Imputed data saved to {final_imputed_path}")
                            #normalization
                            print("Proceeding to normalization...")
                            selected_data = df_final_imputed.copy()
                            # Create a directory for normalization results
                            normalization_dir = os.path.join(results_folder, "normalization")
                            os.makedirs(normalization_dir, exist_ok=True)
                            print(f"Normalization results will be saved to: {normalization_dir}")
                            #split selected_data based on PC column - error handling for missing 'PC' index
                            if 'PC' not in selected_data.index.names:
                                raise ValueError("'PC' index not found in selected data. Ensure 'PC' is part of the index.")
                            pc_groups = selected_data.groupby(level='PC', group_keys=False)
                            df_reverted = pc_groups.apply(lambda x: x)
                            data_3_path = os.path.join(results_folder, "df_proc_3.csv")
                            df_reverted.to_csv(data_3_path)
                            print(f"Final processed data saved to {data_3_path}")
                            #running over each group and performing normalization
                            normalization_methods = {
                                "min-max_normalization": lambda x: (x - x.min()) / (x.max() - x.min()),
                                "central_log_normalization": lambda x: np.log(x + np.sqrt(x ** 2 + 1)),
                                "z-score_normalization": lambda x: (x - x.mean()) / x.std()
                            }
                            for pc_value, pc_group in pc_groups:
                                pc_dir = os.path.join(normalization_dir, f"df_{pc_value}")
                                os.makedirs(pc_dir, exist_ok=True)
                                print(f"Normalization directory created at: {pc_dir}")
                                #original data
                                original_data_path = os.path.join(pc_dir, "original_data.csv")
                                pc_group.to_csv(original_data_path)
                                print(f"Original data saved to {original_data_path}")
                                #normalization methods and saving results
                                normalized_data_dict = {}
                                for method_name, method_func in normalization_methods.items():
                                    try:
                                        normalized_data = method_func(pc_group)
                                        normalized_data.to_csv(os.path.join(pc_dir, f"{method_name.lower().replace(' ', '_')}.csv"))
                                        normalized_data_dict[method_name] = normalized_data
                                        print(f"{method_name} applied and saved to {method_name.lower().replace(' ', '_')}.csv")
                                    except Exception as e:
                                        logging.error(f"Error applying {method_name} for PC value {pc_value}: {e}")
                                #box-cox transformation - data must be positive
                                try:
                                    transformed_data = pd.DataFrame()
                                    for column in pc_group.columns:
                                        if (pc_group[column] <= 0).any():
                                            raise ValueError(
                                                f"Box-Cox transformation requires positive data, found non-positive values in {column}.")
                                        transformed_column, _ = boxcox(pc_group[column])
                                        transformed_data[column] = transformed_column
                                    transformed_data.index = pc_group.index
                                    transformed_data.to_csv(os.path.join(pc_dir, "box_cox_normalized_data.csv"))
                                    print(f"Box-Cox normalized data saved to box_cox_normalized_data.csv")
                                except Exception as e:
                                    logging.error(f"Error applying Box-Cox transformation for PC value {pc_value}: {e}")
                                #log normalization
                                try:
                                    log_data = pd.DataFrame()
                                    for column in pc_group.columns:
                                        #handle negative or zero values for log normalization
                                        if (pc_group[column] <= 0).any():
                                            raise ValueError(
                                                f"Log normalization requires positive data, found non-positive values in {column}.")
                                        log_column = np.log(pc_group[column])
                                        log_data[column] = log_column
                                    log_data.index = pc_group.index
                                    log_data.to_csv(os.path.join(pc_dir, "log_normalized_data.csv"))
                                    print(f"Log normalized data saved to log_normalized_data.csv")
                                except Exception as e:
                                    logging.error(f"Error applying log normalization for PC value {pc_value}: {e}")
                                #QQ plots and histograms for each feature in each normalized data frame
                                normalized_dfs = [
                                    ("Original Data", pc_group),
                                    ("min-max_normalization", normalized_data_dict.get("min-max_normalization")),
                                    ("central_log_normalization", normalized_data_dict.get("central_log_normalization")),
                                    ("z-score_normalization", normalized_data_dict.get("z-score_normalization")),
                                    ("Box_Cox Normalized Data", transformed_data),
                                    ("Log Normalized Data", log_data)

                                ]
                                for df_name, df in normalized_dfs:
                                    if df is None:
                                        logging.warning(f"{df_name} is not available for PC value {pc_value}. Skipping plots.")
                                        continue
                                    df_dir = os.path.join(pc_dir, df_name.replace(' ', '_').lower())
                                    os.makedirs(df_dir, exist_ok=True)
                                    print(f"Plot directory created at: {df_dir}")
                                    for feature in df.columns:
                                        try:
                                            ## qq
                                            plt.figure()
                                            scipy.stats.probplot(df[feature], dist="norm", plot=plt)
                                            plt.title(f"Q-Q-{feature}")
                                            qq_plot_path_2 = os.path.join(df_dir, f"qq_2_plot_{feature}.png")
                                            plt.savefig(qq_plot_path_2)
                                            plt.close()
                                            plt.clf()
                                        except Exception as e:
                                            logging.error(f"Error plotting Q-Q for {feature} in {df_name}: {e}")
                                        # histogram
                                        try:
                                            plt.figure(figsize=(18, 12))
                                            sns.histplot(df[feature], bins=20, kde=True, color='skyblue')
                                            plt.title(f"Histogram for {feature} ({df_name})")
                                            plt.xlabel("Value")
                                            plt.ylabel("Frequency")
                                            plt.tight_layout()
                                            histogram_path = os.path.join(df_dir, f"histogram_{feature}.png")
                                            plt.savefig(histogram_path)
                                            plt.close()
                                            print(f"Histogram saved to: {histogram_path}")
                                        except Exception as e:
                                            logging.error(f"Error plotting histogram for {feature} in {df_name}: {e}")
                            #user to select the normalized data frame to continue with
                            try:
                                normalized_options = ["min-max_normalization", "central_log_normalization",
                                                      "z-score_normalization",
                                                      "Box_Cox Normalized Data", "Log Normalized Data"]
                                print("Available normalization methods:")
                                for idx, option in enumerate(normalized_options):
                                    print(f"{idx + 1}: {option}")
                                selected_normalization = int(
                                    input("Check normalized data and enter the number corresponding to the desired normalized data frame: ")) - 1
                                if selected_normalization < 0 or selected_normalization >= len(normalized_options):
                                    raise ValueError("Selection out of range.")
                                selected_normalization_name = normalized_options[selected_normalization]
                                print(f"Selected normalization method: {selected_normalization_name}")
                                #combine all PC values for the selected normalization into one final normalized data frame
                                final_normalized_data = pd.DataFrame()
                                for pc_value, pc_group in pc_groups:
                                    pc_dir = os.path.join(normalization_dir, f"df_{pc_value}")
                                    file_path = os.path.join(pc_dir, f"{selected_normalization_name.lower().replace(' ', '_')}.csv")
                                    if not os.path.exists(file_path):
                                        logging.warning(f"Normalized file not found: {file_path}. Skipping.")
                                        continue
                                    normalized_df = pd.read_csv(file_path, index_col=0)
                                    final_normalized_data = pd.concat([final_normalized_data, normalized_df])
                                    logging.debug(f"Data from {file_path} added to final normalized data.")
                                #save final normalized data
                                final_normalized_data_path = os.path.join(results_folder, "final_normalized_data.csv")
                                final_normalized_data.to_csv(final_normalized_data_path)
                                print(f"Final normalized data saved to: {final_normalized_data_path}")
                            except Exception as e:
                                logging.error(f"An error occurred during the imputation or normalization process: {e}")
                                return
                        except Exception as e:
                            logging.error(f"An unexpected error occurred: {e}")
                            return
                        finally:
                            print('All processing steps completed.')
                        #dashboard HTML generation
                        #use relative path for links
                        print("creating dashboard for results summary")
                        dashboard_page = os.path.join(main_results_dir, f'{item_name}_dashboard.html')
                        dashboard_content = f"""
                        <html>
                        <head><title>Dashboard for Folder {item_name}</title></head>
                        <body>
                            <h1>Results for Folder: {item_name}</h1>
                            <p>Results stored in: {results_folder}</p>
                        """
                        #add a section for each  file in the folder
                        dashboard_content += "<h2>Processed Plates</h2>"
                        for file_name in os.listdir(item_path):
                            if file_name.endswith(".csv"):
                                results_folder = os.path.join(item_path, f"results_{file_name}_res")
                                barplot_dir = os.path.join(results_folder, "box_plots")
                                summary_dir = os.path.join(results_folder, "general_results")
                                barplot_relative_path = os.path.relpath(barplot_dir, main_results_dir)
                                summary_relative_path = os.path.relpath(summary_dir, main_results_dir)
                                dashboard_content += f"<h3>Results for Plate: {file_name}</h3>"
                                dashboard_content += "<ul>"
                                #bar plots for important features
                                for feature in important_features:
                                    bar_plot_html = os.path.join(barplot_relative_path, f'{feature}_bar_plot.html')
                                    dashboard_content += f'<li><a href="{bar_plot_html}">Bar Plot for {feature}</a></li>'
                                #PCA plot
                                pca_plot_html_id = os.path.join(summary_relative_path, f'pca_plot_with_id_{file_name[:-4]}.html')
                                dashboard_content += f'<li><a href="{pca_plot_html_id}">PCA Plot for {file_name}</a></li>'
                                #PCA plot2
                                pca_plot_html_reg = os.path.join(summary_relative_path, f'pca_plot_wo_id_{file_name[:-4]}.html')
                                dashboard_content += f'<li><a href="{pca_plot_html_reg}">PCA Plot for {file_name}</a></li>'
                                #number of cells plot for the current plate
                                count_plot = os.path.join(summary_relative_path,
                                                                f'number_of_cells_bar_plot.pdf')
                                dashboard_content += f'<li><a href="{count_plot}">number of cells Plot for {file_name}</a></li>'
                                #sum file
                                sum_file = os.path.join(summary_relative_path,
                                                                f'init_statistical_summary.xlsx')
                                dashboard_content += f'<li><a href="{sum_file}">stat_summary for {file_name}</a></li>'
                                #heatmap
                                heat_plot_pdf_reg = os.path.join(summary_relative_path,
                                                                 f'heatmap.pdf')
                                dashboard_content += f'<li><a href="{heat_plot_pdf_reg}">heatmap for {file_name}</a></li>'

                                dashboard_content += "</ul>"

                        #closing the HTML content
                        dashboard_content += "</body></html>"

                        #updated dashboard content to the HTML file
                        with open(dashboard_page, 'w') as f:
                            f.write(dashboard_content)
                        print(f"Dashboard updated with  plots and saved for {item_name} at {dashboard_page}")
                    print(f"{file_name} proccessed succesfully")
    except Exception as e:
        logging.error(f"failed proccessing {plate_df}")
#####################
def gather_user_lists():
    print("Choose appropriate columns for your analysis!")
    #options based on loaded data
    print("usual unwanted columns: Number of Analyzed Fields, Height, Plane, Time, Timepoint, Cell Count, Target CO2, CO2, Target Temperature, Temperature")
    print("enter unwanted columns:")
    list1 = input("enter unwanted columns:(comma-separated): ").split(',')
    print("usual important columns: spots_chanel_3_final - Corrected Spot Intensity - Mean per Well, spots_chanel_3_final - Region Intensity - Mean per Well, spots_chanel_3_final - Spot to Region Intensity - Mean per Well, ir_chanel_3_at - Intensity ir_chanel_3_at chanel_3 Mean - Mean per Well, all_cells - chanel_3 Gabor Max 2 px w2 - Mean per Well, all_cells - chanel_3 SER Ridge 1 px - Mean per Well, all_cells - chanel_3_intensity Mean - Mean per Well, all_cells - chanel_2 Area, all_cells - chanel_1 Area, chan3_intensitymean_div_chan2area, spotschan3_intensitymean_div_chan2area")
    print("enter important_features:")
    list2 = input("enter important_features:(comma-separated): ").split(',')
    print("must index columns!!!: PC, cell_type, cell_id, compound, concentration")
    print("enter index_columns:")
    list3 = input("enter index_columns: (comma-separated): ").split(',')
    #list of strings
    list1 = [item.strip() for item in list1]
    list2 = [item.strip() for item in list2]
    list3 = [item.strip() for item in list3]
    return list1, list2, list3
#RUN MAIN
def main(main_folder, main_folder_path, final_folder_path):
    print(Fore.GREEN + "Starting ROS analysis for combined file of all plates in folder in your main directory:)")
    print(Fore.GREEN + "STEP 1 - processing text files...")
    process_folders(main_folder)
    prompt_for_review()
    print(Fore.GREEN + "STEP 2 - combining csv files...")
    process_data_combine(main_folder_path)
    prompt_for_review()
    print(Fore.GREEN + "STEP 3 - choosing column names for analysis...")
    print(Fore.RED + "==== Please make sure you choose correct column names! ====")
    list1, list2, list3 = gather_user_lists()
    process_plate_data_on_folders(final_folder_path, list1, list2, list3)
    print("Completed combined plates IF analysis, good luck with statistics:)")
    print(Fore.RED + Style.BRIGHT + "==== IMPORTANT MESSAGE ====" +
          "PLEASE REMOVE YOUR FILES FROM THE MAIN DIRECTORY ====")
#########
if __name__ == "__main__":
    main_folder = input("Enter folder name(containing a folder with .txt files): ")
    main_folder_path = input("Enter: the main folder name/results ")
    final_folder_path = input("Enter: the main folder name/ROS_result_combined ")
    main(main_folder, main_folder_path, final_folder_path)