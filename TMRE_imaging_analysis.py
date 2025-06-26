### TMRE- Combined Plates
## takes as input a folder that includes an additional folder inside it with txt files
## there should be numeric features in these files & index columns
## must index columns are - PC, cell_type, cell_id, compound, concentration
## follow the steps until a final normalized data frame ready for statistics will be generated.
## all output files will be saved in the main folder including all steps in the analysis, in addition a dashboard with the main results will be saved.
## all text files will be merged and analyzed.
##libraries
import itertools
import logging
import pandas as pd
import numpy as np
from sklearn.decomposition import PCA
import matplotlib.pyplot as plt
import os
import glob
import matplotlib.pyplot as plt
import seaborn as sns
import os
import pandas as pd
import matplotlib.pyplot as plt
from itertools import combinations
import os
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from scipy.stats import zscore
from sklearn.decomposition import PCA
from sklearn.discriminant_analysis import LinearDiscriminantAnalysis as LDA
from sklearn.cluster import KMeans
from scipy.cluster.hierarchy import dendrogram, linkage
import plotly.express as px
from sklearn.decomposition import PCA
from sklearn.feature_selection import VarianceThreshold
from sklearn.neighbors import LocalOutlierFactor
from sklearn.preprocessing import LabelEncoder
from sklearn.preprocessing import StandardScaler
import plotly.express as px
from sklearn.impute import SimpleImputer
import plotly
from plotly.subplots import make_subplots
import matplotlib.patches as mpatches
import plotly.graph_objects as go
import kaleido
import matplotlib.patches as mpatches
from sklearn.feature_selection import VarianceThreshold
from sklearn.model_selection import cross_val_score
from sklearn.ensemble import RandomForestRegressor
from scipy.stats import spearmanr
from scipy.stats import zscore
import scipy
import scipy.stats as stats
import statsmodels.api as sm
import statsmodels.graphics.gofplots as smqq
from scipy.stats import boxcox
from pptx import Presentation
from pptx.util import Inches
from openpyxl import Workbook
from openpyxl.utils.dataframe import dataframe_to_rows
from sklearn.manifold import TSNE
#list of index columns
# index_columns = ['experiment', 'plate', 'measurment', 'row', 'column', 'field', 'case', 'id', 'compound', 'concentration',
#                  'cell_count_seed', 'incubation_condition', 'disease_stage', 'cell_type', 'passage', 'biopsy',
#                  'age', 'gender', 'race']

#index_columns = ['ep', 'position', 'ID']
##todo fix drop
## todo add more runs for statistical sum
## todo finish with multiple batch column choosing .
## todo change index column- cell count to cell seed..
## todo more normalization methodds
## todo remove indexes from outlier methods
# def process_csv(file_path):
#     # Step 1: Read the CSV file
#     df = pd.read_csv(file_path)
#
#     # Step 2: Print starting message
#     print("starting field proc")
#
#     # Step 3: Open or create a results folder
#     results_folder = 'results_field_GSD3'
#     if not os.path.exists(results_folder):
#         os.makedirs(results_folder)
#
#     # Step 4: Prompt the user to drop columns
#     print("Columns in the CSV file:")
#     print(df.columns)
#     print("Please review index columns and remove redundant columns")
#     columns_to_drop = input("Enter columns to drop separated by commas (or press Enter to skip): ")
#     #columns_to_drop = [col.strip() for col in columns_to_drop].split(',')
#     if columns_to_drop:
#         columns_to_drop = [col.strip() for col in columns_to_drop.split(',')]
#         df.drop(columns=columns_to_drop, inplace=True)
#         index_columns[:] = [col for col in index_columns if col not in columns_to_drop]
#
#     # Step 5: Create new columns based on all possible combinations of the index columns
#     new_index_columns = []
#     for r in range(2, len(index_columns) + 1):
#         for combo in combinations(index_columns, r):
#             col_name = '_'.join(combo)
#             df[col_name] = df[list(combo)].apply(lambda row: '_'.join(row.values.astype(str)), axis=1)
#             new_index_columns.append(col_name)
#
#     # Print indexing complete
#     print("indexing complete")
#
#     # Update index columns list
#     index_columns.extend(new_index_columns)
#
#     # Step 6: Save the DataFrame
#     processed_file_path = os.path.join(results_folder, 'processed_data.csv')
#     df.to_csv(processed_file_path, index=False)
#
#     # Step 7: Prompt the user to choose one of the index columns
#     print("Index columns available:")
#     print(index_columns)
#
#     chosen_index_col = input("Choose one of the index columns to create a statistical summary: ")
#
#     if chosen_index_col not in index_columns:
#         print(f"Invalid choice. Please choose from {index_columns}.")
#         return
#
#     # Step 8: Create a statistical summary for all numeric features based on the chosen index column
#     numeric_cols = df.select_dtypes(include='number').columns
#     summary = df.groupby(chosen_index_col)[numeric_cols].describe()
#
#     # Step 9: Save the statistical summary
#     summary_file_path = os.path.join(results_folder, 'statistical_summary.csv')
#     summary.to_csv(summary_file_path)
#
#     print("Statistical summary saved to", summary_file_path)
#
#     # Step 10: Plot a histogram of the "cell_count" feature and save it
#     if 'cell_count' in df.columns:
#         plt.figure()
#         df['cell_count'].hist(bins=30)
#         plt.title('Histogram of Cell Count')
#         plt.xlabel('Cell Count')
#         plt.ylabel('Frequency')
#         histogram_file_path = os.path.join(results_folder, 'cell_count_histogram.png')
#         plt.savefig(histogram_file_path)
#         plt.close()
#
#         # Step 11: Save a text file with statistics of the "cell_count" feature
#         cell_count_stats = df['cell_count'].describe()
#         stats_file_path = os.path.join(results_folder, 'cell_count_stats.txt')
#         with open(stats_file_path, 'w') as f:
#             f.write(cell_count_stats.to_string())
#
#         print("Cell count histogram and statistics saved.")
#
#         # Step 12: Prompt the user to choose a minimum and maximum value for "cell_count"
#         min_value = float(input("Enter the minimum value for cell_count: "))
#         max_value = float(input("Enter the maximum value for cell_count: "))
#
#         # Step 13: Filter the DataFrame based on the specified minimum and maximum values
#         filtered_df = df[(df['cell_count'] >= min_value) & (df['cell_count'] <= max_value)]
#
#         # Step 14: Save the filtered DataFrame
#         filtered_file_path = os.path.join(results_folder, 'filtered_data.csv')
#         filtered_df.to_csv(filtered_file_path, index=False)
#
#         print("Filtered DataFrame saved to", filtered_file_path)
#     else:
#         print("Column 'cell_count' not found in the DataFrame.")
#         filtered_df = df.copy()  # Use the original DataFrame if 'cell_count' is not found
#
#     # Step 15: Create a new column called 'field_count' and assign 1 to each row
#     filtered_df['field_count'] = 1
#     print("field count initiated")
#
#     # Define the outlier detection functions
#     def compute_percentiles(df, output_file):
#         up_bounds = []
#         low_bounds = []
#         above_count = []
#         below_count = []
#         numeric_cols = df.select_dtypes(include='number').columns
#         for col in numeric_cols:
#             percentile_low = df[col].quantile(0.01)
#             percentile_high = df[col].quantile(0.99)
#             up_bound = percentile_high
#             low_bound = percentile_low
#             above = df[df[col] > percentile_high]
#             below = df[df[col] < percentile_low]
#             above_count.append(len(above))
#             below_count.append(len(below))
#             df = df[(df[col] >= percentile_low) & (df[col] <= percentile_high)]
#             up_bounds.append(up_bound)
#             low_bounds.append(low_bound)
#         output = pd.DataFrame(
#             {'Column': numeric_cols, 'Up Bound': up_bounds, 'Low Bound': low_bounds, 'Above Count': above_count,
#              'Below Count': below_count})
#         output.to_csv(output_file, index=False)
#         return df
#
#     def compute_iqr(df, output_file):
#         up_bounds = []
#         low_bounds = []
#         above_count = []
#         below_count = []
#         numeric_cols = df.select_dtypes(include='number').columns
#         for col in numeric_cols:
#             percentile25th = df[col].quantile(0.25)
#             percentile75th = df[col].quantile(0.75)
#             iqr = percentile75th - percentile25th
#             up_bound = percentile75th + 1.5 * iqr
#             low_bound = percentile25th - 1.5 * iqr
#             above = df[df[col] > up_bound]
#             below = df[df[col] < low_bound]
#             above_count.append(len(above))
#             below_count.append(len(below))
#             df = df[(df[col] >= low_bound) & (df[col] <= up_bound)]
#             up_bounds.append(up_bound)
#             low_bounds.append(low_bound)
#         output = pd.DataFrame(
#             {'Column': numeric_cols, 'Up Bound': up_bounds, 'Low Bound': low_bounds, 'Above Count': above_count,
#              'Below Count': below_count})
#         output.to_csv(output_file, index=False)
#         return df
#
#     def compute_scores(df, output_file):
#         up_scores = []
#         low_scores = []
#         above_count = []
#         below_count = []
#         numeric_cols = df.select_dtypes(include='number').columns
#         for col in numeric_cols:
#             up_score = df[col].mean() + 3 * df[col].std()
#             low_score = df[col].mean() - 3 * df[col].std()
#             above = df[df[col] > up_score]
#             below = df[df[col] < low_score]
#             above_count.append(len(above))
#             below_count.append(len(below))
#             df = df[(df[col] >= low_score) & (df[col] <= up_score)]
#             up_scores.append(up_score)
#             low_scores.append(low_score)
#         output = pd.DataFrame(
#             {'Column': numeric_cols, 'Up Score': up_scores, 'Low Score': low_scores, 'Above Count': above_count,
#              'Below Count': below_count})
#         output.to_csv(output_file, index=False)
#         return df
#
#     # Apply the outlier detection functions on filtered_df
#     df_percentiles = compute_percentiles(filtered_df, os.path.join(results_folder, 'percentiles_outliers.csv'))
#     df_iqr = compute_iqr(filtered_df, os.path.join(results_folder, 'iqr_outliers.csv'))
#     df_scores = compute_scores(filtered_df, os.path.join(results_folder, 'scores_outliers.csv'))
#
#     print("Please review field outlier detection and choose one method to continue with:")
#     method_choice = input("Enter 1 for percentiles, 2 for IQR, 3 for scores: ")
#
#     if method_choice == '1':
#         final_df = df_percentiles
#     elif method_choice == '2':
#         final_df = df_iqr
#     elif method_choice == '3':
#         final_df = df_scores
#     else:
#         print("Invalid choice. Defaulting to original DataFrame.")
#         final_df = filtered_df
#
#     # Save the final filtered DataFrame
#     final_file_path = os.path.join(results_folder, 'final_filtered_data.csv')
#     final_df.to_csv(final_file_path, index=False)
#
#     print("Final filtered DataFrame saved to", final_file_path)
#
#     # Step 16: Prompt the user to choose a main index column to group by and other index columns to keep
#     print("Index columns available for grouping:")
#     print(index_columns)
#
#     main_index_col = input("Choose the main index column to group by-well: ").strip()
#     if main_index_col not in index_columns:
#         print(f"Invalid choice. Please choose from {index_columns}.")
#         return
#
#     #additional_index_cols = input("Enter additional index columns to keep, separated by commas: ")
#     #if additional_index_cols:
#     #    additional_index_cols = [col.strip() for col in additional_index_cols.split(',')]
#     #else:
#     #    additional_index_cols = []
#
#     # Step 17: Group the DataFrame by the main index column and additional index columns, aggregating numeric features by their sum
#     #Group by user-specified index column
#     #index_col = input("Enter the index column for grouping: ").strip()
#     numeric_cols = final_df.select_dtypes(include=np.number).columns
#     agg_dict = {col: 'sum' for col in numeric_cols}
#     agg_dict.update({col: 'first' for col in index_columns if col != main_index_col})
#     aggregated_df = final_df.groupby(main_index_col).agg(agg_dict)
#     aggregated_df = aggregated_df.reset_index()
#     #groupby_cols = [main_index_col] + additional_index_cols
#     #aggregated_df = final_df.groupby(groupby_cols).sum().reset_index()
#     print("well_res_created")
#     # Step 18: Create a new results folder called "by well results"
#     by_well_results_folder = os.path.join(results_folder, 'by_well_results')
#     if not os.path.exists(by_well_results_folder):
#         os.makedirs(by_well_results_folder)
#
#     print("New folder 'by well results' created.")
#     filtered_well_file_path = os.path.join(by_well_results_folder, 'aggregated_df_data.csv')
#     aggregated_df.to_csv(filtered_well_file_path, index=False)
#
#     # Step 19: Perform the same process for the "field_count" column as previously done for "cell_count"
#     if 'field_count' in aggregated_df.columns:
#         plt.figure()
#         aggregated_df['field_count'].hist(bins=30)
#         plt.title('Histogram of Field Count')
#         plt.xlabel('Field Count')
#         plt.ylabel('Frequency')
#         field_count_histogram_path = os.path.join(by_well_results_folder, 'field_count_histogram.png')
#         plt.savefig(field_count_histogram_path)
#         plt.close()
#
#         # Save a text file with statistics of the "field_count" column
#         field_count_stats = aggregated_df['field_count'].describe()
#         field_count_stats_path = os.path.join(by_well_results_folder, 'field_count_stats.txt')
#         with open(field_count_stats_path, 'w') as f:
#             f.write(field_count_stats.to_string())
#
#         print("Field count histogram and statistics saved.")
#
#         # Prompt the user to choose a minimum and maximum value for "field_count"
#         min_field_count_value = float(input("Enter the minimum value for field_count: "))
#         max_field_count_value = float(input("Enter the maximum value for field_count: "))
#
#         # Filter the DataFrame based on the specified minimum and maximum values
#         filtered_field_df = aggregated_df[(aggregated_df['field_count'] >= min_field_count_value) & (
#                     aggregated_df['field_count'] <= max_field_count_value)]
#
#         # Save the filtered DataFrame
#         filtered_field_file_path = os.path.join(by_well_results_folder, 'filtered_field_data.csv')
#         filtered_field_df.to_csv(filtered_field_file_path, index=False)
#
#         print("Filtered field DataFrame saved to", filtered_field_file_path)
#     else:
#         print("Column 'field_count' not found in the aggregated DataFrame.")
#
#     batch_index_col = input("Choose the batch index column to group by: ").strip()
#     if batch_index_col not in index_columns:
#         print(f"Invalid choice. Please choose from {index_columns}.")
#         return
#
#     numeric_cols = filtered_field_df.select_dtypes(include=np.number).columns
#     agg_dict = {col: 'mean' for col in numeric_cols}
#     agg_dict.update({col: 'first' for col in index_columns if col != batch_index_col})
#     batch_df = filtered_field_df.groupby(batch_index_col).agg(agg_dict)
#     batch_df = batch_df.reset_index()
#     # Create meta_data
#     meta_data_columns = index_columns
#     meta_data = batch_df[meta_data_columns].drop_duplicates()
#     print(meta_data)
#     # meta_data.to_csv(os.path.join(result_batch_dir, "meta_data.csv"), index=False)
#     #for cols in index_columns:
#     #    col_name = '_'.join(cols)
#     #    meta_data[col_name] = meta_data[list(cols)].apply(lambda x: '_'.join(map(str, x)), axis=1)
#     meta_data.to_csv(os.path.join(by_well_results_folder, "meta_data.csv"), index=False)
#     # Ask for index column to keep in feat_data
#     # index_to_keep = input("Enter index column to keep in feat_data: ")
#     if batch_index_col in batch_df.columns:
#         # Create feat_data
#         feat_data = batch_df.set_index(batch_index_col)
#         # Drop indexes
#         feat_data = feat_data.drop(
#             columns=[col_name for col_name in index_columns if col_name not in [batch_index_col]])
#         # Save feat_data
#         feat_data.to_csv(os.path.join(by_well_results_folder, "feat_data.csv"))
#     else:
#         print(f"Index column '{main_index_col}' not found in DataFrame columns.")
#
#     return aggregated_df
#
# # run-IRF
# #file_path = 'IRF/field_data_final_true.csv'
# #process_csv(file_path)
# #print("Processed DataFrame:")
# # run-GSD3
# file_path = 'GSD3/final_field_data_gsd3.csv'
# process_csv(file_path)
# print("Processed DataFrame:")
#print(df.head())

### next segment


## todo fix inc time problem, change to int
## todo decide on feature selcetion
## todo if batch what to do?
#### post batch well proc
# def process_data(csv_file_path):
#     #  folder path and filename
#     folder_path, filename = os.path.split(csv_file_path)
#
#     #results folder path
#     results_folder_path = os.path.join(folder_path, "results_proc_try")
#
#     # Readfile
#     data = pd.read_csv(csv_file_path)
#
#     #make results folder
#     os.makedirs(results_folder_path, exist_ok=True)
#
#     #edit column names
#     data.columns = [col.split("[")[0].strip() if "[" in col else col for col in data.columns]
#     data.to_csv(os.path.join(results_folder_path, "tool_processed_data_1.csv"))
#     #columns
#     print("Columns in the dataset:")
#     for idx, column in enumerate(data.columns):
#         print(f"{idx}: {column}")
#
#     #user to specify columns to drop
#     columns_to_drop = input("Enter columns to drop (comma-separated): ").split(',')
#     # Remove any leading/trailing whitespace from column names
#     columns_to_drop = [col.strip() for col in columns_to_drop]
#
#     #user to specify columns to set as indexes
#     ## add PC as must.
#     columns_to_set_index = input("Enter columns to set as indexes (comma-separated): ").split(',')
#     # Remove any leading/trailing whitespace from column names
#     columns_to_set_index = [col.strip() for col in columns_to_set_index]
#
#     # Drop specified columns
#     if columns_to_drop:
#         data = data.drop(columns=columns_to_drop)
#
#     # Set specified columns as indexes
#     if columns_to_set_index:
#         data = data.set_index(columns_to_set_index, drop = False)
#
#     # Save current df
#     data.to_csv(os.path.join(results_folder_path, "tool_processed_data_2.csv"))
#
#     # Generate histogram for cell count feature
#     fig = px.histogram(data, x="all_cell_count", title="Histogram of Cell Count")
#     histogram_file_path = os.path.join(results_folder_path, "histogram.html")
#     fig.write_html(histogram_file_path)
#
#     # Create statistics for "cell count" feature
#     stats = data["all_cell_count"].describe()
#     stats_file_path = os.path.join(results_folder_path, "cell_count-statistics.txt")
#     with open(stats_file_path, "w") as f:
#         f.write(stats.to_string())
#     ## add summary saved bla nla
#     # Ask user for minimum and maximum values
#     min_value = float(input("Enter the minimum value for cell count: "))
#     max_value = float(input("Enter the maximum value for cell count: "))
#
#     # Remove rows with cell count outside the specified range
#     data = data[(data["all_cell_count"] >= min_value) & (data["all_cell_count"] <= max_value)]
#
#     # Save DataFrame in results folder
#     data.to_csv(os.path.join(results_folder_path, "tool_processed_data_3.csv"))
#
#     # Prompt user to choose index column for outlier detection
#     print("Columns available for outlier detection:")
#     print(data.columns)
#     index_column = input("Choose one column as index for outlier detection: ")
#
#     # Outlier detection using Local Outlier Factor
#     types = data[index_column].unique()
#     print(types)
#     mask = []
#     features = data.select_dtypes(include='number').columns
#
#     detector_list = [
#         ("Local Outlier Factor 30", LocalOutlierFactor(n_neighbors=30))
#     ]
#
#     for name, algorithm in detector_list:
#         errors = np.full(len(data),fill_value=np.nan)
#         outliers = np.full(len(data),fill_value=np.nan)
#
#         for type in types:
#             x = data.loc[:,features].values
#             F = x.sum(1)
#             mask = np.zeros(x.shape[0])
#             mask[np.isfinite(F)] = 1
#             mask_type = mask * np.array(data[index_column] == type)
#             Curr_df = data.loc[mask_type == 1, features]
#             x = Curr_df.values
#
#
#             if name == "Local Outlier Factor 30":
#                 algorithm.fit(x)
#                 errors[mask_type==1]  = algorithm.negative_outlier_factor_
#                 outliers[mask_type==1] = algorithm.fit_predict(x)
#
#         data[name] = errors
#         data[f'{name}_outliers'] = outliers
#         data.set_index(name, inplace=True,
#                                                append=True, drop=False)
#         data.to_csv(os.path.join(results_folder_path, "tool_processed_data_outliers_1.csv"))
#     # Exclude rows that were defined as outliers
#         for col in data.columns:
#             if col.endswith("_outliers"):
#                 data = data[data[col] != -1]
#
#
#         data.to_csv(os.path.join(results_folder_path, "tool_processed_data_outliers_2.csv"))
#         # Apply feature selection for each subgroup
#
#         # Define subgroups and their features (replace with your actual subgroup feature lists)
#         # subgroups = {
#         #     'df_network': ['network_count', 'network_total_area'],
#         #     'df_rod': ['rod_count', 'rod_total_area'],
#         #     'df_rounded': ['rounded_count', 'rounded_total_area']
#         #
#         #
#         # }
#         #
#         #
#         # # Calculate variance of each feature across the entire dataset
#         # #features = data.select_dtypes(include='number').columns
#         # #variance_across_dataset = features.var()
#         # # Select top features for each subgroup
#         # def select_top_features(data, subgroups):
#         #     top_features = {}
#         #     for subgroup_name, subgroup_features in subgroups.items():
#         #         # Ensure that we only include features that exist in the dataframe
#         #         existing_features = [feature for feature in subgroup_features if feature in data.columns]
#         #         #
#         #         if not existing_features:
#         #             print(f"No valid features found for subgroup: {subgroup_name}")
#         #             continue
#         #         subgroup_data = data[existing_features]
#         #         # Drop features with all missing values
#         #         subgroup_data = subgroup_data.dropna(axis=1, how='all')
#         #         # Fill remaining missing values with column means
#         #         subgroup_data = subgroup_data.fillna(subgroup_data.mean())
#         #         variance_within_subgroup = subgroup_data.var()
#         #         top_feature = variance_within_subgroup.idxmax()
#         #         top_features[subgroup_name] = top_feature
#         #     return top_features
#         #
#         # selected_features = select_top_features(data, subgroups)
#         #
#         # # Print selected top features and original features of each subgroup
#         # print("Selected top features:")
#         # for subgroup_name, top_feature in selected_features.items():
#         #     print(f"Top feature for {subgroup_name}: {top_feature}")
#         #     print(f"Original features for {subgroup_name}: {subgroups[subgroup_name]}")
#         #
#         # # Allow the user to decide if they want to swap features
#         # for subgroup_name, top_feature in selected_features.items():
#         #     swap_choice = input(f"Do you want to swap features for {subgroup_name}? (yes/no): ")
#         #     if swap_choice.lower() == 'yes':
#         #         selected_feature = input(f"Enter the selected feature to replace '{top_feature}' for {subgroup_name}: ")
#         #         if selected_feature in subgroups[subgroup_name]:
#         #             selected_features[subgroup_name] = selected_feature
#         #             print(f"Feature '{selected_feature}' replaced '{top_feature}' for {subgroup_name}")
#         #
#         # # Update data based on the selected top features
#         # for subgroup_name, top_feature in selected_features.items():
#         #     if top_feature in subgroups[subgroup_name]:
#         #         data[top_feature] = data.pop(top_feature)
#         #         # If the selected feature is not already in the DataFrame,
#         #         # we should add it.
#         #         if selected_features[subgroup_name] not in data.columns:
#         #             data[selected_features[subgroup_name]] = data[
#         #                 subgroups[subgroup_name][subgroups[subgroup_name].index(selected_features[subgroup_name])]]
#         #         else:
#         #             data[selected_features[subgroup_name]] = data.pop(
#         #                 subgroups[subgroup_name][subgroups[subgroup_name].index(selected_features[subgroup_name])])
#         #
#         #         print(f"Updated {subgroup_name} with selected feature: {selected_features[subgroup_name]}")
#         # # # Remove columns that were not selected or swapped
#         # # for subgroup_name, subgroup_features in subgroups.items():
#         # #     # Get the features that were not selected or swapped
#         # #     unselected_features = [feature for feature in subgroup_features if
#         # #                            feature not in selected_features.values()]
#         # #     # Drop these features from the DataFrame
#         # #     data.drop(unselected_features, axis=1, inplace=True)
#         # # Create a new DataFrame to store the selected features
#         # selected_data = pd.DataFrame()
#         #
#         # # Update data based on the selected top features
#         # for subgroup_name, top_feature in selected_features.items():
#         #     if top_feature in subgroups[subgroup_name]:
#         #         data[top_feature] = data.pop(top_feature)
#         #         # If the selected feature is not already in the DataFrame,
#         #         # we should add it.
#         #         if selected_features[subgroup_name] not in data.columns:
#         #             data[selected_features[subgroup_name]] = data[
#         #                 subgroups[subgroup_name][subgroups[subgroup_name].index(selected_features[subgroup_name])]]
#         #         else:
#         #             data[selected_features[subgroup_name]] = data.pop(
#         #                 subgroups[subgroup_name][subgroups[subgroup_name].index(selected_features[subgroup_name])])
#         #
#         #         print(f"Updated {subgroup_name} with selected feature: {selected_features[subgroup_name]}")
#         #         # Add the selected feature to the selected_data DataFrame
#         #         selected_data[selected_features[subgroup_name]] = data[selected_features[subgroup_name]]
#         #
#         # # Copy the indexes from the original data DataFrame to the selected_data DataFrame
#         # selected_data.index = data.index
#         #
#         # print("Updated data with selected top features:")
#         # print(selected_data)
#
#
#         selected_data = data.copy()
#         ##todo fix cv for feature qc
#         ##todo other vis / meth
#         selected_data.to_csv(os.path.join(results_folder_path, "processed_data_4.csv"))
#         # Assess feature quality
#         # Prompt the user to select the correct index title as the target variable
#         print("Select the correct index title to be the target variable:")
#         for idx, index_title in enumerate(selected_data.index.names):
#             print(f"{idx + 1}: {index_title}")
#
#         selected_index = int(input("Enter the number corresponding to the correct index title: ")) - 1
#
#         # Set the selected index title as the target variable
#         target_variable = selected_data.index.names[selected_index]
#
#         print(f"Target variable selected: {target_variable}")
#         target_variable_labels, _ = pd.factorize(selected_data.index.get_level_values(target_variable))
#         print(target_variable_labels)
#         ##
#         # Create a directory for feature QC results
#         feature_qc_dir = os.path.join(results_folder_path, "feature_qc")
#         os.makedirs(feature_qc_dir, exist_ok=True)
#
#         # 1. Statistical Summary CSV
#         statistical_summary = selected_data.describe()
#         statistical_summary.to_csv(os.path.join(feature_qc_dir, "statistical_summary.csv"))
#
#         # 2. Feature Stability Scores (CV)
#         stability_scores = selected_data.var() / selected_data.mean()
#
#         # 3. Correlation Analysis
#         correlation_matrix = selected_data.corr(method='spearman')
#         correlation_matrix.to_csv(os.path.join(feature_qc_dir, "correlation_matrix.csv"))
#
#         # Plot Correlation Matrix
#         plt.figure(figsize=(15, 12))
#         sns.heatmap(correlation_matrix, annot=True, cmap='coolwarm', fmt=".2f")
#         plt.title("Correlation Matrix")
#         plt.tight_layout()
#         correlation_plot_path = os.path.join(feature_qc_dir, "correlation_matrix_plot.png")
#         plt.savefig(correlation_plot_path)
#         plt.close()
#         print(f"Correlation matrix plot saved to: {correlation_plot_path}")
#
#         # 4. CV Plot (if stability scores are available)
#         if stability_scores is not None:
#             plt.figure(figsize=(12, 6))
#             plt.bar(stability_scores.index, stability_scores.values, color='skyblue')
#             plt.title("Feature Stability Scores (Coefficient of Variation)")
#             plt.xlabel("Feature")
#             plt.ylabel("Stability Score")
#             plt.xticks(rotation=45, ha='right')
#             plt.tight_layout()
#             cv_plot_path = os.path.join(feature_qc_dir, "stability_scores_plot.png")
#             plt.savefig(cv_plot_path)
#             plt.close()
#             print(f"CV plot saved to: {cv_plot_path}")
#         else:
#             print("Stability scores not available. Skipping CV plot.")
#
#         # 5. Random Forest Score computation
#         target_variable_labels, _ = pd.factorize(selected_data.index.get_level_values(target_variable))
#         rf_scores = {}
#         for feature in selected_data.columns:
#             model = RandomForestRegressor()
#             feature_values = selected_data[feature].values.reshape(-1, 1)
#             model.fit(feature_values, target_variable_labels)
#             rf_scores[feature] = model.score(feature_values, target_variable_labels)
#
#         # Save RF scores to CSV
#         rf_scores_df = pd.DataFrame(rf_scores, index=['RF_Score'])
#         rf_scores_df.to_csv(os.path.join(feature_qc_dir, "random_forest_scores.csv"))
#
#         # Plot RF Scores
#         plt.figure(figsize=(10, 6))
#         plt.bar(rf_scores.keys(), rf_scores.values(), color='lightgreen')
#         plt.title("Random Forest Scores for each Feature")
#         plt.xlabel("Feature")
#         plt.ylabel("Random Forest Score")
#         plt.xticks(rotation=45, ha='right')
#         plt.tight_layout()
#         rf_plot_path = os.path.join(feature_qc_dir, "random_forest_scores_plot.png")
#         plt.savefig(rf_plot_path)
#         plt.close()
#         print(f"Random Forest scores plot saved to: {rf_plot_path}")
#         ##
#         # Ask the user if they want to change feature names
#         print("final_proc_pre_norm")
#         rename_columns = input("Do you want to rename any columns? (yes/no): ")
#         if rename_columns.lower() == 'yes':
#             print("Current Columns:")
#             for idx, column in enumerate(selected_data.columns):
#                 print(f"{idx}: {column}")
#
#             renaming = input(
#                 "Enter the column name and new name separated by ':', each pair separated by commas (e.g., col1:new_col1,col2:new_col2): ")
#             rename_dict = dict(pair.split(':') for pair in renaming.split(','))
#             selected_data.rename(columns=rename_dict, inplace=True)
#             print("Names changed successfully:)")
#         ##
#         print("ready_for_norm")
#         # Normalize features using different techniques
#         # Create a directory for normalization results
#         # Create a directory for normalization results
#         # Create a directory for normalization results
#         normalization_dir = os.path.join(results_folder_path, "normalization")
#         os.makedirs(normalization_dir, exist_ok=True)
#         # Split selected_data based on the 'PC' index
#         pc_groups = selected_data.groupby(level='incubation_time', group_keys=False)
#         df_reverted = pc_groups.apply(lambda x: x)
#         df_splited = selected_data.groupby(level='incubation_time', group_keys=False).reset_index()
#         df_reverted.to_csv(os.path.join(results_folder_path, "processed_data_5.csv"))
#         # Iterate over each group and perform normalization
#         # Perform normalization for each group
#         selected_data.to_csv(os.path.join(normalization_dir, "original_data.csv"))
#         # For Min-Max normalization
#         min_max_normalized_data = (selected_data - selected_data.min()) / (selected_data.max() - selected_data.min())
#         min_max_normalized_data.to_csv(os.path.join(normalization_dir, "min_max_normalized_data.csv"))
#
#         # For Central Logarithmic normalization
#         central_log_normalized_data = np.log(selected_data + np.sqrt(selected_data ** 2 + 1))
#         central_log_normalized_data.to_csv(os.path.join(normalization_dir, "central_log_normalized_data.csv"))
#
#         # For Z-score normalization
#         z_score_normalized_data = (selected_data - selected_data.mean()) / selected_data.std()
#         z_score_normalized_data.to_csv(os.path.join(normalization_dir, "z_score_normalized_data.csv"))
#
#         # For Box-Cox transformation
#         transformed_data = pd.DataFrame()
#         for column in selected_data.columns:
#             transformed_column, _ = boxcox(selected_data[column])
#             transformed_data[column] = transformed_column
#         transformed_data.index = selected_data.index
#         transformed_data.to_csv(os.path.join(normalization_dir, "box_cox_normalized_data.csv"))
#
#         # For log normalization
#         log_data = pd.DataFrame()
#         for column in selected_data.columns:
#             log_column = np.log(selected_data[column])
#             log_data[column] = log_column
#         log_data.index = selected_data.index
#         log_data.to_csv(os.path.join(normalization_dir, "log_normalized_data.csv"))
#
#         for df_name, df in [("Original Data", selected_data),
#                                 ("Min-Max Normalized Data", min_max_normalized_data),
#                                 ("Central Log Normalized Data", central_log_normalized_data),
#                                 ("Z-score Normalized Data", z_score_normalized_data),
#                                 ("Box_Cox Normalized Data", transformed_data),
#                                 ("Log Normalized Data", log_data)]:
#
#
#
#
#
#                 # Create a directory for the current data frame
#                 df_dir = os.path.join(normalization_dir, df_name)
#                 os.makedirs(df_dir, exist_ok=True)
#
#                 for feature in df.columns:
#
#                         # Convert Series to array for plotting QQ plot
#                     data_array = (df[feature]).astype(int)
#                         #type(data_array)
#                     print(data_array)
#                     new_array = data_array.to_numpy
#                         #type(new_array)
#                     print(new_array)
#
#                         # QQ plot
#                         #fig = plt.figure(figsize=(8, 6))
#
#                     sm.qqplot(df[feature], line='45')
#                         #stats.probplot(new_array, dist="norm", plot=plt)
#                     plt.title(f"QQ Plot for {feature} ({df_name})")
#                     qq_plot_path = os.path.join(df_dir, f"qq_plot_{feature}.png")
#                     plt.savefig(qq_plot_path)
#
#                     plt.close()
#                     plt.clf()
#
#                         # Histogram
#                     fig = plt.figure(figsize=(8, 6))
#
#                     plt.hist(df[feature], bins=20, color='skyblue', edgecolor='black')
#                     plt.title(f"Histogram for {feature} ({df_name})")
#                     plt.xlabel("Value")
#                     plt.ylabel("Frequency")
#                     plt.tight_layout()
#                     histogram_path = os.path.join(df_dir, f"histogram_{feature}.png")
#                     plt.savefig(histogram_path)
#
#                     plt.close(fig)
#                     plt.clf()
#
#
#         print("Normalization and visualization completed. Results saved in 'normalization' directory.")
#
#
#
#         return data
#         # for pc_value, pc_group in pc_groups:
#         #     # Create a directory for the current PC value
#         #     pc_dir = os.path.join(normalization_dir, f"df_{pc_value}")
#         #     os.makedirs(pc_dir, exist_ok=True)
#         #     # Perform normalization for each group
#         #     pc_group.to_csv(os.path.join(results_folder_path, "original_data.csv"))
#         #     # For Min-Max normalization
#         #     min_max_normalized_data = (pc_group - pc_group.min()) / (pc_group.max() - pc_group.min())
#         #     min_max_normalized_data.to_csv(os.path.join(pc_dir, "min_max_normalized_data.csv"))
#         #
#         #     # For Central Logarithmic normalization
#         #     central_log_normalized_data = np.log(pc_group + np.sqrt(pc_group ** 2 + 1))
#         #     central_log_normalized_data.to_csv(os.path.join(pc_dir, "central_log_normalized_data.csv"))
#         #
#         #     # For Z-score normalization
#         #     z_score_normalized_data = (pc_group - pc_group.mean()) / pc_group.std()
#         #     z_score_normalized_data.to_csv(os.path.join(pc_dir, "z_score_normalized_data.csv"))
#         #
#         #     # For Box-Cox transformation
#         #     transformed_data = pd.DataFrame()
#         #     for column in pc_group.columns:
#         #         transformed_column, _ = boxcox(pc_group[column])
#         #         transformed_data[column] = transformed_column
#         #     transformed_data.index = pc_group.index
#         #     transformed_data.to_csv(os.path.join(pc_dir, "box_cox_normalized_data.csv"))
#         #
#         #     # For log normalization
#         #     log_data = pd.DataFrame()
#         #     for column in pc_group.columns:
#         #         log_column = np.log(pc_group[column])
#         #         log_data[column] = log_column
#         #     log_data.index = pc_group.index
#         #     log_data.to_csv(os.path.join(pc_dir, "log_normalized_data.csv"))
#
#
#
#
#         # # Plot QQ plots and histograms for each feature in each normalized data frame
#         #     for df_name, df in [("Original Data", pc_group),
#         #                         ("Min-Max Normalized Data", min_max_normalized_data),
#         #                         ("Central Log Normalized Data", central_log_normalized_data),
#         #                         ("Z-score Normalized Data", z_score_normalized_data),
#         #                         ("Box_Cox Normalized Data", transformed_data),
#         #                         ("Log Normalized Data", log_data)]:
#         #
#         #
#         #
#         #
#         #
#         #         # Create a directory for the current data frame
#         #         df_dir = os.path.join(pc_dir, df_name)
#         #         os.makedirs(df_dir, exist_ok=True)
#         #
#         #         for feature in df.columns:
#         #
#         #                 # Convert Series to array for plotting QQ plot
#         #             data_array = (df[feature]).astype(int)
#         #                 #type(data_array)
#         #             print(data_array)
#         #             new_array = data_array.to_numpy
#         #                 #type(new_array)
#         #             print(new_array)
#         #
#         #                 # QQ plot
#         #                 #fig = plt.figure(figsize=(8, 6))
#         #
#         #             sm.qqplot(df[feature], line='45')
#         #                 #stats.probplot(new_array, dist="norm", plot=plt)
#         #             plt.title(f"QQ Plot for {feature} ({df_name})")
#         #             qq_plot_path = os.path.join(df_dir, f"qq_plot_{feature}.png")
#         #             plt.savefig(qq_plot_path)
#         #
#         #             plt.close()
#         #             plt.clf()
#         #
#         #                 # Histogram
#         #             fig = plt.figure(figsize=(8, 6))
#         #
#         #             plt.hist(df[feature], bins=20, color='skyblue', edgecolor='black')
#         #             plt.title(f"Histogram for {feature} ({df_name})")
#         #             plt.xlabel("Value")
#         #             plt.ylabel("Frequency")
#         #             plt.tight_layout()
#         #             histogram_path = os.path.join(df_dir, f"histogram_{feature}.png")
#         #             plt.savefig(histogram_path)
#         #
#         #             plt.close(fig)
#         #             plt.clf()
#         #
#         #
#         # print("Normalization and visualization completed. Results saved in 'normalization' directory.")
#         #
#         #
#         #
#         # return data
#
# #
# # # RUN
# csv_file_path = "IRF/results_irf_first/by_well_results/filtered_field_data.csv"
# data = process_data(csv_file_path)

### todo save meta + feat - following removal for batch verification

## imp v2
#### post batch well proc - good!
# def process_data(csv_file_path):
#     #  folder path and filename
#     folder_path, filename = os.path.split(csv_file_path)
#
#     #results folder path
#     results_folder_path = os.path.join(folder_path, "results_proc_true")
#
#     # Readfile
#     data = pd.read_csv(csv_file_path)
#
#     #make results folder
#     os.makedirs(results_folder_path, exist_ok=True)
#
#     #edit column names
#     data.columns = [col.split("[")[0].strip() if "[" in col else col for col in data.columns]
#     data.to_csv(os.path.join(results_folder_path, "tool_processed_data_1.csv"))
#
#     # User to remove rows based on previous batch analysis
#     while True:
#         print("Columns in the dataset:")
#         for idx, column in enumerate(data.columns):
#             print(f"{idx}: {column}")
#
#         column_to_filter = input("Enter the column name to filter by (or 'no' to skip): ").strip()
#         if column_to_filter.lower() == 'no':
#             break
#
#         if column_to_filter not in data.columns:
#             print("Invalid column name. Please try again.")
#             continue
#
#         print(f"Values in the selected column '{column_to_filter}':")
#         print(data[column_to_filter].unique())
#
#         values_to_remove = input("Enter the values to remove (comma-separated): ").split(',')
#         values_to_remove = [val.strip() for val in values_to_remove]
#
#         data = data[~data[column_to_filter].isin(values_to_remove)]
#         print(f"Rows with values {values_to_remove} in column '{column_to_filter}' have been removed.")
#
#     #columns
#     print("Columns in the dataset:")
#     for idx, column in enumerate(data.columns):
#         print(f"{idx}: {column}")
#
#     #user to specify columns to drop
#     columns_to_drop = input("Enter columns to drop (comma-separated): ").split(',')
#     # Remove any leading/trailing whitespace from column names
#     columns_to_drop = [col.strip() for col in columns_to_drop]
#
#     #user to specify columns to set as indexes
#     ## add PC as must.
#     columns_to_set_index = input("Enter columns to set as indexes (comma-separated): ").split(',')
#     # Remove any leading/trailing whitespace from column names
#     columns_to_set_index = [col.strip() for col in columns_to_set_index]
#
#     # Drop specified columns
#     if columns_to_drop:
#         data = data.drop(columns=columns_to_drop)
#
#     # Set specified columns as indexes
#     if columns_to_set_index:
#         data = data.set_index(columns_to_set_index, drop = False)
#
#     # Save current df
#     data.to_csv(os.path.join(results_folder_path, "tool_processed_data_2.csv"))
#
#     # Generate histogram for cell count feature
#     fig = px.histogram(data, x="all_cell_count", title="Histogram of Cell Count")
#     histogram_file_path = os.path.join(results_folder_path, "histogram.html")
#     fig.write_html(histogram_file_path)
#
#     # Create statistics for "cell count" feature
#     stats = data["all_cell_count"].describe()
#     stats_file_path = os.path.join(results_folder_path, "cell_count-statistics.txt")
#     with open(stats_file_path, "w") as f:
#         f.write(stats.to_string())
#     ## add summary saved bla nla
#     # Ask user for minimum and maximum values
#     min_value = float(input("Enter the minimum value for cell count: "))
#     max_value = float(input("Enter the maximum value for cell count: "))
#
#     # Remove rows with cell count outside the specified range
#     data = data[(data["all_cell_count"] >= min_value) & (data["all_cell_count"] <= max_value)]
#
#     # Save DataFrame in results folder
#     data.to_csv(os.path.join(results_folder_path, "tool_processed_data_3.csv"))
#
#     # Prompt user to choose index column for outlier detection
#     print("Columns available for outlier detection:")
#     print(data.columns)
#     index_column = input("Choose one column as index for outlier detection: ")
#
#     # Outlier detection using Local Outlier Factor
#     types = data[index_column].unique()
#     print(types)
#     mask = []
#     features = data.select_dtypes(include='number').columns
#
#     detector_list = [
#         ("Local Outlier Factor 30", LocalOutlierFactor(n_neighbors=30))
#     ]
#
#     for name, algorithm in detector_list:
#         errors = np.full(len(data),fill_value=np.nan)
#         outliers = np.full(len(data),fill_value=np.nan)
#
#         for type in types:
#             x = data.loc[:,features].values
#             F = x.sum(1)
#             mask = np.zeros(x.shape[0])
#             mask[np.isfinite(F)] = 1
#             mask_type = mask * np.array(data[index_column] == type)
#             Curr_df = data.loc[mask_type == 1, features]
#             x = Curr_df.values
#
#
#             if name == "Local Outlier Factor 30":
#                 algorithm.fit(x)
#                 errors[mask_type==1]  = algorithm.negative_outlier_factor_
#                 outliers[mask_type==1] = algorithm.fit_predict(x)
#
#         data[name] = errors
#         data[f'{name}_outliers'] = outliers
#         data.set_index(name, inplace=True,
#                                                append=True, drop=False)
#         data.to_csv(os.path.join(results_folder_path, "tool_processed_data_outliers_1.csv"))
#     # Exclude rows that were defined as outliers
#         for col in data.columns:
#             if col.endswith("_outliers"):
#                 data = data[data[col] != -1]
#
#
#         data.to_csv(os.path.join(results_folder_path, "tool_processed_data_outliers_2.csv"))
#
#         # Feature creation
#         print("Features are being computed.")
#         classes = ["network", "rod", "rounded"]
#         for class_type in classes:
#             class_total_area = f"{class_type}_total_area"
#             class_count = f"{class_type}_count"
#             all_cell_count = "all_cell_count"
#
#             data[f"{class_type}_average_size"] = data[class_total_area] / data[class_count]
#             data[f"{class_type}_average_size_per_cell"] = data[f"{class_type}_average_size"] / data[all_cell_count]
#             data[f"{class_type}_count_per_cell"] = data[class_count] / data[all_cell_count]
#             data[f"{class_type}_area_per_cell"] = data[class_total_area] / data[all_cell_count]
#
#         data["mitochondria_area_per_cell"] = data[[f"{class_type}_area_per_cell" for class_type in classes]].sum(axis=1)
#
#         for class_type in classes:
#             data[f"{class_type}_average_size_per_cell_norm"] = data[f"{class_type}_average_size_per_cell"] / data[
#                 "mitochondria_area_per_cell"]
#             data[f"{class_type}_count_per_cell_norm"] = data[f"{class_type}_count_per_cell"] / data[
#                 "mitochondria_area_per_cell"]
#             data[f"{class_type}_fraction"] = data[f"{class_type}_area_per_cell"] / data["mitochondria_area_per_cell"]
#
#         # Save data with new features
#         data.to_csv(os.path.join(results_folder_path, "tool_processed_data_features.csv"))
#
#         # Feature selection
#         print("Performing advanced method testing for variance contribution and ranking features.")
#         # Calculate variance of each feature across the entire dataset
#         features = data.select_dtypes(include='number').columns
#         existing_features = [feature for feature in features if feature in features]
#         subgroup_data = data[existing_features]
#         # Drop features with all missing values
#         subgroup_data = subgroup_data.dropna(axis=1, how='all')
#         # Fill remaining missing values with column means
#         subgroup_data = subgroup_data.fillna(subgroup_data.mean())
#         feature_variances = subgroup_data.var()
#         ranked_features = feature_variances.sort_values(ascending=False).index.tolist()
#         ranked_features_df = pd.DataFrame({"Feature": ranked_features, "Variance": feature_variances[ranked_features]})
#         ranked_features_df.to_csv(os.path.join(results_folder_path, "feature_ranking.csv"), index=False)
#
#         selected_data = data[existing_features]
#         selected_data.index = data.index
#         selected_data.to_csv(os.path.join(results_folder_path, "processed_data_4.csv"))
#         # Assess feature quality
#         # Prompt the user to select the correct index title as the target variable
#         print("Select the correct index title to be the target variable:")
#         for idx, index_title in enumerate(selected_data.index.names):
#             print(f"{idx + 1}: {index_title}")
#
#         selected_index = int(input("Enter the number corresponding to the correct index title: ")) - 1
#
#         # Set the selected index title as the target variable
#         target_variable = selected_data.index.names[selected_index]
#
#         print(f"Target variable selected: {target_variable}")
#         target_variable_labels, _ = pd.factorize(selected_data.index.get_level_values(target_variable))
#         print(target_variable_labels)
#         ##
#         # Create a directory for feature QC results
#         feature_qc_dir = os.path.join(results_folder_path, "feature_qc")
#         os.makedirs(feature_qc_dir, exist_ok=True)
#
#         # 1. Statistical Summary CSV
#         statistical_summary = selected_data.describe()
#         statistical_summary.to_csv(os.path.join(feature_qc_dir, "statistical_summary.csv"))
#
#         # 2. Feature Stability Scores (CV)
#         stability_scores = selected_data.var() / selected_data.mean()
#
#         # 3. Correlation Analysis
#         correlation_matrix = selected_data.corr(method='spearman')
#         correlation_matrix.to_csv(os.path.join(feature_qc_dir, "correlation_matrix.csv"))
#
#         # Plot Correlation Matrix
#         plt.figure(figsize=(15, 12))
#         sns.heatmap(correlation_matrix, annot=True, cmap='coolwarm', fmt=".2f")
#         plt.title("Correlation Matrix")
#         plt.tight_layout()
#         correlation_plot_path = os.path.join(feature_qc_dir, "correlation_matrix_plot.png")
#         plt.savefig(correlation_plot_path)
#         plt.close()
#         print(f"Correlation matrix plot saved to: {correlation_plot_path}")
#
#         # 4. CV Plot (if stability scores are available)
#         if stability_scores is not None:
#             plt.figure(figsize=(12, 6))
#             plt.bar(stability_scores.index, stability_scores.values, color='skyblue')
#             plt.title("Feature Stability Scores (Coefficient of Variation)")
#             plt.xlabel("Feature")
#             plt.ylabel("Stability Score")
#             plt.xticks(rotation=45, ha='right')
#             plt.tight_layout()
#             cv_plot_path = os.path.join(feature_qc_dir, "stability_scores_plot.png")
#             plt.savefig(cv_plot_path)
#             plt.close()
#             print(f"CV plot saved to: {cv_plot_path}")
#         else:
#             print("Stability scores not available. Skipping CV plot.")
#
#         # 5. Random Forest Score computation
#         target_variable_labels, _ = pd.factorize(selected_data.index.get_level_values(target_variable))
#         rf_scores = {}
#         for feature in selected_data.columns:
#             model = RandomForestRegressor()
#             feature_values = selected_data[feature].values.reshape(-1, 1)
#             model.fit(feature_values, target_variable_labels)
#             rf_scores[feature] = model.score(feature_values, target_variable_labels)
#
#         # Save RF scores to CSV
#         rf_scores_df = pd.DataFrame(rf_scores, index=['RF_Score'])
#         rf_scores_df.to_csv(os.path.join(feature_qc_dir, "random_forest_scores.csv"))
#
#         # Plot RF Scores
#         plt.figure(figsize=(10, 6))
#         plt.bar(rf_scores.keys(), rf_scores.values(), color='lightgreen')
#         plt.title("Random Forest Scores for each Feature")
#         plt.xlabel("Feature")
#         plt.ylabel("Random Forest Score")
#         plt.xticks(rotation=45, ha='right')
#         plt.tight_layout()
#         rf_plot_path = os.path.join(feature_qc_dir, "random_forest_scores_plot.png")
#         plt.savefig(rf_plot_path)
#         plt.close()
#         print(f"Random Forest scores plot saved to: {rf_plot_path}")
#         ##
#         # Ask the user if they want to change feature names
#         print("final_proc_pre_norm")
#         rename_columns = input("Do you want to rename any columns? (yes/no): ")
#         if rename_columns.lower() == 'yes':
#             print("Current Columns:")
#             for idx, column in enumerate(selected_data.columns):
#                 print(f"{idx}: {column}")
#
#             renaming = input(
#                 "Enter the column name and new name separated by ':', each pair separated by commas (e.g., col1:new_col1,col2:new_col2): ")
#             rename_dict = dict(pair.split(':') for pair in renaming.split(','))
#             selected_data.rename(columns=rename_dict, inplace=True)
#             print("Names changed successfully:)")
#         ##
# #imputation segment
#         # Initialize the SimpleImputer for missing values (NaN)
#         imputer = SimpleImputer(strategy='median')
#         # Impute missing values (NaN)
#         df_imputed = pd.DataFrame(imputer.fit_transform(selected_data), columns=selected_data.columns)
#         # Replace zero values with NaN to use the same imputer
#         df_imputed.replace(0, np.nan, inplace=True)
#         # Impute zero values (now NaNs) using the same strategy
#         df_final_imputed = pd.DataFrame(imputer.fit_transform(df_imputed), columns=selected_data.columns)
#         df_final_imputed.index = selected_data.index
#         df_final_imputed.to_csv(os.path.join(results_folder_path, "df_final_imputed.csv"))
#
#
#         print("ready_for_norm")
#         selected_data = df_final_imputed.copy()
#         print("ready_for_norm")
#         #Normalize features using different techniques
#         normalization_dir = os.path.join(results_folder_path, "normalization")
#         os.makedirs(normalization_dir, exist_ok=True)
#         pc_groups = selected_data.groupby(level='incubation_condition', group_keys=False)
#         df_reverted = pc_groups.apply(lambda x: x)
#         df_reverted.to_csv(os.path.join(results_folder_path, "processed_data_5.csv"))
#         for pc_value, pc_group in pc_groups:
#             # Create a directory for the current PC value
#             pc_dir = os.path.join(normalization_dir, f"df_{pc_value}")
#             os.makedirs(pc_dir, exist_ok=True)
#             # Perform normalization for each group
#             pc_group.to_csv(os.path.join(results_folder_path, "original_data.csv"))
#             # For Min-Max normalization
#             min_max_normalized_data = (pc_group - pc_group.min()) / (pc_group.max() - pc_group.min())
#             min_max_normalized_data.to_csv(os.path.join(pc_dir, "min_max_normalized_data.csv"))
#
#             # For Central Logarithmic normalization
#             central_log_normalized_data = np.log(pc_group + np.sqrt(pc_group ** 2 + 1))
#             central_log_normalized_data.to_csv(os.path.join(pc_dir, "central_log_normalized_data.csv"))
#
#             # For Z-score normalization
#             z_score_normalized_data = (pc_group - pc_group.mean()) / pc_group.std()
#             z_score_normalized_data.to_csv(os.path.join(pc_dir, "z_score_normalized_data.csv"))
#
#             # For Box-Cox transformation
#             transformed_data = pd.DataFrame()
#             m = pc_group.lt(0).any()
#             pc_group = pc_group.loc[:, ~m]
#             pc_group = pc_group.drop(columns="Local Outlier Factor 30_outliers")
#             print(pc_group)
#             for column in pc_group.columns:
#                 #pc_group = pc_group.loc[:, pc_group.ge(0).all()]
#                 if column not in columns_to_set_index:
#                     print(column)
#                     transformed_column, _ = boxcox(pc_group[column])
#                     transformed_data[column] = transformed_column
#                 transformed_data.index = pc_group.index
#                 transformed_data.to_csv(os.path.join(pc_dir, "box_cox_normalized_data.csv"))
#
#             # For log normalization
#             log_data = pd.DataFrame()
#             for column in pc_group.columns:
#                 log_column = np.log(pc_group[column])
#                 log_data[column] = log_column
#             log_data.index = pc_group.index
#             log_data.to_csv(os.path.join(pc_dir, "log_normalized_data.csv"))
#
#
#
#
#         # # Plot QQ plots and histograms for each feature in each normalized data frame
#             for df_name, df in [("Original Data", pc_group),
#                                 ("Min-Max Normalized Data", min_max_normalized_data),
#                                 ("Central Log Normalized Data", central_log_normalized_data),
#                                 ("Z-score Normalized Data", z_score_normalized_data),
#                                 ("Box_Cox Normalized Data", transformed_data),
#                                 ("Log Normalized Data", log_data)]:
#
#
#
#
#
#                 # Create a directory for the current data frame
#                 df_dir = os.path.join(pc_dir, df_name)
#                 os.makedirs(df_dir, exist_ok=True)
#                 # Drop features with all missing values
#                 df = df.dropna(axis=1, how='all')
#                 # Fill remaining missing values with column means
#                 df = df.fillna(df.mean())
#                 for feature in df:
#                     ## qq imp
#                     plt.figure()
#                     scipy.stats.probplot(df[feature], dist="norm", plot=plt)
#                     plt.title(f"Q-Q-{feature}")
#                     qq_plot_path_2 = os.path.join(df_dir, f"qq_2_plot_{feature}.png")
#                     plt.savefig(qq_plot_path_2)
#                     plt.close()
#                     plt.clf()
#
#                 for feature in df.columns:
#                     # Convert Series to array for plotting QQ plot
#                     data_array = df[feature].to_numpy()
#                     # QQ plot
#                     sm.qqplot(data_array, line='45')
#                     plt.title(f"QQ Plot for {feature} ({df_name})")
#                     qq_plot_path = os.path.join(df_dir, f"qq_plot_{feature}.png")
#                     plt.savefig(qq_plot_path)
#                     plt.close()
#                     plt.clf()
#
#                     # Histogram
#                     fig = plt.figure(figsize=(8, 6))
#                     plt.hist(df[feature], bins=20, color='skyblue', edgecolor='black')
#                     plt.title(f"Histogram for {feature} ({df_name})")
#                     plt.xlabel("Value")
#                     plt.ylabel("Frequency")
#                     plt.tight_layout()
#                     histogram_path = os.path.join(df_dir, f"histogram_{feature}.png")
#                     plt.savefig(histogram_path)
#                     plt.close(fig)
#                     plt.clf()
#
#         print("Normalization and visualization completed. Results saved in 'normalization' directory.")
#
#         # Prompt user to select the normalized data frame to continue with
#         print("Select the normalized data frame you want to continue with:")
#         normalized_options = ["Min-Max Normalized Data", "Central Log Normalized Data", "Z-score Normalized Data",
#                               "Box_Cox Normalized Data", "Log Normalized Data"]
#         for idx, option in enumerate(normalized_options):
#             print(f"{idx + 1}: {option}")
#
#         selected_normalization = int(input("Enter the number corresponding to the desired normalized data frame: ")) - 1
#         selected_normalized_data = normalized_options[selected_normalization]
#
#         # Combine all PC values for the selected normalization into one final normalized data frame
#         final_normalized_data = pd.DataFrame()
#         for pc_value, pc_group in pc_groups:
#             pc_dir = os.path.join(normalization_dir, f"df_{pc_value}")
#             file_path = os.path.join(pc_dir, f"{selected_normalized_data.lower().replace(' ', '_')}.csv")
#             normalized_df = pd.read_csv(file_path, index_col=0)
#             final_normalized_data = pd.concat([final_normalized_data, normalized_df])
#
#         final_normalized_data.to_csv(os.path.join(results_folder_path, "final_normalized_data.csv"))
#         print(f"Final normalized data saved to: {os.path.join(results_folder_path, 'final_normalized_data.csv')}")
#
#         return final_normalized_data
#         # normalization_dir = os.path.join(results_folder_path, "normalization")
#         # os.makedirs(normalization_dir, exist_ok=True)
#         #
#         #
#         #
#         #
#         # # Split selected_data based on the 'PC' index
#         # pc_groups = selected_data.groupby(level='incubation_time', group_keys=False)
#         # df_reverted = pc_groups.apply(lambda x: x)
#         # df_reverted.to_csv(os.path.join(results_folder_path, "processed_data_5.csv"))
#         # # Iterate over each group and perform normalization
#         # # Perform normalization for each group
#         # selected_data.to_csv(os.path.join(normalization_dir, "original_data.csv"))
#         # # For Min-Max normalization
#         # min_max_normalized_data = (selected_data - selected_data.min()) / (selected_data.max() - selected_data.min())
#         # min_max_normalized_data.to_csv(os.path.join(normalization_dir, "min_max_normalized_data.csv"))
#         #
#         # # For Central Logarithmic normalization
#         # central_log_normalized_data = np.log(selected_data + np.sqrt(selected_data ** 2 + 1))
#         # central_log_normalized_data.to_csv(os.path.join(normalization_dir, "central_log_normalized_data.csv"))
#         #
#         # # For Z-score normalization
#         # z_score_normalized_data = (selected_data - selected_data.mean()) / selected_data.std()
#         # z_score_normalized_data.to_csv(os.path.join(normalization_dir, "z_score_normalized_data.csv"))
#         #
#         # # For Box-Cox transformation
#         # transformed_data = pd.DataFrame()
#         # for column in selected_data.columns:
#         #     transformed_column, _ = boxcox(selected_data[column])
#         #     transformed_data[column] = transformed_column
#         # transformed_data.index = selected_data.index
#         # transformed_data.to_csv(os.path.join(normalization_dir, "box_cox_normalized_data.csv"))
#         #
#         # # For log normalization
#         # log_data = pd.DataFrame()
#         # for column in selected_data.columns:
#         #     log_column = np.log(selected_data[column])
#         #     log_data[column] = log_column
#         # log_data.index = selected_data.index
#         # log_data.to_csv(os.path.join(normalization_dir, "log_normalized_data.csv"))
#         #
#         # for df_name, df in [("Original Data", selected_data),
#         #                         ("Min-Max Normalized Data", min_max_normalized_data),
#         #                         ("Central Log Normalized Data", central_log_normalized_data),
#         #                         ("Z-score Normalized Data", z_score_normalized_data),
#         #                         ("Box_Cox Normalized Data", transformed_data),
#         #                         ("Log Normalized Data", log_data)]:
#         #
#         #
#         #
#         #
#         #
#         #         # Create a directory for the current data frame
#         #         df_dir = os.path.join(normalization_dir, df_name)
#         #         os.makedirs(df_dir, exist_ok=True)
#         #
#         #         for feature in df.columns:
#         #
#         #                 # Convert Series to array for plotting QQ plot
#         #             data_array = (df[feature]).astype(int)
#         #                 #type(data_array)
#         #             print(data_array)
#         #             new_array = data_array.to_numpy
#         #                 #type(new_array)
#         #             print(new_array)
#         #
#         #                 # QQ plot
#         #                 #fig = plt.figure(figsize=(8, 6))
#         #
#         #             sm.qqplot(df[feature], line='45')
#         #                 #stats.probplot(new_array, dist="norm", plot=plt)
#         #             plt.title(f"QQ Plot for {feature} ({df_name})")
#         #             qq_plot_path = os.path.join(df_dir, f"qq_plot_{feature}.png")
#         #             plt.savefig(qq_plot_path)
#         #
#         #             plt.close()
#         #             plt.clf()
#         #
#         #                 # Histogram
#         #             fig = plt.figure(figsize=(8, 6))
#         #
#         #             plt.hist(df[feature], bins=20, color='skyblue', edgecolor='black')
#         #             plt.title(f"Histogram for {feature} ({df_name})")
#         #             plt.xlabel("Value")
#         #             plt.ylabel("Frequency")
#         #             plt.tight_layout()
#         #             histogram_path = os.path.join(df_dir, f"histogram_{feature}.png")
#         #             plt.savefig(histogram_path)
#         #
#         #             plt.close(fig)
#         #             plt.clf()
#         #
#         #
#         # print("Normalization and visualization completed. Results saved in 'normalization' directory.")
#         #
#         #
#         #
#         # return data
#         # for pc_value, pc_group in pc_groups:
#         #     # Create a directory for the current PC value
#         #     pc_dir = os.path.join(normalization_dir, f"df_{pc_value}")
#         #     os.makedirs(pc_dir, exist_ok=True)
#         #     # Perform normalization for each group
#         #     pc_group.to_csv(os.path.join(results_folder_path, "original_data.csv"))
#         #     # For Min-Max normalization
#         #     min_max_normalized_data = (pc_group - pc_group.min()) / (pc_group.max() - pc_group.min())
#         #     min_max_normalized_data.to_csv(os.path.join(pc_dir, "min_max_normalized_data.csv"))
#         #
#         #     # For Central Logarithmic normalization
#         #     central_log_normalized_data = np.log(pc_group + np.sqrt(pc_group ** 2 + 1))
#         #     central_log_normalized_data.to_csv(os.path.join(pc_dir, "central_log_normalized_data.csv"))
#         #
#         #     # For Z-score normalization
#         #     z_score_normalized_data = (pc_group - pc_group.mean()) / pc_group.std()
#         #     z_score_normalized_data.to_csv(os.path.join(pc_dir, "z_score_normalized_data.csv"))
#         #
#         #     # For Box-Cox transformation
#         #     transformed_data = pd.DataFrame()
#         #     for column in pc_group.columns:
#         #         transformed_column, _ = boxcox(pc_group[column])
#         #         transformed_data[column] = transformed_column
#         #     transformed_data.index = pc_group.index
#         #     transformed_data.to_csv(os.path.join(pc_dir, "box_cox_normalized_data.csv"))
#         #
#         #     # For log normalization
#         #     log_data = pd.DataFrame()
#         #     for column in pc_group.columns:
#         #         log_column = np.log(pc_group[column])
#         #         log_data[column] = log_column
#         #     log_data.index = pc_group.index
#         #     log_data.to_csv(os.path.join(pc_dir, "log_normalized_data.csv"))
#         #
#         #
#         #
#         #
#         # # # Plot QQ plots and histograms for each feature in each normalized data frame
#         #     for df_name, df in [("Original Data", pc_group),
#         #                         ("Min-Max Normalized Data", min_max_normalized_data),
#         #                         ("Central Log Normalized Data", central_log_normalized_data),
#         #                         ("Z-score Normalized Data", z_score_normalized_data),
#         #                         ("Box_Cox Normalized Data", transformed_data),
#         #                         ("Log Normalized Data", log_data)]:
#         #
#         #
#         #
#         #
#         #
#         #         # Create a directory for the current data frame
#         #         df_dir = os.path.join(pc_dir, df_name)
#         #         os.makedirs(df_dir, exist_ok=True)
#         #
#         #         for feature in df.columns:
#         #
#         #                 # Convert Series to array for plotting QQ plot
#         #             data_array = (df[feature]).astype(int)
#         #                 #type(data_array)
#         #             print(data_array)
#         #             new_array = data_array.to_numpy
#         #                 #type(new_array)
#         #             print(new_array)
#         #
#         #                 # QQ plot
#         #                 #fig = plt.figure(figsize=(8, 6))
#         #
#         #             sm.qqplot(df[feature], line='45')
#         #                 #stats.probplot(new_array, dist="norm", plot=plt)
#         #             plt.title(f"QQ Plot for {feature} ({df_name})")
#         #             qq_plot_path = os.path.join(df_dir, f"qq_plot_{feature}.png")
#         #             plt.savefig(qq_plot_path)
#         #
#         #             plt.close()
#         #             plt.clf()
#         #
#         #                 # Histogram
#         #             fig = plt.figure(figsize=(8, 6))
#         #
#         #             plt.hist(df[feature], bins=20, color='skyblue', edgecolor='black')
#         #             plt.title(f"Histogram for {feature} ({df_name})")
#         #             plt.xlabel("Value")
#         #             plt.ylabel("Frequency")
#         #             plt.tight_layout()
#         #             histogram_path = os.path.join(df_dir, f"histogram_{feature}.png")
#         #             plt.savefig(histogram_path)
#         #
#         #             plt.close(fig)
#         #             plt.clf()
#         #
#         #
#         # print("Normalization and visualization completed. Results saved in 'normalization' directory.")
#         #
#         #
#         #
#         # return data
#
#
# # # RUN
# # csv_file_path = "IRF/results_irf_first/by_well_results/filtered_field_data.csv"
# # data = process_data(csv_file_path)
# csv_file_path = "GSD3/results_field_GSD3/by_well_results/filtered_field_data_gsd3.csv"
# data = process_data(csv_file_path)


### next step
## statistics in R - improve

## todo miltivar
## todo LDA, PCA, fraction - optimize. add corleation function,



###multi var
###### imp try - with groupby - work on it.
# def load_data(file_path):
#     return pd.read_csv(file_path)
#
# def prompt_columns_removal(df):
#     print("Current columns:", df.columns.tolist())
#     cols_to_remove = input("Enter columns to remove (comma-separated) or press Enter to skip: ").split(',')
#     cols_to_remove = [col.strip() for col in cols_to_remove if col.strip() in df.columns]
#     return df.drop(columns=cols_to_remove) if cols_to_remove else df
#
# def prompt_index_column(df):
#     print("Current columns:", df.columns.tolist())
#     index_col = input("Enter the column to use as index or press Enter to skip: ")
#     return index_col if index_col in df.columns else None
#
# def prompt_secondary_index_columns(df):
#     print("Current columns:", df.columns.tolist())
#     secondary_index_cols = input("Enter any secondary index columns (comma-separated) or press Enter to skip: ").split(',')
#     secondary_index_cols = [col.strip() for col in secondary_index_cols if col.strip() in df.columns]
#     return secondary_index_cols
#
# def prompt_groupby_index(segment_name):
#     groupby = input(f"Do you want to group by the index column for {segment_name}? (yes/no): ").strip().lower()
#     return groupby == 'yes'
#
# def zscore_normalize(df, index_col, secondary_index_cols):
#     exclude_cols = [index_col] + secondary_index_cols
#     numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
#     numeric_cols = [col for col in numeric_cols if col not in exclude_cols]
#     df[numeric_cols] = df[numeric_cols].apply(zscore)
#     return df
#
# def create_line_plots(df, index_col, secondary_index_cols, results_dir):
#     os.makedirs(results_dir, exist_ok=True)
#
#     pc_col = 'PC'
#     if pc_col not in df.columns:
#         print(f"Column '{pc_col}' not found in the dataframe.")
#         return
#
#     for feature in df.columns:
#         if feature not in [index_col, pc_col] + secondary_index_cols:
#             plt.figure()
#             sns.lineplot(data=df, x=index_col, y=feature, hue=pc_col, n_boot=1000, errorbar=('ci', 90), dashes=True)
#             plt.title(f"Line Plot for Feature: {feature}")
#             plt.xticks(rotation=45, fontsize=8)
#             plt.savefig(os.path.join(results_dir, f"line_plot_{feature}.png"))
#             plt.close()
#     print("Line plots done!")
#
# def bootstrap_ci(data, num_bootstrap=1000, ci=90):
#     boot_means = np.random.choice(data, (num_bootstrap, len(data))).mean(axis=1)
#     lower_bound = np.percentile(boot_means, (100 - ci) / 2)
#     upper_bound = np.percentile(boot_means, 100 - (100 - ci) / 2)
#     return lower_bound, upper_bound
#
# def create_interactive_line_plots(df, index_col, secondary_index_cols, results_dir):
#     os.makedirs(results_dir, exist_ok=True)
#     pc_col = 'PC'
#     if (pc_col not in df.columns) or (index_col not in df.columns):
#         print(f"Column '{pc_col}' or '{index_col}' not found in the dataframe.")
#         return
#
#     for feature in df.columns:
#         if feature not in [index_col, pc_col] + secondary_index_cols:
#             fig = make_subplots()
#
#             for pc in df[pc_col].unique():
#                 df_pc = df[df[pc_col] == pc]
#                 means = df_pc.groupby(index_col)[feature].mean()
#                 cis = df_pc.groupby(index_col)[feature].apply(bootstrap_ci)
#
#                 lower_bounds = [ci[0] for ci in cis]
#                 upper_bounds = [ci[1] for ci in cis]
#
#                 fig.add_trace(go.Scatter(
#                     x=means.index,
#                     y=means,
#                     mode='lines+markers',
#                     name=f'{pc_col} {pc}',
#                     line=dict(shape='linear')
#                 ))
#
#                 fig.add_trace(go.Scatter(
#                     x=means.index,
#                     y=lower_bounds,
#                     mode='lines',
#                     line=dict(width=0),
#                     showlegend=False
#                 ))
#
#                 fig.add_trace(go.Scatter(
#                     x=means.index,
#                     y=upper_bounds,
#                     mode='lines',
#                     fill='tonexty',
#                     fillcolor='rgba(0,100,80,0.2)',
#                     line=dict(width=0),
#                     showlegend=False
#                 ))
#
#             fig.update_layout(
#                 title=f"Line Plot for Feature: {feature}",
#                 xaxis_title=index_col,
#                 yaxis_title=feature,
#                 legend_title=pc_col
#             )
#             fig.write_html(os.path.join(results_dir, f"interactive_line_plot_{feature}.html"))
#     print("Interactive line plots done!")
#
# def create_heatmap(df, index_col, secondary_index_cols, results_dir):
#     print("Heatmap creation has begun!")
#     df_heat = df.copy()
#     df_heat = prompt_columns_removal(df_heat)
#
#     print("Index column values:", df_heat[index_col].unique())
#     order_x = input(f"Enter the order for {index_col} values (comma-separated) or press Enter to skip: ").split(',')
#     order_y = input("Enter the order for features (comma-separated) or press Enter to skip: ").split(',')
#     order_y = [col.strip() for col in order_y if col.strip() in df_heat.columns]
#
#     if prompt_groupby_index("heatmap"):
#         df_grouped = df_heat.groupby(index_col).mean().reset_index()
#     else:
#         df_grouped = df_heat.copy()
#
#     if order_x:
#         df_grouped = df_grouped.set_index(index_col).loc[order_x].reset_index()
#     if order_y:
#         df_grouped = df_grouped[[index_col] + order_y]
#
#     lan = plt.rcParams['font.family'] = ['Arial']
#     plt.figure(figsize=(20, 12))
#     font = {'family': 'Arial', 'size': 22}
#     plt.rc('font', **font)
#     plt.yticks(family='Arial', rotation="horizontal", size=10)
#     plt.xticks(family='Arial', rotation="vertical", size=8)
#     plt.margins(0.1)
#     plt.subplots_adjust(bottom=0.3)
#     sns.heatmap(df_grouped.set_index(index_col), cmap='viridis')
#     plt.title("Heatmap")
#     plt.savefig(os.path.join(results_dir, "heatmap.png"))
#     plt.close()
#     print("Heatmap complete!")
#
#     fig = px.imshow(df_grouped.set_index(index_col), title="Interactive Heatmap")
#     fig.write_html(os.path.join(results_dir, "interactive_heatmap.html"))
#     print("Interactive Heatmap complete!")
#
# def perform_pca(df, index_col, secondary_index_cols, results_dir):
#     print("PCA has begun!")
#     df_pca = df.copy()
#     df_pca = prompt_columns_removal(df_pca)
#     os.makedirs(results_dir, exist_ok=True)
#
#     pca = PCA(n_components=2)
#     exclude_cols = [index_col] + secondary_index_cols
#     numeric_cols = df_pca.select_dtypes(include=[np.number]).columns.tolist()
#     numeric_cols = [col for col in numeric_cols if col not in exclude_cols]
#     pca_result = pca.fit_transform(df_pca[numeric_cols])
#     explained_variance = pca.explained_variance_ratio_
#
#     unique_labels = df_pca[index_col].unique()
#     colors = pd.factorize(df_pca[index_col])[0]
#
#     plt.figure(figsize=(20, 12))
#     scatter = plt.scatter(pca_result[:, 0], pca_result[:, 1], c=colors, cmap='viridis')
#     plt.xlabel(f"PC1 ({explained_variance[0] * 100:.2f}%)")
#     plt.ylabel(f"PC2 ({explained_variance[1] * 100:.2f}%)")
#
#     colors_legend = plt.cm.viridis(np.linspace(0, 1, len(unique_labels)))
#     legend = plt.legend(handles=[mpatches.Circle((0, 0), color=colors_legend[i], label=unique_labels[i])
#                                  for i in range(len(unique_labels))], loc='upper right')
#     legend.set_title('index')
#     plt.title("PCA - All Data")
#     plt.savefig(os.path.join(results_dir, "pca_all.png"))
#     plt.close()
#
#     df_mean = df_pca.groupby(index_col).mean()
#     pca_result_mean = pca.fit_transform(df_mean)
#     explained_variance_mean = pca.explained_variance_ratio_
#
#     plt.figure(figsize=(20, 12))
#     scatter_mean = plt.scatter(pca_result_mean[:, 0], pca_result_mean[:, 1], c=np.arange(len(df_mean)), cmap='viridis')
#     plt.xlabel(f"PC1 ({explained_variance_mean[0] * 100:.2f}%)")
#     plt.ylabel(f"PC2 ({explained_variance_mean[1] * 100:.2f}%)")
#
#     colors_legend_mean = plt.cm.viridis(np.linspace(0, 1, len(unique_labels)))
#     legend_mean = plt.legend(handles=[mpatches.Circle((0, 0), color=colors_legend_mean[i], label=unique_labels[i])
#                                       for i in range(len(unique_labels))], loc='upper right')
#     legend_mean.set_title('index')
#     plt.title("PCA - Mean Data")
#     plt.savefig(os.path.join(results_dir, "pca_mean.png"))
#     plt.close()
#     print("PCA complete!")
#
# def perform_lda(df, index_col, secondary_index_cols, results_dir):
#     print("LDA has begun!")
#     df_lda = df.copy()
#     df_lda = prompt_columns_removal(df_lda)
#     os.makedirs(results_dir, exist_ok=True)
#
#     lda = LDA(n_components=2)
#     exclude_cols = [index_col] + secondary_index_cols
#     numeric_cols = df_lda.select_dtypes(include=[np.number]).columns.tolist()
#     numeric_cols = [col for col in numeric_cols if col not in exclude_cols]
#     lda_result = lda.fit_transform(df_lda[numeric_cols], df_lda[index_col])
#     explained_variance = lda.explained_variance_ratio_
#
#     unique_labels = df_lda[index_col].unique()
#     colors = pd.factorize(df_lda[index_col])[0]
#
#     plt.figure(figsize=(20, 12))
#     scatter = plt.scatter(lda_result[:, 0], lda_result[:, 1], c=colors, cmap='viridis')
#     plt.xlabel(f"LD1 ({explained_variance[0] * 100:.2f}%)")
#     plt.ylabel(f"LD2 ({explained_variance[1] * 100:.2f}%)")
#
#     colors_legend = plt.cm.viridis(np.linspace(0, 1, len(unique_labels)))
#     legend = plt.legend(handles=[mpatches.Circle((0, 0), color=colors_legend[i], label=unique_labels[i])
#                                  for i in range(len(unique_labels))], loc='upper right')
#     legend.set_title('index')
#     plt.title("LDA")
#     plt.savefig(os.path.join(results_dir, "lda.png"))
#     plt.close()
#     print("LDA complete!")
#
# def perform_kmeans(df, index_col, secondary_index_cols, results_dir):
#     print("k-means clustering has begun!")
#     df_kmeans = df.copy()
#     df_kmeans = prompt_columns_removal(df_kmeans)
#     os.makedirs(results_dir, exist_ok=True)
#
#     kmeans = KMeans(n_clusters=2)
#     exclude_cols = [index_col] + secondary_index_cols
#     numeric_cols = df_kmeans.select_dtypes(include=[np.number]).columns.tolist()
#     numeric_cols = [col for col in numeric_cols if col not in exclude_cols]
#     kmeans_result = kmeans.fit_transform(df_kmeans[numeric_cols])
#     clusters = kmeans.labels_
#
#     plt.figure(figsize=(20, 12))
#     scatter = plt.scatter(kmeans_result[:, 0], kmeans_result[:, 1], c=clusters, cmap='viridis')
#     plt.xlabel("K-means Component 1")
#     plt.ylabel("K-means Component 2")
#     plt.title("K-means Clustering")
#     plt.savefig(os.path.join(results_dir, "kmeans.png"))
#     plt.close()
#     print("K-means clustering complete!")
#
# def perform_hierarchical_clustering(df, index_col, secondary_index_cols, results_dir):
#     print("Hierarchical clustering has begun!")
#     df_hc = df.copy()
#     df_hc = prompt_columns_removal(df_hc)
#     os.makedirs(results_dir, exist_ok=True)
#
#     exclude_cols = [index_col] + secondary_index_cols
#     numeric_cols = df_hc.select_dtypes(include=[np.number]).columns.tolist()
#     numeric_cols = [col for col in numeric_cols if col not in exclude_cols]
#
#     df_mean = df_hc.groupby(index_col).mean()
#     z = linkage(df_mean[numeric_cols], method='ward')
#
#     plt.figure(figsize=(20, 12))
#     dendrogram(z, labels=df_mean.index.tolist())
#     plt.title("Hierarchical Clustering Dendrogram")
#     plt.savefig(os.path.join(results_dir, "hierarchical_clustering.png"))
#     plt.close()
#     print("Hierarchical clustering complete!")
#
# def bootstrap_ci_2(data, n_iterations=1000, ci=95):
#     means = []
#     for _ in range(n_iterations):
#         sample = np.random.choice(data, size=len(data), replace=True)
#         means.append(np.mean(sample))
#     lower_bound = np.percentile(means, (100 - ci) / 2)
#     upper_bound = np.percentile(means, 100 - (100 - ci) / 2)
#     return lower_bound, upper_bound
#
#
# def compute_and_plot_ci(df, index_col, result_dir, feature_list):
#     ci_dict = {}
#     for feature in feature_list:
#         ci_dict[f"{feature}_lower"] = []
#         ci_dict[f"{feature}_upper"] = []
#
#     for idx in df[index_col].unique():
#         subset = df[df[index_col] == idx]
#         for feature in feature_list:
#             lower_bound, upper_bound = bootstrap_ci_2(subset[feature])
#             ci_dict[f"{feature}_lower"].append(lower_bound)
#             ci_dict[f"{feature}_upper"].append(upper_bound)
#
#     for feature in feature_list:
#         df[f"{feature}_lower"] = df[index_col].map(
#             dict(zip(df[index_col].unique(), ci_dict[f"{feature}_lower"]))
#         )
#         df[f"{feature}_upper"] = df[index_col].map(
#             dict(zip(df[index_col].unique(), ci_dict[f"{feature}_upper"]))
#         )
#
#     df.to_csv(f"{result_dir}/ci_dataframe.csv", index=False)
#     df_ci = df.groupby(index_col).mean(numeric_only=True)
#     # df_ci.set_index([index_col], inplace=True,
#     #                       append=True, drop=True)
#     df_ci = df_ci * 100
#     df_ci = df_ci.reset_index()
#     df_ci.to_csv(f"{result_dir}/ci_dataframe-2.csv", index=False)
#     fig = go.Figure()
#     for feature in feature_list:
#         feature_lower_col = f"{feature}_lower"
#         feature_upper_col = f"{feature}_upper"
#
#         fig.add_trace(go.Bar(
#             x=df_ci[index_col],
#             y=df_ci[feature],
#             name=feature,
#             error_y=dict(
#                 type='data',
#                 symmetric=True,
#                 array=df_ci[feature_upper_col] * 2 - df_ci[feature] * 2,
#                 arrayminus=df_ci[feature] * 2 - df_ci[feature_lower_col] * 2,
#                 visible=True,
#                 thickness=3,
#                 color="purple",
#             ),
#         ))
#
#     fig.update_layout(
#         title="Total Mitochondrias by Group",
#         xaxis_title=index_col,
#         yaxis_title="Total Mitochondrias",
#         barmode='stack'
#     )
#
#     fig.write_image(f"{result_dir}/bar_plot.pdf", engine="kaleido")
#     fig.show()
#
#
#
#
# #
# def perform_multivariate_analysis(file_path, results_dir):
#     df = load_data(file_path)
#     df = prompt_columns_removal(df)
#     index_col = prompt_index_column(df)
#     secondary_index_cols = prompt_secondary_index_columns(df)
#     compute_and_plot_ci(df, index_col, results_dir,
#                         feature_list=["network_fraction", "rod_fraction", "rounded_fraction"])
#     df = zscore_normalize(df, index_col, secondary_index_cols)
#
#     #create_line_plots(df, index_col, secondary_index_cols, results_dir)
#     #create_interactive_line_plots(df, index_col, secondary_index_cols, results_dir)
#     create_heatmap(df, index_col, secondary_index_cols, results_dir)
#     perform_pca(df, index_col, secondary_index_cols, results_dir)
#     #perform_lda(df, index_col, secondary_index_cols, results_dir)
#     #perform_kmeans(df, index_col, secondary_index_cols, results_dir)
#     #perform_hierarchical_clustering(df, index_col, secondary_index_cols, results_dir)
#
#
# #
# #
# perform_multivariate_analysis('processed_data_5.csv', 'try_1/')

#incubation_time, Local Outlier Factor 30, row, column, field, id, passage, age, all_cell_count, field_count, incubation_time.1, Local Outlier Factor 30.1, Local Outlier Factor 30_outliers


# def bootstrap_ci(data, n_iterations=1000, ci=95):
#     means = []
#     for _ in range(n_iterations):
#         sample = np.random.choice(data, size=len(data), replace=True)
#         means.append(np.mean(sample))
#     lower_bound = np.percentile(means, (100 - ci) / 2)
#     upper_bound = np.percentile(means, 100 - (100 - ci) / 2)
#     return lower_bound, upper_bound
#
#
# def compute_and_plot_ci(df, index_col, feature_list, result_dir):
#     ci_dict = {}
#     for feature in feature_list:
#         ci_dict[f"{feature}_lower"] = []
#         ci_dict[f"{feature}_upper"] = []
#
#     for idx in df[index_col].unique():
#         subset = df[df[index_col] == idx]
#         for feature in feature_list:
#             lower_bound, upper_bound = bootstrap_ci(subset[feature])
#             ci_dict[f"{feature}_lower"].append(lower_bound)
#             ci_dict[f"{feature}_upper"].append(upper_bound)
#
#     for feature in feature_list:
#         df[f"{feature}_lower"] = df[index_col].map(
#             dict(zip(df[index_col].unique(), ci_dict[f"{feature}_lower"]))
#         )
#         df[f"{feature}_upper"] = df[index_col].map(
#             dict(zip(df[index_col].unique(), ci_dict[f"{feature}_upper"]))
#         )
#
#     df.to_csv(f"{result_dir}/ci_dataframe.csv", index=False)
#     df_ci = df.groupby(index_col).mean(numeric_only=True)
#     # df_ci.set_index([index_col], inplace=True,
#     #                       append=True, drop=True)
#     df_ci = df_ci * 100
#     df_ci = df_ci.reset_index()
#     df_ci.to_csv(f"{result_dir}/ci_dataframe-2.csv", index=False)
#     fig = go.Figure()
#     for feature in feature_list:
#         feature_lower_col = f"{feature}_lower"
#         feature_upper_col = f"{feature}_upper"
#
#         fig.add_trace(go.Bar(
#             x=df_ci[index_col],
#             y=df_ci[feature],
#             name=feature,
#             error_y=dict(
#                 type='data',
#                 symmetric=True,
#                 array=df_ci[feature_upper_col] * 2 - df_ci[feature] * 2,
#                 arrayminus=df_ci[feature] * 2 - df_ci[feature_lower_col] * 2,
#                 visible=True,
#                 thickness=3,
#                 color="purple",
#             ),
#         ))
#
#     fig.update_layout(
#         title="Total Mitochondrias by Group",
#         xaxis_title=index_col,
#         yaxis_title="Total Mitochondrias",
#         barmode='stack'
#     )
#
#     fig.write_image(f"{result_dir}/bar_plot.pdf", engine="kaleido")
#     fig.show()
#
# # Example usage:
# df = pd.read_csv("processed_data_5.csv")
# compute_and_plot_ci(df, index_col="case", feature_list=["network_fraction", "rod_fraction", "rounded_fraction"], result_dir="results")

#Local Outlier Factor 30_outliers, Local Outlier Factor 30.1, Local Outlier Factor 30, row, column, field, id.1, all_cell_count, cell_count, field_count

# case, id, incubation_condition, disease_stage, id_incubation_condition, incubation_condition_disease_stage
# case, id, incubation_condition, disease_stage, incubation_condition_disease_stage, case_incubation_condition

######## here - multivar tool
# def load_data(file_path):
#     return pd.read_csv(file_path)
#
# def prompt_columns_removal(df):
#     print("Current columns:", df.columns.tolist())
#     cols_to_remove = input("Enter columns to remove (comma-separated) or press Enter to skip: ").split(',')
#     cols_to_remove = [col.strip() for col in cols_to_remove if col.strip() in df.columns]
#     return df.drop(columns=cols_to_remove) if cols_to_remove else df
#
# ## groupby definition
# def prompt_groupby_index(segment_name):
#     groupby = input(f"Do you want to group by the index column for {segment_name}? (yes/no): ").strip().lower()
#     return groupby == 'yes'
#
# def prompt_index_column(df):
#     print("Current columns:", df.columns.tolist())
#     index_col = input("Enter the column to use as index or press Enter to skip: ")
#     return index_col if index_col in df.columns else None
#
# def prompt_secondary_index_columns(df):
#     print("Current columns:", df.columns.tolist())
#     secondary_index_cols = input("Enter any secondary index columns (comma-separated) or press Enter to skip: ").split(',')
#     secondary_index_cols = [col.strip() for col in secondary_index_cols if col.strip() in df.columns]
#     return secondary_index_cols
#
# def zscore_normalize(df, index_col, secondary_index_cols):
#     exclude_cols = [index_col] + secondary_index_cols
#     numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
#     numeric_cols = [col for col in numeric_cols if col not in exclude_cols]
#     df[numeric_cols] = df[numeric_cols].apply(zscore)
#     return df
#
# def create_line_plots(df, index_col, secondary_index_cols, results_dir):
#     os.makedirs(results_dir, exist_ok=True)
#
#     pc_col = 'incubation_condition'
#     if pc_col not in df.columns:
#         print(f"Column '{pc_col}' not found in the dataframe.")
#         return
#
#     for feature in df.columns:
#         if feature not in [index_col, pc_col] + secondary_index_cols:
#             plt.figure()
#             sns.lineplot(data=df, x=index_col, y=feature, hue=pc_col, n_boot = 1000, errorbar=('ci', 90), dashes=True, )
#             plt.title(f"Line Plot for Feature: {feature}")
#             plt.xticks(rotation=45, fontsize=8)
#             plt.savefig(os.path.join(results_dir, f"line_plot_{feature}.png"))
#             plt.close()
#     print("line_plots_done!")
#
#
#
# def bootstrap_ci(data, num_bootstrap=1000, ci=90):
#     """
#     Compute the bootstrap confidence interval.
#     """
#     boot_means = np.random.choice(data, (num_bootstrap, len(data))).mean(axis=1)
#     lower_bound = np.percentile(boot_means, (100 - ci) / 2)
#     upper_bound = np.percentile(boot_means, 100 - (100 - ci) / 2)
#     return lower_bound, upper_bound
#
#
# def create_interactive_line_plots(df, index_col, secondary_index_cols, results_dir):
#     os.makedirs(results_dir, exist_ok=True)
#     pc_col = 'incubation_condition'
#     if pc_col not in df.columns:
#         print(f"Column '{pc_col}' not found in the dataframe.")
#         return
#
#     for feature in df.columns:
#         if feature not in [index_col, pc_col] + secondary_index_cols:
#             fig = make_subplots()
#
#             for pc in df[pc_col].unique():
#                 df_pc = df[df[pc_col] == pc]
#                 means = df_pc.groupby(index_col)[feature].mean()
#                 cis = df_pc.groupby(index_col)[feature].apply(bootstrap_ci)
#
#                 lower_bounds = [ci[0] for ci in cis]
#                 upper_bounds = [ci[1] for ci in cis]
#
#                 fig.add_trace(go.Scatter(
#                     x=means.index,
#                     y=means,
#                     mode='lines+markers',
#                     name=f'{pc_col} {pc}',
#                     line=dict(shape='linear')
#                 ))
#
#                 fig.add_trace(go.Scatter(
#                     x=means.index,
#                     y=lower_bounds,
#                     mode='lines',
#                     line=dict(width=0),
#                     showlegend=False
#                 ))
#
#                 fig.add_trace(go.Scatter(
#                     x=means.index,
#                     y=upper_bounds,
#                     mode='lines',
#                     fill='tonexty',
#                     fillcolor='rgba(0,100,80,0.2)',
#                     line=dict(width=0),
#                     showlegend=False
#                 ))
#
#             fig.update_layout(
#                 title=f"Line Plot for Feature: {feature}",
#                 xaxis_title=index_col,
#                 yaxis_title=feature,
#                 legend_title=pc_col
#             )
#             fig.write_html(os.path.join(results_dir, f"interactive_line_plot_{feature}.html"))
#     print("Interactive line plots done!")
#
# def create_heatmap(df, index_col, secondary_index_cols, results_dir):
#     print("heatmap_began!")
#     df_heat = df.copy()
#     df_heat = prompt_columns_removal(df_heat)
#
#     print("Index column values:", df_heat[index_col].unique())
#     order_x = input(f"Enter the order for {index_col} values (comma-separated) or press Enter to skip: ").split(',')
#     print("Feature columns available:", df_heat.columns)
#     order_y = input("Enter the order for features (comma-separated) or press Enter to skip: ").split(',')
#     order_y = [col.strip() for col in order_y if col.strip() in df_heat.columns]
#
#     df_grouped = df_heat.groupby(index_col).mean().reset_index()
#     if order_x:
#         df_grouped = df_grouped.set_index(index_col).loc[order_x].reset_index()
#     if order_y:
#         df_grouped = df_grouped[[index_col] + order_y]
#
#     lan = plt.rcParams['font.family'] = ['Arial']
#     plt.figure(figsize=(20, 12))
#     font = {'family': 'Arial',
#             'size': 22}
#
#     plt.rc('font', **font)
#     plt.yticks(family='Arial', rotation="horizontal", size=10)
#     plt.xticks(family='Arial', rotation="vertical", size=8)
#     plt.margins(0.1)
#     plt.subplots_adjust(bottom=0.3)
#     sns.heatmap(df_grouped.set_index(index_col), cmap='viridis')
#     plt.title("Heatmap")
#     plt.savefig(os.path.join(results_dir, "heatmap.png"))
#     plt.close()
#     print("Heatmap-Complete!")
#
#     fig = px.imshow(df_grouped.set_index(index_col), title="Interactive Heatmap")
#     fig.write_html(os.path.join(results_dir, "interactive_heatmap.html"))
#     print("Interactive Heatmap-Complete!")
#
#
# ## pca
# def perform_pca(df, index_col, secondary_index_cols, results_dir):
#     print("pca_began!")
#     df_pca = df.copy()
#     df_pca[index_col] = df_pca[index_col].astype(str)
#     df_pca = prompt_columns_removal(df_pca)
#     os.makedirs(results_dir, exist_ok=True)
#
#     pca = PCA(n_components=2)
#     exclude_cols = [index_col] + secondary_index_cols
#     numeric_cols = df_pca.select_dtypes(include=[np.number]).columns.tolist()
#     numeric_cols = [col for col in numeric_cols if col not in exclude_cols]
#     pca_result = pca.fit_transform(df[numeric_cols])
#     explained_variance = pca.explained_variance_ratio_
#
#     unique_labels = df_pca[index_col].unique()
#     colors = pd.factorize(df_pca[index_col])[0]
#
#
#
#     plt.figure(figsize=(20, 12))
#     scatter = plt.scatter(pca_result[:, 0], pca_result[:, 1], c=colors, cmap='viridis')
#     plt.xlabel(f"PC1 ({explained_variance[0] * 100:.2f}%)")
#     plt.ylabel(f"PC2 ({explained_variance[1] * 100:.2f}%)")
#
#     # Get legend elements correctly
#     colors_legend = plt.cm.viridis(np.linspace(0, 1, len(unique_labels)))
#     legend = plt.legend(handles=[mpatches.Circle((0, 0), color=colors_legend[i],
#                                                  label=unique_labels[i])
#                                  for i in range(len(unique_labels))], loc='upper right')
#     legend.set_title('index')
#     plt.title("PCA - All Data")
#     plt.savefig(os.path.join(results_dir, "pca_all.png"))
#     plt.close()
#     print("pc1 done")
#     df_mean = df_pca.groupby(index_col).mean()
#     pca_result_mean = pca.fit_transform(df_mean)
#     explained_variance_mean = pca.explained_variance_ratio_
#
#     plt.figure(figsize=(20, 12))
#     scatter_mean = plt.scatter(pca_result_mean[:, 0], pca_result_mean[:, 1], c=np.arange(len(df_mean)), cmap='viridis')
#     plt.xlabel(f"PC1 ({explained_variance_mean[0] * 100:.2f}%)")
#     plt.ylabel(f"PC2 ({explained_variance_mean[1] * 100:.2f}%)")
#
#     # Get legend elements correctly
#     colors_legend = plt.cm.viridis(np.linspace(0, 1, len(unique_labels)))
#     legend = plt.legend(handles=[mpatches.Circle((0, 0), color=colors_legend[i],
#                                                  label=unique_labels[i])
#                                  for i in range(len(unique_labels))], loc='upper right')
#     legend.set_title('index')
#
#     plt.title("PCA - Grouped by Mean")
#     plt.savefig(os.path.join(results_dir, "pca_mean.png"))
#     plt.close()
#     print("pca_mean done")
#     # Bar plots for top 10 importance features for PC1
#     top_10_indices_pc1 = np.argsort(np.abs(pca.components_[0]))[-10:]
#     plt.figure(figsize=(32, 12))
#     plt.yticks(family='Arial', rotation="horizontal", size=4)
#     plt.barh(np.array(numeric_cols)[top_10_indices_pc1], np.abs(pca.components_[0][top_10_indices_pc1]))
#     plt.xlabel("PCA Feature Importance")
#     plt.title("Top 10 Important Features for PCA - PC1")
#     plt.tight_layout()
#     plt.savefig(os.path.join(results_dir, "pca_top_10_features_pc1.png"))
#     plt.close()
#
#     # Bar plots for top 10 importance features for PC2
#     top_10_indices_pc2 = np.argsort(np.abs(pca.components_[1]))[-10:]
#     plt.figure(figsize=(32, 12))
#     plt.yticks(family='Arial', rotation="horizontal", size=6)
#     plt.barh(np.array(numeric_cols)[top_10_indices_pc2], np.abs(pca.components_[1][top_10_indices_pc2]))
#     plt.xlabel("PCA Feature Importance")
#     plt.title("Top 10 Important Features for PCA - PC2")
#     plt.tight_layout()
#     plt.savefig(os.path.join(results_dir, "pca_top_10_features_pc2.png"))
#     plt.close()
#     print("Top 5 feature importances (PCA - All Data):", np.argsort(np.abs(pca.components_), axis=1)[:, -5:])
#     print("Top 5 feature importances (PCA - Grouped by Mean):", np.argsort(np.abs(pca.components_), axis=1)[:, -5:])
#
#     ##inter
#     pca = PCA(n_components=2)
#     exclude_cols = [index_col] + secondary_index_cols
#     numeric_cols = df_pca.select_dtypes(include=[np.number]).columns.tolist()
#     numeric_cols = [col for col in numeric_cols if col not in exclude_cols]
#     pca_result = pca.fit_transform(df_pca[numeric_cols])
#     explained_variance = pca.explained_variance_ratio_
#
#     fig = px.scatter(x=pca_result[:, 0], y=pca_result[:, 1], color=df_pca[index_col], title="PCA - All Data",
#                      labels={"x": f"PC1 ({explained_variance[0] * 100:.2f}%)", "y": f"PC2 ({explained_variance[1] * 100:.2f}%)"})
#     fig.write_html(os.path.join(results_dir, "interactive_pca_all.html"))
#
#     df_mean = df_pca.groupby(index_col).mean()
#     pca_result_mean = pca.fit_transform(df_mean)
#     explained_variance_mean = pca.explained_variance_ratio_
#
#     fig_mean = px.scatter(x=pca_result_mean[:, 0], y=pca_result_mean[:, 1], color=df_mean.index, title="PCA - Grouped by Mean",
#                           labels={"x": f"PC1 ({explained_variance_mean[0] * 100:.2f}%)", "y": f"PC2 ({explained_variance_mean[1] * 100:.2f}%)"})
#     fig_mean.write_html(os.path.join(results_dir, "interactive_pca_mean.html"))
#
#
# def perform_lda(df, index_col, secondary_index_cols, results_dir):
#     print("lda_began!")
#     df_lda = df.copy()
#     df_lda[index_col] = df_lda[index_col].astype(str)
#     df_lda = prompt_columns_removal(df_lda)
#     os.makedirs(results_dir, exist_ok=True)
#
#     lda = LDA(n_components=2)
#     exclude_cols = [index_col] + secondary_index_cols
#     numeric_cols = df_lda.select_dtypes(include=[np.number]).columns.tolist()
#     numeric_cols = [col for col in numeric_cols if col not in exclude_cols]
#     lda_result = lda.fit_transform(df_lda[numeric_cols], df_lda[index_col])
#     explained_variance = lda.explained_variance_ratio_
#
#     unique_labels = df_lda[index_col].unique()
#     colors = df_lda[index_col].astype('category').cat.codes
#
#     plt.figure(figsize=(20, 12))
#     scatter = plt.scatter(lda_result[:, 0], lda_result[:, 1], c=colors, cmap='viridis')
#     plt.xlabel(f"LD1 ({explained_variance[0] * 100:.2f}%)")
#     plt.ylabel(f"LD2 ({explained_variance[1] * 100:.2f}%)")
#
#     # Get legend elements correctly
#     colors_legend = plt.cm.viridis(np.linspace(0, 1, len(unique_labels)))
#     legend = plt.legend(handles=[mpatches.Circle((0, 0), color=colors_legend[i],
#                                                  label=unique_labels[i])
#                                  for i in range(len(unique_labels))], loc='upper right')
#     legend.set_title('index')
#     plt.title("LDA")
#     plt.savefig(os.path.join(results_dir, "lda.png"))
#     plt.close()
#     print("lda1_done!")
#
#
#     # Bar plots for top 10 importance features for LD1
#     top_10_indices_ld1 = np.argsort(np.abs(lda.coef_[0]))[-10:]
#     plt.figure(figsize=(32, 12))
#     plt.yticks(family='Arial', rotation="horizontal", size=6)
#     plt.barh(np.array(numeric_cols)[top_10_indices_ld1], np.abs(lda.coef_[0][top_10_indices_ld1]))
#     plt.xlabel("LDA Feature Importance")
#     plt.title("Top 10 Important Features for LDA - LD1")
#     plt.tight_layout()
#     plt.savefig(os.path.join(results_dir, "lda_top_10_features_ld1.png"))
#     plt.close()
#
#     # Bar plots for top 10 importance features for LD2 (if applicable)
#     if lda.coef_.shape[0] > 1:
#         top_10_indices_ld2 = np.argsort(np.abs(lda.coef_[1]))[-10:]
#         plt.figure(figsize=(32, 12))
#         plt.yticks(family='Arial', rotation="horizontal", size=6)
#         plt.barh(np.array(numeric_cols)[top_10_indices_ld2], np.abs(lda.coef_[1][top_10_indices_ld2]))
#         plt.xlabel("LDA Feature Importance")
#         plt.title("Top 10 Important Features for LDA - LD2")
#         plt.tight_layout()
#         plt.savefig(os.path.join(results_dir, "lda_top_10_features_ld2.png"))
#         plt.close()
#
#     print("Top 5 feature importances (LDA):", np.argsort(np.abs(lda.coef_), axis=1)[:, -5:])
#
#     lda = LDA(n_components=2)
#     exclude_cols = [index_col] + secondary_index_cols
#     lda_result = lda.fit_transform(df_lda[numeric_cols], df_lda[index_col])
#     explained_variance = lda.explained_variance_ratio_
#
#     fig = px.scatter(x=lda_result[:, 0], y=lda_result[:, 1], color=df_lda[index_col], title="LDA",
#                      labels={"x": f"LD1 ({explained_variance[0] * 100:.2f}%)",
#                              "y": f"LD2 ({explained_variance[1] * 100:.2f}%)"})
#     fig.write_html(os.path.join(results_dir, "interactive_lda.html"))
#
#     print("Top 5 feature importances (LDA):", np.argsort(np.abs(lda.coef_), axis=1)[:, -5:])
#     print("lda_inter1_done!")
#     df_mean_l = df_lda.groupby(index_col).mean()
#     df_mean_l.to_csv(os.path.join(results_dir, "df_mean_l.csv"))
#
#     if prompt_groupby_index("LDA"):
#         print(index_col)
#         types = df_lda[index_col].unique()
#         print(types)
#
#
#         df_lda = df_lda.groupby(index_col).mean().reset_index()
#
#
#         lda = LDA(n_components=2)
#         exclude_cols = [index_col] + secondary_index_cols
#         numeric_cols = df_lda.select_dtypes(include=[np.number]).columns.tolist()
#         numeric_cols = [col for col in numeric_cols if col not in exclude_cols]
#         print(numeric_cols)
#         print(df_lda[index_col])
#         lda_result = lda.fit_transform(df_lda[numeric_cols], df_lda[index_col])
#
#         unique_labels = df_lda[index_col].unique()
#         colors = pd.factorize(df_lda[index_col])[0]
#
#         plt.figure(figsize=(20, 12))
#         scatter = plt.scatter(lda_result[:, 0], lda_result[:, 1], c=colors, cmap='viridis')
#         plt.colorbar(scatter, ticks=range(len(unique_labels)))
#         plt.title(f"LDA: {index_col}")
#         plt.xlabel("LD1")
#         plt.ylabel("LD2")
#         plt.legend(handles=scatter.legend_elements()[0], labels=unique_labels)
#         plt.savefig(os.path.join(results_dir, "LDA-group.png"))
#         plt.close()
#         print("LDA group done!")
#         ## inter
#         lda = LDA(n_components=2)
#         exclude_cols = [index_col] + secondary_index_cols
#         lda_result = lda.fit_transform(df_lda[numeric_cols], df_lda[index_col])
#         explained_variance = lda.explained_variance_ratio_
#
#         fig = px.scatter(x=lda_result[:, 0], y=lda_result[:, 1], color=df_lda[index_col], title="LDA",
#                          labels={"x": f"LD1 ({explained_variance[0] * 100:.2f}%)",
#                                  "y": f"LD2 ({explained_variance[1] * 100:.2f}%)"})
#         fig.write_html(os.path.join(results_dir, "interactive_lda-group.html"))
#
#         print("Top 5 feature importances (LDA):", np.argsort(np.abs(lda.coef_), axis=1)[:, -5:])
#         print("lda_inter-group_done!")
#
#
# def tsne_kmeans_scatter(df, index_col, secondary_index_cols, results_dir, tsne_components=2, random_state=42):
#     print("K-Means started!")
#     df_kmeans = df.copy()
#     df_kmeans[index_col] = df_kmeans[index_col].astype(str)
#     df_kmeans = prompt_columns_removal(df_kmeans)
#     os.makedirs(results_dir, exist_ok=True)
#     n_clusters = int(input("Enter the number of clusters for K-Means: "))
#     exclude_cols = [index_col] + secondary_index_cols
#     numeric_cols = df_kmeans.select_dtypes(include=[np.number]).columns.tolist()
#     numeric_cols = [col for col in numeric_cols if col not in exclude_cols]
#
#
#     #feature matrix
#     X = df_kmeans[numeric_cols].values
#
#     # t-SNE for dimensionality reduction
#     tsne = TSNE(n_components=tsne_components, random_state=random_state)
#     X_tsne = tsne.fit_transform(X)
#
#     # Apply K-Means clustering
#     kmeans = KMeans(n_clusters=n_clusters, random_state=random_state)
#     kmeans_result = kmeans.fit_predict(X_tsne)
#     df_kmeans['cluster'] = kmeans_result
#     df_kmeans.to_csv(os.path.join(results_dir, "kmeans_clustering_results-group.csv"))
#
#     # Add cluster and t-SNE dimensions to the DataFrame
#     #df_kmeans['clusters'] = kmeans_result
#     df_kmeans['tsne1'] = X_tsne[:, 0]
#     if tsne_components > 1:
#         df_kmeans['tsne2'] = X_tsne[:, 1]
#     if tsne_components > 2:
#         df_kmeans['tsne3'] = X_tsne[:, 2]
#
#     # Create a scatter plot
#     plt.figure(figsize=(32, 18))
#     sns.scatterplot(
#         x='tsne1', y='tsne2', hue='cluster', style=index_col,
#         data=df_kmeans, palette='viridis', s=100
#     )
#     plt.title(f't-SNE with K-Means Clustering:{index_col}')
#     plt.tight_layout()
#     plt.savefig(os.path.join(results_dir, "KMeans-group.png"))
#     plt.close()
#     print("K-Means Clustering done!")
#     ## groupby
#     if prompt_groupby_index("K-Means Clustering"):
#         df_kmeans = df_kmeans.groupby(index_col).mean().reset_index()
#         # feature matrix
#         X = df_kmeans[numeric_cols].values
#
#         # t-SNE for dimensionality reduction
#         tsne = TSNE(n_components=tsne_components, random_state=random_state)
#         X_tsne = tsne.fit_transform(X)
#
#         # Apply K-Means clustering
#         kmeans = KMeans(n_clusters=n_clusters, random_state=random_state)
#         kmeans_result = kmeans.fit_predict(X_tsne)
#         df_kmeans['cluster'] = kmeans_result
#         df_kmeans.to_csv(os.path.join(results_dir, "kmeans_clustering_results-group.csv"))
#
#         # Add cluster and t-SNE dimensions to the DataFrame
#         #df_kmeans['clusters'] = kmeans_result
#         df_kmeans['tsne1'] = X_tsne[:, 0]
#         if tsne_components > 1:
#             df_kmeans['tsne2'] = X_tsne[:, 1]
#         if tsne_components > 2:
#             df_kmeans['tsne3'] = X_tsne[:, 2]
#
#         # Create a scatter plot
#         plt.figure(figsize=(32, 18))
#         sns.scatterplot(
#             x='tsne1', y='tsne2', hue='cluster', style=index_col,
#             data=df_kmeans, palette='viridis', s=100
#         )
#         plt.title(f't-SNE with K-Means Clustering:{index_col}')
#         plt.tight_layout()
#         plt.savefig(os.path.join(results_dir, "KMeans-group.png"))
#         plt.close()
#         print("K-Means Clustering done!")
#
#
# def perform_correlation_heatmap(df, results_dir):
#     print("correlation_began!")
#     df_cor = df.copy()
#     print("keep only numeric columns")
#     df_cor = prompt_columns_removal(df_cor)
#     os.makedirs(results_dir, exist_ok=True)
#     df_cor.corr().to_csv(os.path.join(results_dir, "correlation_results.csv"))
#     lan = plt.rcParams['font.family'] = ['Arial']
#     plt.figure(figsize=(20, 12))
#     font = {'family': 'Arial',
#             'size': 12}
#
#     plt.rc('font', **font)
#     #cax = plt.gcf().axes[0]
#     #cax.tick_params(labelsize=6)
#     plt.yticks(family='Arial', rotation="horizontal", size=10)
#     plt.xticks(family='Arial', rotation="vertical", size=8)
#     plt.margins(0.1)
#     plt.subplots_adjust(bottom=0.3)
#     sns.heatmap(df_cor.corr(), annot=True, cmap='viridis')
#     plt.title("Correlation Heatmap")
#     plt.tight_layout()
#     plt.savefig(os.path.join(results_dir, "correlation_heatmap.png"))
#     plt.close()
#     print("Correlation Heatmap-Complete!")
#
# def perform_hierarchical_clustering(df, index_col, secondary_index_cols, results_dir):
#     print("Hierarchical Clustering started!")
#     df_hc = df.copy()
#     df_hc[index_col] = df_hc[index_col].astype(str)
#     df_hc = prompt_columns_removal(df_hc)
#     os.makedirs(results_dir, exist_ok=True)
#     linkage_method = input("Enter the linkage method for Hierarchical Clustering (ward, complete, average, single): ")
#     if linkage_method not in ['ward', 'complete', 'average', 'single']:
#         print("Invalid linkage method. Defaulting to 'ward'.")
#         linkage_method = 'ward'
#
#     exclude_cols = [index_col] + secondary_index_cols
#     numeric_cols = df_hc.select_dtypes(include=[np.number]).columns.tolist()
#     numeric_cols = [col for col in numeric_cols if col not in exclude_cols]
#
#     linked = linkage(df_hc[numeric_cols], method=linkage_method)
#     plt.figure(figsize=(20, 12))
#     dendrogram(linked, orientation='top', labels=df_hc[index_col].tolist(), distance_sort='descending')
#     plt.title(f"Hierarchical Clustering: {index_col}")
#     plt.xticks(family='Arial', rotation="vertical", size=6)
#     plt.savefig(os.path.join(results_dir, "HierarchicalClustering.png"))
#     plt.close()
#     print("Hierarchical Clustering done!")
#
#     if prompt_groupby_index("Hierarchical Clustering"):
#         df_hc = df_hc.groupby(index_col).mean().reset_index()
#
#         linkage_method = input("Enter the linkage method for Hierarchical Clustering (ward, complete, average, single): ")
#         if linkage_method not in ['ward', 'complete', 'average', 'single']:
#             print("Invalid linkage method. Defaulting to 'ward'.")
#             linkage_method = 'ward'
#
#         exclude_cols = [index_col] + secondary_index_cols
#         numeric_cols = df_hc.select_dtypes(include=[np.number]).columns.tolist()
#         numeric_cols = [col for col in numeric_cols if col not in exclude_cols]
#
#         linked = linkage(df_hc[numeric_cols], method=linkage_method)
#         plt.figure(figsize=(20, 12))
#         dendrogram(linked, orientation='top', labels=df_hc[index_col].tolist(), distance_sort='descending')
#         plt.title(f"Hierarchical Clustering: {index_col}")
#         plt.xticks(family='Arial', rotation="vertical", size=6)
#         plt.savefig(os.path.join(results_dir, "HierarchicalClustering-group.png"))
#         plt.close()
#         print("Hierarchical Clustering done!")
#
#
# def bootstrap_ci_2(data, n_iterations=1000, ci=95):
#     means = []
#     for _ in range(n_iterations):
#         sample = np.random.choice(data, size=len(data), replace=True)
#         means.append(np.mean(sample))
#     lower_bound = np.percentile(means, (100 - ci) / 2)
#     upper_bound = np.percentile(means, 100 - (100 - ci) / 2)
#     return lower_bound, upper_bound
#
#
# def compute_and_plot_ci(df, index_col, result_dir, feature_list):
#     df_temp = df.copy()
#     ci_dict = {}
#     for feature in feature_list:
#         ci_dict[f"{feature}_lower"] = []
#         ci_dict[f"{feature}_upper"] = []
#
#     for idx in df[index_col].unique():
#         subset = df[df[index_col] == idx]
#         for feature in feature_list:
#             lower_bound, upper_bound = bootstrap_ci_2(subset[feature])
#             ci_dict[f"{feature}_lower"].append(lower_bound)
#             ci_dict[f"{feature}_upper"].append(upper_bound)
#
#     for feature in feature_list:
#         df_temp[f"{feature}_lower"] = df_temp[index_col].map(
#             dict(zip(df_temp[index_col].unique(), ci_dict[f"{feature}_lower"]))
#         )
#         df_temp[f"{feature}_upper"] = df_temp[index_col].map(
#             dict(zip(df_temp[index_col].unique(), ci_dict[f"{feature}_upper"]))
#         )
#
#     df_temp.to_csv(f"{result_dir}/ci_dataframe.csv", index=False)
#     df_ci = df_temp.groupby(index_col).mean(numeric_only=True)
#     # df_ci.set_index([index_col], inplace=True,
#     #                       append=True, drop=True)
#     df_ci = df_ci * 100
#     df_ci = df_ci.reset_index()
#     df_ci.to_csv(f"{result_dir}/ci_dataframe-2.csv", index=False)
#     fig = go.Figure()
#     for feature in feature_list:
#         feature_lower_col = f"{feature}_lower"
#         feature_upper_col = f"{feature}_upper"
#
#         fig.add_trace(go.Bar(
#             x=df_ci[index_col],
#             y=df_ci[feature],
#             name=feature,
#             error_y=dict(
#                 type='data',
#                 symmetric=True,
#                 array=df_ci[feature_upper_col] * 2 - df_ci[feature] * 2,
#                 arrayminus=df_ci[feature] * 2 - df_ci[feature_lower_col] * 2,
#                 visible=True,
#                 thickness=3,
#                 color="purple",
#             ),
#         ))
#
#     fig.update_layout(
#         title="Total Mitochondrias by Group",
#         xaxis_title=index_col,
#         yaxis_title="Total Mitochondrias",
#         barmode='stack'
#     )
#
#     fig.write_image(f"{result_dir}/fraction_bar_plot.pdf", engine="kaleido")
#     fig.write_html(f"{result_dir}/fraction_bar_plot.htmal")
#     fig.show()
#
#
#
# def main(file_path, results_dir):
#     df = load_data(file_path)
#     # make results folder
#     os.makedirs(results_dir, exist_ok=True)
#     df = prompt_columns_removal(df)
#     index_col = prompt_index_column(df)
#     secondary_index_cols = prompt_secondary_index_columns(df)
#     compute_and_plot_ci(df, index_col, results_dir,
#                         feature_list=["network_fraction", "rod_fraction", "rounded_fraction"])
#
#     if index_col:
#         df = zscore_normalize(df, index_col, secondary_index_cols)
#
#     df.to_csv(os.path.join(results_dir, "normalized_data.csv"))
#     create_line_plots(df, index_col, secondary_index_cols, os.path.join(results_dir, "line_plots"))
#     create_interactive_line_plots(df, index_col, secondary_index_cols, os.path.join(results_dir, "interactive_line_plots"))
#     #df = prompt_columns_removal(df)
#     print("choose params for heatmap(keep only index column")
#     index_col = prompt_index_column(df)
#     secondary_index_cols = prompt_secondary_index_columns(df)
#     create_heatmap(df, index_col, secondary_index_cols, results_dir)
#     #create_interactive_heatmap(df, index_col, secondary_index_cols, results_dir)
#     print("choose params for pca(keep only index column")
#     #df = prompt_columns_removal(df)
#     index_col = prompt_index_column(df)
#     secondary_index_cols = prompt_secondary_index_columns(df)
#     perform_pca(df, index_col, secondary_index_cols, results_dir)
#     #perform_interactive_pca(df, index_col, secondary_index_cols, results_dir)
#     print("choose params for lda(keep only index column")
#     #df = prompt_columns_removal(df)
#     index_col = prompt_index_column(df)
#     secondary_index_cols = prompt_secondary_index_columns(df)
#     perform_lda(df, index_col, secondary_index_cols, results_dir)
#     #perform_interactive_lda(df, index_col, secondary_index_cols, results_dir)
#     #df = prompt_columns_removal(df)
#     print("choose params for kmeans(keep only index column")
#     index_col = prompt_index_column(df)
#     secondary_index_cols = prompt_secondary_index_columns(df)
#     tsne_kmeans_scatter(df, index_col, secondary_index_cols, results_dir)
#     print("choose params for correlation(no need for index column")
#     index_col = prompt_index_column(df)
#     secondary_index_cols = prompt_secondary_index_columns(df)
#     perform_correlation_heatmap(df, results_dir)
#     print("choose params for H_C(keep only index column")
#     index_col = prompt_index_column(df)
#     secondary_index_cols = prompt_secondary_index_columns(df)
#     perform_hierarchical_clustering(df, index_col, secondary_index_cols, results_dir)
#
# if __name__ == "__main__":
#     file_path = input("Enter the path to the CSV file: ")
#     results_dir = input("Enter the path to the results directory: ")
#     main(file_path, results_dir)
## todo add html save
## todo replacamnet to lda - pls ?
## todo whiskers above plot with sum.
## todo do not include the boot CI in df!!!!



###MA
## TMRE
## text function imp- 140224-final-use this! :)
# def process_folders(main_folder):
#     for subfolder_name in os.listdir(main_folder):
#         subfolder_path = os.path.join(main_folder, subfolder_name)
#         if os.path.isdir(subfolder_path):
#             print(f"Processing subfolder: {subfolder_name}")
#
#             # Prompt user for the value for the added column
#             pc_value = input("Enter the value for the 'PC' column:")
#
#             # Iterate through text files in the subfolder
#             for txt_file in glob.glob(os.path.join(subfolder_path, '*.txt')):
#                 df = pd.read_csv(txt_file, skiprows=8, engine='python', sep='\t')
#
#                 # Add PC column with specified value
#                 df['PC'] = pc_value
#                 # Extract plate name from the CSV file name
#                 plate_name = os.path.splitext(os.path.basename(txt_file))[0]
#                 print(f"Processing plate: {plate_name}")
#
#                 # Create output folder if it doesn't exist
#                 output_folder = os.path.join(main_folder, "results", subfolder_name)
#                 os.makedirs(output_folder, exist_ok=True)
#
#                 # Convert and save as CSV
#                 csv_filename = plate_name + '_csv.csv'
#
#                 # Ask the user for flagged wells to remove
#                 flagged_wells_input = input(f"Enter flagged wells to remove for plate '{plate_name}' (row_column, ...) or 'none' if none: ")
#                 if flagged_wells_input.lower() == 'none':
#                     flagged_wells = []
#                 else:
#                     flagged_wells = flagged_wells_input.split(',')
#
#                 # Remove flagged wells from DataFrame
#                 for flagged_well in flagged_wells:
#                     if flagged_well != '':
#                         row, column = flagged_well.split('_')
#                         row = int(row)
#                         column = int(column)
#                         df = df.loc[~((df['Row'] == row) & (df['Column'] == column))]
#
#                 # Save the modified DataFrame to CSV
#                 csv_path = os.path.join(output_folder, csv_filename)
#                 df.to_csv(csv_path, index=None)
#                 print(f"File '{csv_filename}' saved in folder '{output_folder}'")
#
#     # Custom message after processing all subfolders
#     print("All subfolders processed. Ready to continue.")
# # #
# # # # Run the function
# main_folder_path = 'MA-TMRE'
# process_folders(main_folder_path)
# #
# # # User to signal when ready to proceed
# input("Please combine separated csv files and Press Enter when ready to continue...")
#
# ## inter plate combine
# def process_data_combine(folder_path):
#     # Create output result folder
#     result_combine_dir = os.path.join(os.getcwd(), 'MA-TMRE/MA-TMRE_result_combined/final')
#     if not os.path.exists(result_combine_dir):
#         os.makedirs(result_combine_dir)
#
#     #     # List immediate subdirectories
#     subdirectories = [subdir for subdir in os.listdir(folder_path) if os.path.isdir(os.path.join(folder_path, subdir))]
#
#     combined_dfs = []  # List to store DataFrames from all subdirectories
#
#     # Process each immediate subdirectory
#     for subdir in subdirectories:
#         sub_dir_path = os.path.join(folder_path, subdir)
#         print("Processing files in:", sub_dir_path)
#         dfs = []
#         for file in os.listdir(sub_dir_path):
#             if file.endswith('.csv'):
#                 file_path = os.path.join(sub_dir_path, file)
#                 df = pd.read_csv(file_path, encoding='unicode_escape')
#                 df['plate'] = os.path.splitext(file)[0]  # Add plate column
#                 dfs.append(df)
#
#
#                 # Combine DataFrames from this subdirectory
#     combined_df = pd.concat(dfs, ignore_index=True)
#     combined_dfs.append(combined_df)  # Append to the list of DataFrames
#
#     # Combine all DataFrames from all subdirectories
#     combined_df = pd.concat(combined_dfs, ignore_index=True)
#     print("Data frames combined!")
#     print(combined_df)
#     # 'plate' & 'well' column to index_columns
#     combined_df['well'] = combined_df['Row'].astype('str') + combined_df['Column'].astype('str')
#     combined_df['well'] = combined_df['well'].astype('str')
#     # Save combined DataFrame
#     combined_df.to_csv(os.path.join(result_combine_dir, "combined_df_data.csv"), index=False)
#
# folder_path = "MA-TMRE/results"
# process_data_combine(folder_path)
#
# # # User to signal when ready to proceed
# input("Please combine separated csv files and Press Enter when ready to continue...")

## new - use this!!! ma tmre
# def process_plate_data_on_folders(main_folder, unwanted_columns, important_features, index_columns):
#     #iterating each item in main input folder
#     for item_name in os.listdir(main_folder):
#         item_path = os.path.join(main_folder, item_name)
#         if os.path.isdir(item_path):
#             #iterate each csv in each subfolder
#             for file_name in os.listdir(item_path):
#                 if file_name.endswith(".csv"):
#                     file_path = os.path.join(item_path, file_name)
#
#                     #loading csv
#                     plate_df = pd.read_csv(file_path)
#                     print(f"reading : {plate_df}")
#                     plate_df.columns = plate_df.columns.str.strip()
#                     #open results folder dir with the name of csv
#                     results_folder = os.path.join(item_path, f"results_{file_name}_res")
#                     os.makedirs(results_folder, exist_ok=True)
#                     print("results folder generated!")
#                     #drop unwanted columns & start proc
#                     plate_df = plate_df.drop(columns=unwanted_columns)
#                     print("unwanted columns dropped!")
#
#                     #'well' column added to index_columns
#                     #plate_df['well'] = plate_df['Row'].astype('str') + plate_df['Column'].astype('str')
#                     #plate_df['well'] = plate_df['well'].astype('str')
#                     # index_columns.append('well')
#                     # index_columns.append('plate')
#                     #creating all index options
#                     new_index_columns = []
#                     for r in range(1, len(index_columns) + 1):
#                         for comb in itertools.combinations(index_columns, r):
#                             new_index_columns.append(comb)
#                     for cols in new_index_columns:
#                         col_name = '_'.join(cols)
#                         plate_df[col_name] = plate_df[list(cols)].apply(lambda x: '_'.join(map(str, x)), axis=1)
#                     print("indexing complete!")
#                     #prompting feature removal
#                     features_to_remove = input("Enter comma-separated column names to remove (or 'none'): ").strip()
#                     if features_to_remove.lower() != 'none':
#                         features_to_remove = [col.strip() for col in features_to_remove.split(',')]
#                         plate_df.drop(columns=features_to_remove, inplace=True)
#                         print("Features removed:", features_to_remove)
#                     ##prompting ids removal
#                     remove_id = input("Enter ID value to remove (or 'none'): ").strip()
#                     if remove_id.lower() != 'none':
#                         plate_df = plate_df[~plate_df['ID'].astype(str).str.contains(remove_id)]
#                         print(f"Rows with ID '{remove_id}' removed.")
#                     print(" main proc done!")
#                     #removing rows with 90% or more zero/na values
#                     plate_df = plate_df.dropna(thresh=0.1 * len(plate_df), axis=0)
#                     plate_df = plate_df.replace(0, np.nan)
#                     plate_df = plate_df.dropna(thresh=0.1 * len(plate_df.columns), axis=1)
#                     print(new_index_columns)
#                     #prompting index column for statistical summary
#                     index_col = input("Enter the index column for statistical summary: ")
#                     summary_df = plate_df.groupby(index_col).describe()
#                     summary_df.to_excel(os.path.join(results_folder, "statistical_summary.xlsx"))
#                     print("Statistical summary saved")
#                     #save first df
#                     plate_df.to_csv(os.path.join(results_folder, "df_proc_1.csv"), index=False)
#                     #generating histogram for cell count
#                     fig = px.histogram(plate_df, x="all_cells - Number of Objects", title="Histogram of Cell Count")
#                     histogram_file_path = os.path.join(results_folder, "histogram.html")
#                     fig.write_html(histogram_file_path)
#                     #statistics for cell count
#                     stats = plate_df["all_cells - Number of Objects"].describe()
#                     stats_file_path = os.path.join(results_folder, "cell_count-statistics.txt")
#                     with open(stats_file_path, "w") as f:
#                         f.write(stats.to_string())
#                     print("take a look at cell count stats and hist...")
#                     ##prompting minimum and maximum values for cell count
#                     min_value = float(input("Enter the minimum value for cell count: "))
#                     max_value = float(input("Enter the maximum value for cell count: "))
#                     #excluding cell count outside the specified range
#                     plate_df = plate_df[(plate_df["all_cells - Number of Objects"] >= min_value) & (plate_df["all_cells - Number of Objects"] <= max_value)]
#                     print("low/high cell counts removed")
#                     plate_df.to_csv(os.path.join(results_folder, "df_proc_2.csv"))
#                     print(new_index_columns)
#                     #prompting index column for box plots
#                     index_col = input("Enter the index column for box plot vis: ")
#                     bar_df = plate_df.copy()
#                     # box plots for selected index- selected features only
#                     print("box plotting selected features!")
#                     for col in bar_df:
#                         if col in important_features:
#                             bp = px.box(bar_df, x=index_col, y=bar_df[col], points="all", color=index_col, notched=True)
#                             bp.update_traces(quartilemethod="inclusive")  # or "inclusive", or "linear" by default
#                             bp.update_layout(
#                                 font_family="Arial",
#                                 font_color="Black",
#                                 font_size=20,
#                                 font=dict(
#                                     family="Arial",
#                                     size=20,
#                                     color="Black"
#                                 )
#                             )
#                             print("one_down")
#                             bp.write_html(os.path.join(results_folder, f'{col}_bar_plot.html'))
#                             bp.write_image(os.path.join(results_folder, f"box_plot{col}.pdf"), engine="kaleido")
#                             bp.show()
#                     print("box plots done!")
#                     #creating excel  with sheets with split data and summary stats
#                     print("creating excel files and statistical summary & ppt for combination of groups!")
#                     group_columns = ['Cell Type', 'Compound', 'Concentration']
#                     excel_file_path = os.path.join(results_folder, f'summary.xlsx')
#                     with pd.ExcelWriter(excel_file_path, engine='openpyxl') as writer:
#                         for name, group in plate_df.groupby(group_columns):
#                             try:
#                                 sheet_name = f"{name[0]}_{name[1]}_{name[2]}"
#                                 group.to_excel(writer, sheet_name=sheet_name, index=False)
#                                 ####summary statistics to another sheet
#                                 summary_df = group.describe(include='all').T[
#                                     ['mean', 'std', 'min', '25%', '50%', '75%', 'max']]
#                                 summary_df.to_excel(writer, sheet_name=f"{sheet_name}_Summary")
#                             except ValueError:
#                                 pass
#
#                         if not writer.sheets:
#
#                             pd.DataFrame({'Dummy': [1]}).to_excel(writer, sheet_name='DummySheet', index=False)
#
#                     #ppt with slides for each unique combination
#                     print("ppt file generated!")
#                     ppt_file_path = os.path.join(results_folder, f'slides.pptx')
#                     prs = Presentation()
#                     for name, _ in plate_df.groupby(group_columns):
#                         slide = prs.slides.add_slide(prs.slide_layouts[5])
#                         title = slide.shapes.title
#                         title.text = f"Combination: {name[0]}, {name[1]}, {name[2]}"
#                     prs.save(ppt_file_path)
#                     #bar plot for Number of Cells
#                     print("working on cell count vis!")
#                     group_columns = ['Cell Type', 'Compound', 'Concentration']
#                     plt.figure(figsize=(20, 18))
#                     for name, group in plate_df.groupby(group_columns):
#                         try:
#                             plt.bar(str(name), group['all_cells - Number of Objects'].mean(),
#                                     yerr=group['all_cells - Number of Objects'].std(), capsize=5,
#                                     label=str(name))
#                         except ValueError:
#                             pass
#                     plt.title(f'Bar Plot for Number of Cells - {file_name}')
#                     plt.xlabel('Cell Type, Compound, Concentration')
#                     plt.ylabel('Number of Cells')
#                     plt.xticks(rotation=45, ha='right')
#                     plt.legend()
#                     plt.savefig(os.path.join(results_folder, f'number_of_cells_bar_plot.pdf'))
#                     plt.close()
#
#                     #heatmap for selected features (mean values & normalized)
#                     heat_df = plate_df.copy()
#                     index_col = input("Enter the index column for statistical summary: ")
#                     print("heatmap prep!")
#                     heat_df = heat_df.set_index(index_col, drop=True)
#                     numeric_features = heat_df[important_features].select_dtypes(include=np.number)
#                     numeric_features_normalized = (numeric_features - numeric_features.mean()) / numeric_features.std()
#                     numeric_features_normalized = numeric_features_normalized.groupby(index_col).mean()
#                     plt.figure(figsize=(48, 28))
#                     sns.heatmap(numeric_features_normalized, annot=True, cmap='coolwarm')
#                     plt.title(f'Heatmap for Selected Numeric Features - {file_name}')
#                     plt.savefig(os.path.join(results_folder, f'heatmap.pdf'))
#                     plt.close()
#                     print("heatmap done!")
#
#                     perform_outlier_detection = input("Do you want to perform outlier detection? (yes/no): ").strip().lower()
#                     if perform_outlier_detection == 'yes':
#                         print("Applying basic Outlier methods...")
#
#
#
#
#
#                         #prompting index column for outliers detecetion
#                         print("Columns available for outlier detection:")
#                         print(plate_df.columns)
#                         index_column = input("Choose one column as index for outlier detection: ")
#                         out_dat_init = plate_df.copy()
#                         # Define the outlier detection functions
#                         def compute_percentiles(df, output_file):
#                             up_bounds = []
#                             low_bounds = []
#                             above_count = []
#                             below_count = []
#                             numeric_cols = df.select_dtypes(include='number').columns
#                             for col in numeric_cols:
#                                 percentile_low = df[col].quantile(0.01)
#                                 percentile_high = df[col].quantile(0.99)
#                                 up_bound = percentile_high
#                                 low_bound = percentile_low
#                                 above = df[df[col] > percentile_high]
#                                 below = df[df[col] < percentile_low]
#                                 above_count.append(len(above))
#                                 below_count.append(len(below))
#                                 df = df[(df[col] >= percentile_low) & (df[col] <= percentile_high)]
#                                 up_bounds.append(up_bound)
#                                 low_bounds.append(low_bound)
#                             output = pd.DataFrame(
#                                 {'Column': numeric_cols, 'Up Bound': up_bounds, 'Low Bound': low_bounds, 'Above Count': above_count,
#                                  'Below Count': below_count})
#                             output.to_csv(output_file, index=False)
#                             return df
#
#                         def compute_iqr(df, output_file):
#                             up_bounds = []
#                             low_bounds = []
#                             above_count = []
#                             below_count = []
#                             numeric_cols = df.select_dtypes(include='number').columns
#                             for col in numeric_cols:
#                                 percentile25th = df[col].quantile(0.25)
#                                 percentile75th = df[col].quantile(0.75)
#                                 iqr = percentile75th - percentile25th
#                                 up_bound = percentile75th + 1.5 * iqr
#                                 low_bound = percentile25th - 1.5 * iqr
#                                 above = df[df[col] > up_bound]
#                                 below = df[df[col] < low_bound]
#                                 above_count.append(len(above))
#                                 below_count.append(len(below))
#                                 df = df[(df[col] >= low_bound) & (df[col] <= up_bound)]
#                                 up_bounds.append(up_bound)
#                                 low_bounds.append(low_bound)
#                             output = pd.DataFrame(
#                                 {'Column': numeric_cols, 'Up Bound': up_bounds, 'Low Bound': low_bounds, 'Above Count': above_count,
#                                  'Below Count': below_count})
#                             output.to_csv(output_file, index=False)
#                             return df
#
#                         def compute_scores(df, output_file):
#                             up_scores = []
#                             low_scores = []
#                             above_count = []
#                             below_count = []
#                             numeric_cols = df.select_dtypes(include='number').columns
#                             for col in numeric_cols:
#                                 up_score = df[col].mean() + 3 * df[col].std()
#                                 low_score = df[col].mean() - 3 * df[col].std()
#                                 above = df[df[col] > up_score]
#                                 below = df[df[col] < low_score]
#                                 above_count.append(len(above))
#                                 below_count.append(len(below))
#                                 df = df[(df[col] >= low_score) & (df[col] <= up_score)]
#                                 up_scores.append(up_score)
#                                 low_scores.append(low_score)
#                             output = pd.DataFrame(
#                                 {'Column': numeric_cols, 'Up Score': up_scores, 'Low Score': low_scores, 'Above Count': above_count,
#                                  'Below Count': below_count})
#                             output.to_csv(output_file, index=False)
#                             return df
#
#                         def LOF_outlier(df, output_file1, output_file2):
#                             # outlier detection using LOF
#                             types = df[index_column].unique()
#                             print(types)
#                             mask = []
#                             features = df.select_dtypes(include='number').columns
#
#                             detector_list = [
#                                 ("Local Outlier Factor 30", LocalOutlierFactor(n_neighbors=30))
#                             ]
#
#                             for name, algorithm in detector_list:
#                                 errors = np.full(len(df), fill_value=np.nan)
#                                 outliers = np.full(len(df), fill_value=np.nan)
#
#                                 for type in types:
#                                     x = df.loc[:, features].values
#                                     F = x.sum(1)
#                                     mask = np.zeros(x.shape[0])
#                                     mask[np.isfinite(F)] = 1
#                                     mask_type = mask * np.array(df[index_column] == type)
#                                     Curr_df = df.loc[mask_type == 1, features]
#                                     x = Curr_df.values
#
#                                     if name == "Local Outlier Factor 30":
#                                         algorithm.fit(x)
#                                         errors[mask_type == 1] = algorithm.negative_outlier_factor_
#                                         outliers[mask_type == 1] = algorithm.fit_predict(x)
#
#                                 df[name] = errors
#                                 df[f'{name}_outliers'] = outliers
#                                 df.set_index(name, inplace=True,
#                                                   append=True, drop=False)
#                                 df.to_csv(os.path.join(output_file1))
#                                 # exclude rows that were defined as outliers
#                                 for col in df.columns:
#                                     if col.endswith("_outliers"):
#                                         df = df[df[col] != -1]
#                             df.to_csv(os.path.join(output_file2))
#                             print("outliers removed!")
#                             return df
#
#
#                         # Apply the outlier detection functions on filtered_df
#                         df_percentiles = compute_percentiles(out_dat_init, os.path.join(results_folder, 'percentiles_outliers.csv'))
#                         df_iqr = compute_iqr(out_dat_init, os.path.join(results_folder, 'iqr_outliers.csv'))
#                         df_scores = compute_scores(out_dat_init, os.path.join(results_folder, 'scores_outliers.csv'))
#                         df_lof = LOF_outlier(out_dat_init, os.path.join(results_folder, 'processed_data_outliers_1.csv'), os.path.join(results_folder, 'processed_data_outliers_2.csv') )
#
#                         print("Please review outlier detection and choose one method to continue with:")
#                         method_choice = input("Enter 1 for percentiles, 2 for IQR, 3 for scores, 4 for lof: ")
#
#                         if method_choice == '1':
#                             out_dat = df_percentiles
#                         elif method_choice == '2':
#                             out_dat = df_iqr
#                         elif method_choice == '3':
#                             out_dat = df_scores
#                         elif method_choice == '4':
#                             out_dat = df_lof
#                         else:
#                             print("Invalid choice. Defaulting to original DataFrame.")
#                             out_dat = out_dat_init
#
#                         # Save the final filtered DataFrame
#                         final_file_path = os.path.join(results_folder, 'final_filtered_data.csv')
#                         out_dat.to_csv(final_file_path, index=False)
#
#                         print("Final filtered DataFrame saved to", final_file_path)
#                     #     #outlier detection using LOF
#                     #     types = out_dat[index_column].unique()
#                     #     print(types)
#                     #     mask = []
#                     #     features = out_dat.select_dtypes(include='number').columns
#                     #
#                     #     detector_list = [
#                     #             ("Local Outlier Factor 30", LocalOutlierFactor(n_neighbors=30))
#                     #         ]
#                     #
#                     #     for name, algorithm in detector_list:
#                     #         errors = np.full(len(out_dat),fill_value=np.nan)
#                     #         outliers = np.full(len(out_dat),fill_value=np.nan)
#                     #
#                     #         for type in types:
#                     #             x = out_dat.loc[:,features].values
#                     #             F = x.sum(1)
#                     #             mask = np.zeros(x.shape[0])
#                     #             mask[np.isfinite(F)] = 1
#                     #             mask_type = mask * np.array(out_dat[index_column] == type)
#                     #             Curr_df = out_dat.loc[mask_type == 1, features]
#                     #             x = Curr_df.values
#                     #
#                     #
#                     #             if name == "Local Outlier Factor 30":
#                     #                 algorithm.fit(x)
#                     #                 errors[mask_type==1]  = algorithm.negative_outlier_factor_
#                     #                 outliers[mask_type==1] = algorithm.fit_predict(x)
#                     #
#                     #         out_dat[name] = errors
#                     #         out_dat[f'{name}_outliers'] = outliers
#                     #         out_dat.set_index(name, inplace=True,
#                     #                                                    append=True, drop=False)
#                     #         out_dat.to_csv(os.path.join(results_folder, "processed_data_outliers_1.csv"))
#                     #         #exclude rows that were defined as outliers
#                     #         for col in out_dat.columns:
#                     #             if col.endswith("_outliers"):
#                     #                 out_dat = out_dat[out_dat[col] != -1]
#                     #     out_dat.to_csv(os.path.join(results_folder, "processed_data_outliers_2.csv"))
#                     #     print("outliers removed!")
#                     # else:
#                     #     print("conitinuing wiothout outlier removal")
#                     #     out_dat = plate_df.copy()
#                     #prompting index column for final proc- "PC" is must!
#                     print(out_dat.columns)
#                     columns_to_set_index = input("Enter columns to set as indexes (comma-separated & PC is must!): ").split(',')
#                     columns_to_set_index = [col.strip() for col in columns_to_set_index]
#                     if columns_to_set_index:
#                         out_dat = out_dat.set_index(columns_to_set_index, drop = False)
#                     #prompting index column for feature QC
#                     print("Select the correct index title to be the target variable:")
#                     for idx, index_title in enumerate(out_dat.index.names):
#                         print(f"{idx + 1}: {index_title}")
#                     selected_index = int(input("Enter the number corresponding to the correct index title: ")) - 1
#                     out_dat = out_dat.select_dtypes(include=np.number)
#                     #set the target variable
#                     feat_df = out_dat[important_features]
#                     target_variable = feat_df.index.names[selected_index]
#                     print(f"Target variable selected: {target_variable}")
#                     target_variable_labels, _ = pd.factorize(feat_df.index.get_level_values(target_variable))
#                     print(target_variable_labels)
#                     #opening directory for feature QC
#                     feature_qc_dir = os.path.join(results_folder, "feature_qc")
#                     os.makedirs(feature_qc_dir, exist_ok=True)
#                     #1- statistical sum csv
#                     statistical_summary = feat_df.describe()
#                     statistical_summary.to_csv(os.path.join(feature_qc_dir, "statistical_summary.csv"))
#                     #2- feature stability scores based on CV
#                     stability_scores = feat_df.var() / feat_df.mean()
#                     #3- correlation analysis - no normality assumptioon
#                     correlation_matrix = feat_df.corr(method='spearman')
#                     correlation_matrix.to_csv(os.path.join(feature_qc_dir, "correlation_matrix.csv"))
#                     plt.figure(figsize=(15, 12))
#                     sns.heatmap(correlation_matrix, annot=True, cmap='coolwarm', fmt=".2f")
#                     plt.title("Correlation Matrix")
#                     plt.tight_layout()
#                     correlation_plot_path = os.path.join(feature_qc_dir, "correlation_matrix_plot.png")
#                     plt.savefig(correlation_plot_path)
#                     plt.close()
#                     print(f"corr matrix plot saved to: {correlation_plot_path}")
#                     if stability_scores is not None:
#                         plt.figure(figsize=(12, 6))
#                         plt.bar(stability_scores.index, stability_scores.values, color='skyblue')
#                         plt.title("Feature Stability Scores (Coefficient of Variation)")
#                         plt.xlabel("Feature")
#                         plt.ylabel("Stability Score")
#                         plt.xticks(rotation=45, ha='right')
#                         plt.tight_layout()
#                         cv_plot_path = os.path.join(feature_qc_dir, "stability_scores_plot.png")
#                         plt.savefig(cv_plot_path)
#                         plt.close()
#                         print(f"CV plot saved to: {cv_plot_path}")
#                     else:
#                         print("stability scores not available")
#
#                     #4-Random Forest Score computation
#                     target_variable_labels, _ = pd.factorize(feat_df.index.get_level_values(target_variable))
#                     rf_scores = {}
#                     for feature in feat_df.columns:
#                         model = RandomForestRegressor()
#                         feature_values = feat_df[feature].values.reshape(-1, 1)
#                         model.fit(feature_values, target_variable_labels)
#                         rf_scores[feature] = model.score(feature_values, target_variable_labels)
#
#                     rf_scores_df = pd.DataFrame(rf_scores, index=['RF_Score'])
#                     rf_scores_df.to_csv(os.path.join(feature_qc_dir, "random_forest_scores.csv"))
#                     plt.figure(figsize=(10, 6))
#                     plt.bar(rf_scores.keys(), rf_scores.values(), color='lightgreen')
#                     plt.title("Random Forest Scores for each Feature")
#                     plt.xlabel("Feature")
#                     plt.ylabel("Random Forest Score")
#                     plt.xticks(rotation=45, ha='right')
#                     plt.tight_layout()
#                     rf_plot_path = os.path.join(feature_qc_dir, "random_forest_scores_plot.png")
#                     plt.savefig(rf_plot_path)
#                     plt.close()
#                     print(f"Random Forest scores plot saved to: {rf_plot_path}")
#                     print("feature qc done!")
#                     ##
#                     ##prompting nmae change
#                     print("final_proc_pre_norm")
#                     rename_columns = input("Do you want to rename any columns? (yes/no): ")
#                     if rename_columns.lower() == 'yes':
#                         print("Current Columns:")
#                         for idx, column in enumerate(out_dat.columns):
#                             print(f"{idx}: {column}")
#                         renaming = input(
#                             "Enter the column name and new name separated by ':', each pair separated by commas (e.g., col1:new_col1,col2:new_col2): ")
#                         rename_dict = dict(pair.split(':') for pair in renaming.split(','))
#                         out_dat.rename(columns=rename_dict, inplace=True)
#                         print("Names changed successfully:)")
#                     out_dat.to_csv(os.path.join(results_folder, "data_with_new_names.csv"))
#         selected_data = out_dat.copy()
#         ##imputation segment
#         # prompting feature removal
#         print("any features to remove?")
#         for idx, column in enumerate(selected_data.columns):
#             print(f"{idx}: {column}")
#         features_to_remove = input("Enter comma-separated column names to remove (or 'none'): ").strip()
#         if features_to_remove.lower() != 'none':
#             features_to_remove = [col.strip() for col in features_to_remove.split(',')]
#             selected_data.drop(columns=features_to_remove, inplace=True)
#             print("Features removed:", features_to_remove)
#         # Initialize the SimpleImputer for missing values (NaN)
#         imputer = SimpleImputer(strategy='median')
#         # Impute missing values (NaN)
#         df_imputed = pd.DataFrame(imputer.fit_transform(selected_data), columns=selected_data.columns)
#         # Replace zero values with NaN to use the same imputer
#         df_imputed.replace(0, np.nan, inplace=True)
#         # Impute zero values (now NaNs) using the same strategy
#         df_final_imputed = pd.DataFrame(imputer.fit_transform(df_imputed), columns=selected_data.columns)
#         df_final_imputed.index = selected_data.index
#         df_final_imputed.to_csv(os.path.join(results_folder, "df_final_imputed.csv"))
#
#
#         print("ready_for_norm")
#         selected_data = df_final_imputed.copy()
#         # Normalize features using different techniques
#         # Create a directory for normalization results
#         normalization_dir = os.path.join(results_folder, "normalization")
#         os.makedirs(normalization_dir, exist_ok=True)
#
#         # Split selected_data based on the 'PC' index
#         pc_groups = selected_data.groupby(level='PC', group_keys=False)
#         df_reverted = pc_groups.apply(lambda x: x)
#         df_reverted.to_csv(os.path.join(results_folder, "processed_data_5.csv"))
#
#         # Iterate over each group and perform normalization
#         for pc_value, pc_group in pc_groups:
#             # Create a directory for the current PC value
#             pc_dir = os.path.join(normalization_dir, f"df_{pc_value}")
#             os.makedirs(pc_dir, exist_ok=True)
#             # Perform normalization for each group
#             pc_group.to_csv(os.path.join(results_folder, "original_data.csv"))
#
#             # For Min-Max normalization
#             min_max_normalized_data = (pc_group - pc_group.min()) / (pc_group.max() - pc_group.min())
#             min_max_normalized_data.to_csv(os.path.join(pc_dir, "min_max_normalized_data.csv"))
#
#             # For Central Logarithmic normalization
#             central_log_normalized_data = np.log(pc_group + np.sqrt(pc_group ** 2 + 1))
#             central_log_normalized_data.to_csv(os.path.join(pc_dir, "central_log_normalized_data.csv"))
#
#             # For Z-score normalization
#             z_score_normalized_data = (pc_group - pc_group.mean()) / pc_group.std()
#             z_score_normalized_data.to_csv(os.path.join(pc_dir, "z-score_normalized_data.csv"))
#
#             # # For Box-Cox transformation
#             # transformed_data = pd.DataFrame()
#             # for column in pc_group.columns:
#             #     transformed_column, _ = boxcox(pc_group[column])
#             #     transformed_data[column] = transformed_column
#             # transformed_data.index = pc_group.index
#             # transformed_data.to_csv(os.path.join(pc_dir, "box_cox_normalized_data.csv"))
#
#             # For log normalization
#             log_data = pd.DataFrame()
#             for column in pc_group.columns:
#                 log_column = np.log(pc_group[column])
#                 log_data[column] = log_column
#             log_data.index = pc_group.index
#             log_data.to_csv(os.path.join(pc_dir, "log_normalized_data.csv"))
#
#             # Plot QQ plots and histograms for each feature in each normalized data frame
#             for df_name, df in [("Original Data", pc_group),
#                                 ("Min-Max Normalized Data", min_max_normalized_data),
#                                 ("Central Log Normalized Data", central_log_normalized_data),
#                                 ("Z-score Normalized Data", z_score_normalized_data),
#           #                      ("Box_Cox Normalized Data", transformed_data),
#                                 ("Log Normalized Data", log_data)]:
#
#                 # Create a directory for the current data frame
#                 df_dir = os.path.join(pc_dir, df_name)
#                 os.makedirs(df_dir, exist_ok=True)
#
#                 for feature in df:
#                     ## qq imp
#                     plt.figure()
#                     scipy.stats.probplot(df[feature], dist="norm", plot=plt)
#                     plt.title(f"Q-Q-{feature}")
#                     qq_plot_path_2 = os.path.join(df_dir, f"qq_2_plot_{feature}.png")
#                     plt.savefig(qq_plot_path_2)
#                     plt.close()
#                     plt.clf()
#
#                 for feature in df.columns:
#                     # Convert Series to array for plotting QQ plot
#                     data_array = df[feature].to_numpy()
#                     # QQ plot
#                     sm.qqplot(data_array, line='45')
#                     plt.title(f"QQ Plot for {feature} ({df_name})")
#                     qq_plot_path = os.path.join(df_dir, f"qq_plot_{feature}.png")
#                     plt.savefig(qq_plot_path)
#                     plt.close()
#                     plt.clf()
#
#                     # Histogram
#                     fig = plt.figure(figsize=(8, 6))
#                     plt.hist(df[feature], bins=20, color='skyblue', edgecolor='black')
#                     plt.title(f"Histogram for {feature} ({df_name})")
#                     plt.xlabel("Value")
#                     plt.ylabel("Frequency")
#                     plt.tight_layout()
#                     histogram_path = os.path.join(df_dir, f"histogram_{feature}.png")
#                     plt.savefig(histogram_path)
#                     plt.close(fig)
#                     plt.clf()
#
#         print("Normalization and visualization completed. Results saved in 'normalization' directory.")
#         #"Box_Cox Normalized Data",
#         # Prompt user to select the normalized data frame to continue with
#         print("Select the normalized data frame you want to continue with:")
#         normalized_options = ["Min-Max Normalized Data", "Central Log Normalized Data", "Z-score Normalized Data",
#                                "Log Normalized Data"]
#         for idx, option in enumerate(normalized_options):
#             print(f"{idx + 1}: {option}")
#
#         selected_normalization = int(input("Enter the number corresponding to the desired normalized data frame: ")) - 1
#         selected_normalized_data = normalized_options[selected_normalization]
#
#         # Combine all PC values for the selected normalization into one final normalized data frame
#         final_normalized_data = pd.DataFrame()
#         for pc_value, pc_group in pc_groups:
#             pc_dir = os.path.join(normalization_dir, f"df_{pc_value}")
#             file_path = os.path.join(pc_dir, f"{selected_normalized_data.lower().replace(' ', '_')}.csv")
#             normalized_df = pd.read_csv(file_path, index_col=0)
#             final_normalized_data = pd.concat([final_normalized_data, normalized_df])
#
#         final_normalized_data.to_csv(os.path.join(results_folder, "final_normalized_data.csv"))
#         print(f"Final normalized data saved to: {os.path.join(results_folder, 'final_normalized_data.csv')}")
#
#         return final_normalized_data
#
# ##this should be the main folder /results
# main_folder_path = "MA-TMRE/MA-TMRE_result_combined"
# ##define redundant columns - any columns that are un-needed upfront
# unwanted_columns=['Number of Analyzed Fields', 'Time [s]', 'Plane', 'Timepoint', 'Cell Count', 'Temperature', 'Target Temperature',	'CO2',	'Target CO2',
#                   'Global Image Binning', 'Row', 'Column']
# ## which columns define the groups tested? select all relvant.
# index_columns = ['Compound', 'Cell Type', 'PC', 'Concentration', 'Cell_id', 'Severity', 'well', 'plate']
# #define important columns - all columns that should be tested for the analysis, in addition to technical columns.
# important_features = [
#                       "modified_spots_chanel_3 - Total Spot Area - Mean per Well", "all_cells - chanel_3 SER Spot 1 px - Mean per Well",
#                       "spots_chanel_3_final - Spot to Region Intensity - Mean per Well", "spots_chanel_3_final - Corrected Spot Intensity - Mean per Well",
#                       "spots_chanel_3_final - Region Intensity - Mean per Well", "ir_chanel_3_at - Intensity ir_chanel_3_at chanel_3 Mean - Mean per Well", "all_cells - chanel_3 Gabor Max 2 px w2 - Mean per Well",
#                       "all_cells - chanel_3 SER Valley 1 px - Mean per Well", "all_cells - chanel_3_intensity Mean - Mean per Well", "all_cells - chanel_2_intensity Mean - Mean per Well", "all_cells - chanel_1_intensity Mean - Mean per Well",
#                       "chan3_intensitymean_div_chan2area", "spotschan3_intensitymean_div_chan2area"]
# ## run constant
# process_plate_data_on_folders(main_folder_path, unwanted_columns, important_features, index_columns)
###################################################################################################


## MA TMRE
# logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
#
# def process_folders(main_folder):
#     try:
#         #check main folder
#         if not os.path.isdir(main_folder):
#             print(f"Error: The folder '{main_folder}' does not exist.")
#             return
#
#         #run through subfolders in main folder
#         for subfolder_name in os.listdir(main_folder):
#             subfolder_path = os.path.join(main_folder, subfolder_name)
#             if os.path.isdir(subfolder_path):
#                 print(f"\nProcessing subfolder: {subfolder_name}")
#
#                 #user to choose value for the 'PC' column- plate condition- must.
#                 pc_value = input("Enter the value for the 'PC' column(e.g- plate condition ), must add (numerical or string): ").strip()
#                 if not pc_value:
#                     print("Error: 'PC' column value cannot be empty. skipping this subfolder.")
#                     continue
#
#                 #run through text files in subfolder
#                 for txt_file in glob.glob(os.path.join(subfolder_path, '*.txt')):
#                     print(f"Reading file: {txt_file}")
#                     try:
#                         #Load file and process
#                         df = pd.read_csv(txt_file, skiprows=8, engine='python', sep='\t')
#                     except Exception as e:
#                         print(f"Error reading file '{txt_file}': {e}")
#                         continue
#
#                     #adding 'PC' column with the value
#                     df['PC'] = pc_value
#
#                     #extract plate name from file name
#                     plate_name = os.path.splitext(os.path.basename(txt_file))[0]
#                     print(f"Processing plate: {plate_name}")
#
#                     #output folder
#                     output_folder = os.path.join(main_folder, "results", f"{subfolder_name}_csv_files")
#                     os.makedirs(output_folder, exist_ok=True)
#
#                     #save CSV
#                     csv_filename = f"{plate_name}_csv.csv"
#
#                     #prompt user for flagged wells to remove
#                     flagged_wells_input = input(
#                         f"Enter flagged wells(bad wells) to remove for plate '{plate_name}' (e.g., row_column, ...) or 'none' if none: ").strip()
#
#                     #process flagged wells
#                     flagged_wells = []
#                     if flagged_wells_input.lower() != 'none':
#                         flagged_wells = [fw.strip() for fw in flagged_wells_input.split(',') if fw.strip()]
#
#                     # remove flagged wells
#                     for flagged_well in flagged_wells:
#                         try:
#                             row, column = flagged_well.split('_')
#                             row = int(row)
#                             column = int(column)
#                             df = df.loc[~((df['Row'] == row) & (df['Column'] == column))]
#                             print(f"Flagged well removed: Row {row}, Column {column}")
#                         except Exception as e:
#                             print(f"Error processing flagged well '{flagged_well}': {e}")
#                             continue
#
#             #editing column names
#                     try:
#                         df.columns = [col.split("[")[0].strip() if "[" in col else col for col in df.columns]
#                         logging.info("Column names cleaned.")
#                     except Exception as e:
#                         logging.error(f"Error cleaning column names: {e}")
#                         return
#                     #save modified DF
#                     csv_path = os.path.join(output_folder, csv_filename)
#                     try:
#                         df.to_csv(csv_path, index=False)
#                         print(f"File '{csv_filename}' saved in folder '{output_folder}'")
#                     except Exception as e:
#                         print(f"Error saving file '{csv_filename}': {e}")
#                         continue
#
#         #done processing all subfolders
#         print("\nAll subfolders processed successfully.")
#         print(f"review the files saved in {main_folder}/results.")
#
#     except Exception as e:
#         print(f"Unexpected error occurred: {e}")
#
# def prompt_for_review():
#     input("review the csv files. when ready, press enter to continue...")
#
#
# def process_data_combine(folder_path):
#     try:
#         #output result folder
#         output_folder = os.path.join(main_folder, "MA-TMRE_result_combined")
#         os.makedirs(output_folder, exist_ok=True)
#         result_combine_dir = os.path.join(main_folder, "MA-TMRE_result_combined/combined_csv_file")
#         if not os.path.exists(result_combine_dir):
#             os.makedirs(result_combine_dir, exist_ok=True)
#         #subdirectories
#         subdirectories = [subdir for subdir in os.listdir(folder_path) if os.path.isdir(os.path.join(folder_path, subdir))]
#         combined_dfs = []
#         #run on each subdirectory
#         for subdir in subdirectories:
#             sub_dir_path = os.path.join(folder_path, subdir)
#             print("Processing files in:", sub_dir_path)
#             dfs = []
#             for file in os.listdir(sub_dir_path):
#                 if file.endswith('.csv'):
#                     file_path = os.path.join(sub_dir_path, file)
#                     df = pd.read_csv(file_path, encoding='unicode_escape')
#                     df['plate'] = os.path.splitext(file)[0]  #add plate column- important
#                     dfs.append(df)
#         print("plate column added to data frame")
#         combined_df = pd.concat(dfs, ignore_index=True)
#         combined_dfs.append(combined_df)
#         combined_df = pd.concat(combined_dfs, ignore_index=True)
#         print("Data frames combined!")
#         combined_df['well'] = combined_df['Row'].astype('str') + combined_df['Column'].astype('str')
#         combined_df['well'] = combined_df['well'].astype('str')
#         print("well column added to data frame")
#         combined_df.to_csv(os.path.join(result_combine_dir, "combined_df_data.csv"), index=False)
#         print(f"review the files saved in {main_folder}/MA-TMRE_result_combined/combined_csv_file.")
#     except Exception as e:
#         logging.error(f"failed to combine plates...{e}")
#
#
# def process_plate_data_on_folders(main_folder, unwanted_columns, important_features, index_columns):
#     try:
#     ## open main folder for dash
#         main_results_dir = os.path.join(main_folder, "main_results_folder")
#         os.makedirs(main_results_dir, exist_ok=True)
#         unwanted_columns = list(set(unwanted_columns))
#         important_features = list(set(important_features))
#         index_columns = list(set(index_columns))
#         #running on each item in the main input folder
#         for item_name in os.listdir(main_folder):
#             item_path = os.path.join(main_folder, item_name)
#             if os.path.isdir(item_path):
#                 #running on each csv in each subfolder
#                 for file_name in os.listdir(item_path):
#                     if file_name.endswith(".csv"):
#                         file_path = os.path.join(item_path, file_name)
#                         logging.info(f"Processing file: {file_path}")
#                         #load each plate
#                         plate_df = pd.read_csv(file_path)
#                         #results folder with the name of the csv file
#                         results_folder = os.path.join(item_path, f"results_{file_name}_res")
#                         os.makedirs(results_folder, exist_ok=True)
#                         #drop unwanted columns & start proc
#                         try:
#                             plate_df = plate_df.drop(columns=unwanted_columns)
#                             logging.info(f"Dropped unwanted columns: {unwanted_columns}")
#                         except KeyError as e:
#                             logging.warning(f"Error dropping unwanted columns: {e}")
#                         #checking important features
#                         print("These are the important columns chosen for the analysis:")
#                         print(important_features)
#                         missing_features = [feature for feature in important_features if feature not in plate_df.columns]
#                         if missing_features:
#                             logging.warning(f"Missing important features: {missing_features}")
#                             important_features = [f for f in important_features if f in plate_df.columns]
#                         else:
#                             logging.info("all selected features found!")
#                         if not important_features:
#                             logging.error("No valid important features available. Skipping.")
#                             continue
#                         try:
#                             print("These are the index columns chosen for the analysis:")
#                             print(index_columns)
#                             index_columns.append('well')
#                             index_columns.append('plate')
#                             logging.info("'well' & 'plate' columns added to index_columns.")
#                         except Exception as e:
#                             logging.error(f"Error adding 'well' and 'plate' columns to index columns: {e}")
#                         #generate combinations of index_columns for new indexes
#                         try:
#                             new_index_columns = []
#                             for r in range(1, len(index_columns) + 1):
#                                 for comb in itertools.combinations(index_columns, r):
#                                     new_index_columns.append(comb)
#                             logging.info(f"Generated {len(new_index_columns)} new index column combinations.")
#                         except Exception as e:
#                             logging.error(f"Error generating new index column combinations: {e}")
#                         #create new index columns based on combinations
#                         try:
#                             new_cols = {}
#                             for cols in new_index_columns:
#                                 col_name = '_'.join(cols)
#                                 if col_name not in plate_df.columns:
#                                     new_cols[col_name] = plate_df[list(cols)].apply(lambda x: '_'.join(map(str, x)), axis=1)
#                                 else:
#                                     print(f"{col_name} already exists and will not be added")
#                             plate_df = pd.concat([plate_df, pd.DataFrame(new_cols)], axis=1)
#                             logging.info("Created new index columns based on combinations.")
#                         except Exception as e:
#                             logging.error(f"Error creating new index columns: {e}")
#                         ## keep only new_index and imp features
#                         ## keep only new_index and imp features
#                         try:
#                             logging.info("subsetting data to focus on important features & index columns...")
#                             columns1 = plate_df[index_columns]
#                             columns2 = plate_df[important_features]
#                             plate_df_int = pd.concat([columns1, columns2, pd.DataFrame(new_cols)], axis=1)
#                         except Exception as e:
#                             logging.error(f"Failed to subset{e}")
#                             plate_df_int = plate_df.copy()
#                         try:
#                             # remove rows and columns with ≥90% zero or NA values
#                             filtered_df = plate_df_int.drop(columns=index_columns, errors='ignore')
#                             filtered_df = filtered_df.dropna(thresh=0.1 * len(filtered_df), axis=0)
#                             filtered_df = filtered_df.replace(0, np.nan)
#                             filtered_df = filtered_df.dropna(thresh=0.1 * len(filtered_df.columns), axis=1)
#                             final_df = pd.concat([columns1, filtered_df], axis=1)
#                             logging.info("Dropped rows and columns with ≥90% zero or NA values.")
#                             plate_df_2 = final_df.copy()
#                         except Exception as e:
#                             logging.error(f"Error dropping rows/columns with high NA/zero values: {e}")
#                         ## add cell count feat
#                         try:
#                             print("usual cell_count_feat names: all_cells - Number of Objects")
#                             cell_count_feat = input("Enter cell_count_feat(check in csv if not sure): ")
#                             plate_df_2[cell_count_feat] = plate_df[cell_count_feat]
#                             plate_df_2[cell_count_feat] = pd.to_numeric(plate_df_2[cell_count_feat])
#                         except Exception as e:
#                             logging.error(f"failed to locate cell count feat {e}")
#                         ##user to remove rows based on any avaialble columns
#                         try:
#                             print("Section to remove rows based on chosen column, choose column following by valuesto remove, otherwise, type no")
#                             while True:
#                                 print("Columns in the dataset:")
#                                 for idx, column in enumerate(plate_df_2.columns):
#                                     print(f"{idx}: {column}")
#                                 column_to_filter = input("Enter the column name to filter by (or 'no' to skip): ").strip()
#                                 if column_to_filter.lower() == 'no':
#                                     break
#                                 if column_to_filter not in plate_df_2.columns:
#                                     print("Invalid column name. Please try again.")
#                                     continue
#                                 print(f"Values in the selected column '{column_to_filter}':")
#                                 print(plate_df_2[column_to_filter].unique())
#                                 values_to_remove = input("Enter the values to remove (comma-separated): ").split(',')
#                                 values_to_remove = [val.strip() for val in values_to_remove]
#                                 plate_df_2 = plate_df_2[~plate_df_2[column_to_filter].isin(values_to_remove)]
#                                 print(f"Rows with values {values_to_remove} in column '{column_to_filter}' have been removed.")
#                         except Exception as e:
#                             logging.error(f"Error removing custom column values: {e}")
#                             return
#                         #remove additional columns if needed
#                         try:
#                             print("Section to remove columns from the data, choose column to remove, otherwise, type none")
#                             for idx, column in enumerate(plate_df_2.columns):
#                                 print(f"{idx + 1}: {column}")
#                             features_to_remove = input("Enter comma-separated column names to remove (or 'none'): ").strip()
#                             if features_to_remove.lower() != 'none':
#                                 features_to_remove = [col.strip() for col in features_to_remove.split(',')]
#                                 plate_df_2.drop(columns=features_to_remove, inplace=True)
#                                 logging.info(f"Features removed: {features_to_remove}")
#                         except KeyError as e:
#                             logging.warning(f"Some features to remove were not found: {e}")
#                         except Exception as e:
#                             logging.error(f"Error removing features: {e}")
#                         #IDs to remove ?
#                         try:
#                             print("Section to remove cell id values from the data, choose id to remove, otherwise, type none")
#                             print(plate_df_2["cell_id"].unique())
#                             remove_id = input("Enter ID value to remove (or 'none'): ").strip()
#                             if remove_id.lower() != 'none':
#                                 plate_df_2 = plate_df_2[~plate_df_2['cell_id'].astype(str).str.contains(remove_id)]
#                                 logging.info(f"Rows with ID '{remove_id}' removed.")
#                         except KeyError:
#                             logging.warning("Column 'ID' not found. Skipping ID removal.")
#                         except Exception as e:
#                             logging.error(f"Error removing rows based on ID: {e}")
#                         #remove rows and columns with ≥90% zero or NA values
#                         # try:
#                         #     if column not in index_columns:
#                         #         plate_df_2 = plate_df_2.dropna(thresh=0.1 * len(plate_df_2), axis=0)
#                         #         plate_df_2 = plate_df_2.replace(0, np.nan)
#                         #         plate_df_2 = plate_df_2.dropna(thresh=0.1 * len(plate_df_2.columns), axis=1)
#                         #         logging.info("Dropped rows and columns with ≥90% zero or NA values.")
#                         # except Exception as e:
#                         #     logging.error(f"Error dropping rows/columns with high NA/zero values: {e}")
#                         #save first df
#                         plate_df_2.to_csv(os.path.join(results_folder, "df_proc_1.csv"), index=False)
#                         print("first processed data frame saved to result directory")
#                         #directory for main results each subfolder
#                         general_results_dir = os.path.join(results_folder, "general_results")
#                         os.makedirs(general_results_dir, exist_ok=True)
#                         #index column for statistical summary
#                         try:
#                             print(new_index_columns)
#                             index_col = input("Enter the index column for statistical summary: seperated by '_' ").strip()
#                             summary_df = plate_df_2.groupby(index_col).describe()
#                             logging.info(f"Generated statistical summary grouped by '{index_col}'.")
#                         except KeyError:
#                             logging.error(f"Index column '{index_col}' not found in DataFrame.")
#                             summary_df = pd.DataFrame()
#                         except Exception as e:
#                             logging.error(f"Error generating statistical summary: {e}")
#                             summary_df = pd.DataFrame()
#                         #save statistical summary
#                         try:
#                             summary_excel_path = os.path.join(general_results_dir, "init_statistical_summary.xlsx")
#                             summary_df.to_excel(summary_excel_path)
#                             logging.info(f"Statistical summary saved to '{summary_excel_path}'")
#                         except Exception as e:
#                             logging.error(f"Error saving statistical summary to Excel: {e}")
#                         #directory for cell count results
#                         print("cell count process initiated")
#                         cell_count_results_dir = os.path.join(results_folder, "cell_count_results")
#                         os.makedirs(cell_count_results_dir, exist_ok=True)
#                         #histogram for cell count
#                         try:
#                             if cell_count_feat in plate_df_2.columns:
#                                 fig = px.histogram(plate_df_2, x=cell_count_feat, title="Histogram of Cell Count")
#                                 histogram_file_path = os.path.join(cell_count_results_dir, "histogram.html")
#                                 fig.write_html(histogram_file_path)
#                                 logging.info(f"Cell count histogram saved to {histogram_file_path}")
#                             else:
#                                 logging.warning(f"{cell_count_feat} column not found. histogram not generated.")
#                         except Exception as e:
#                             logging.error(f"Error generating histogram: {e}")
#                         #statistics for cell count
#                         try:
#                             if cell_count_feat in plate_df_2.columns:
#                                 stats = plate_df_2[cell_count_feat].describe()
#                                 stats_file_path = os.path.join(cell_count_results_dir, "cell_count-statistics.txt")
#                                 with open(stats_file_path, "w") as f:
#                                     f.write(stats.to_string())
#                                 logging.info(f"Cell count statistics saved to {stats_file_path}")
#                             else:
#                                 logging.warning(f"{cell_count_feat} column not found. statistics not generated.")
#                         except Exception as e:
#                             logging.error(f"Error creating cell count statistics: {e}")
#                         #input for minimum and maximum cell count values
#                         try:
#                             min_input = input("Enter the minimum value for cell count (leave blank to skip): ").strip()
#                             max_input = input("Enter the maximum value for cell count (leave blank to skip): ").strip()
#                             min_value = float(min_input) if min_input else None
#                             max_value = float(max_input) if max_input else None
#                             if min_value is not None and max_value is not None:
#                                 plate_df_2 = plate_df_2[(plate_df_2[cell_count_feat] >= min_value) & (
#                                         plate_df_2[cell_count_feat] <= max_value)]
#                                 logging.info(f"Filtered data with cell count between {min_value} and {max_value}")
#                             elif min_value is not None:
#                                 plate_df_2 = plate_df_2[plate_df_2[cell_count_feat] >= min_value]
#                                 logging.info(f"Filtered data with cell count >= {min_value}")
#                             elif max_value is not None:
#                                 plate_df_2 = plate_df_2[plate_df_2[cell_count_feat] <= max_value]
#                                 logging.info(f"Filtered data with cell count <= {max_value}")
#                             else:
#                                 logging.info("No filtering applied to cell count.")
#                         except ValueError:
#                             logging.error("Invalid input for minimum or maximum cell count. Skipping filtering.")
#                         except Exception as e:
#                             logging.error(f"Error filtering cell count: {e}")
#                         plate_df_2.to_csv(os.path.join(results_folder, "df_proc_2.csv"))
#                         print("second processed data frame saved to result directory")
#                         print(new_index_columns)
#                         #directory for box plots results
#                         box_plot_dir = os.path.join(results_folder, "box_plots")
#                         os.makedirs(box_plot_dir, exist_ok=True)
#                         #prompting index column for box plots
#                         index_col = input("Enter the index column for box plot vis: ")
#                         bar_df = plate_df_2.copy()
#                         # box plots for selected index- important features only
#                         try:
#
#                             print("box plotting selected features!")
#                             for col in bar_df:
#                                 if col in important_features:
#                                     bp = px.box(bar_df, x=index_col, y=bar_df[col], points="all", color=index_col, notched=True)
#                                     bp.update_traces(quartilemethod="inclusive")  # or "inclusive", or "linear" by default
#                                     bp.update_layout(
#                                         font_family="Arial",
#                                         font_color="Black",
#                                         font_size=20,
#                                         font=dict(
#                                             family="Arial",
#                                             size=20,
#                                             color="Black"
#                                         )
#                                     )
#                                     print("one_down")
#                                     bp.write_html(os.path.join(box_plot_dir, f'{col}_bar_plot.html'))
#                                     #bp.write_image(os.path.join(box_plot_dir, f"box_plot{col}.pdf"), engine="kaleido")
#                             print("box plots done!")
#                         except Exception as e:
#                             logging.error(f"failed box plotting{e}")
#                         #directory for summary files
#                         summary_files_dir = os.path.join(results_folder, "summary_files")
#                         os.makedirs(summary_files_dir, exist_ok=True)
#                         #creating excel  with sheets with split data and summary stats
#                         try:
#                             print("creating excel files and statistical summary & ppt for combination of groups!")
#                             group_columns = ['cell_type', 'compound', 'concentration']
#                             excel_file_path = os.path.join(summary_files_dir, f'summary_file.xlsx')
#                             with pd.ExcelWriter(excel_file_path, engine='openpyxl') as writer:
#                                 for name, group in plate_df_2.groupby(group_columns):
#                                     try:
#                                         sheet_name = f"{name[0]}_{name[1]}_{name[2]}"
#                                         group.to_excel(writer, sheet_name=sheet_name, index=False)
#                                         ####summary statistics to another sheet
#                                         summary_df = group.describe(include='all').T[
#                                             ['mean', 'std', 'min', '25%', '50%', '75%', 'max']]
#                                         summary_df.to_excel(writer, sheet_name=f"{sheet_name}_Summary")
#                                     except ValueError:
#                                         pass
#                                 if not writer.sheets:
#                                     pd.DataFrame({'Dummy': [1]}).to_excel(writer, sheet_name='DummySheet', index=False)
#                             #ppt with slides for each unique combination
#                             print("ppt file generated!")
#                             ppt_file_path = os.path.join(summary_files_dir, f'slides.pptx')
#                             prs = Presentation()
#                             for name, _ in plate_df_2.groupby(group_columns):
#                                 slide = prs.slides.add_slide(prs.slide_layouts[5])
#                                 title = slide.shapes.title
#                                 title.text = f"Combination: {name[0]}, {name[1]}, {name[2]}"
#                             prs.save(ppt_file_path)
#                         except Exception as e:
#                             logging.error(f"failed creating summary{e}")
#                         #bar plot for Number of Cells
#                         try:
#                             print("working on final cell count vis!")
#                             group_columns = ['cell_type', 'compound', 'concentration', 'cell_id']
#                             plt.figure(figsize=(20, 18))
#                             for name, group in plate_df_2.groupby(group_columns):
#                                 try:
#                                     plt.bar(str(name), group[cell_count_feat].mean(),
#                                             yerr=group[cell_count_feat].std(), capsize=5,
#                                             label=str(name))
#                                 except ValueError:
#                                     pass
#                             plt.title(f'Bar Plot for Number of Cells - {file_name}')
#                             plt.xlabel('cell_type, compound, concentration, cell_id')
#                             plt.ylabel('Number of Cells')
#                             plt.xticks(rotation=45, ha='right')
#                             plt.legend()
#                             plt.savefig(os.path.join(general_results_dir, f'number_of_cells_bar_plot.pdf'))
#                             plt.close()
#                         except Exception as e:
#                             logging.error(f"failed bar plotting cell number{e}")
#                         ## remove cell feat
#                         try:
#                             plate_df_2 = plate_df_2.drop(columns=[cell_count_feat])
#                         except Exception as e:
#                             logging.error(f"failed to remove {cell_count_feat}...{e}")
#                         #heatmap for selected features (mean values & normalized)
#                         try:
#                             heat_df = plate_df_2.copy()
#                             index_col = input("Enter the index column for heatmap: ")
#                             print("heatmap prep!")
#                             heat_df = heat_df.set_index(index_col, drop=True)
#                             numeric_features = heat_df[important_features].select_dtypes(include=np.number)
#                             numeric_features_normalized = (numeric_features - numeric_features.mean()) / numeric_features.std()
#                             numeric_features_normalized = numeric_features_normalized.groupby(index_col).mean()
#                             lan = plt.rcParams['font.family'] = ['Arial']
#                             plt.figure(figsize=(20, 12))
#                             font = {'family': 'Arial',
#                                     'size': 22}
#                             plt.rc('font', **font)
#                             plt.yticks(family='Arial', rotation="horizontal", size=10)
#                             plt.xticks(family='Arial', rotation="vertical", size=8)
#                             plt.margins(0.1)
#                             plt.subplots_adjust(bottom=0.3)
#                             sns.heatmap(numeric_features_normalized, annot=True, cmap='coolwarm')
#                             plt.title(f'Heatmap for Selected Numeric Features - {file_name}')
#                             plt.tight_layout()
#                             plt.savefig(os.path.join(general_results_dir, f'heatmap.pdf'))
#                             plt.close()
#                             print("heatmap done!")
#                         except Exception as e:
#                             logging.error(f"Heatmap failed: {e}")
#                         #PCA based on all numeric features without ID
#                         try:
#                             print("creating pca for group with no ID and with id")
#                             pca = PCA(n_components=2)
#                             # Impute missing values in numeric features
#                             numeric_features = plate_df_2.select_dtypes(include=np.number)
#                             imputer = SimpleImputer(strategy='mean')
#                             numeric_features_imputed = imputer.fit_transform(numeric_features)
#                             #interactive pcas
#                             pca_result_without_id = pca.fit_transform(numeric_features_imputed)
#                             label_without_id = plate_df_2['cell_type'].astype(str) + '_' + \
#                                             plate_df_2['compound'].astype(str) + '_' + \
#                                             plate_df_2['concentration'].astype(str)
#                             color_labels_wo_id = pd.factorize(label_without_id)[0]
#                             pca_df1 = pd.DataFrame({
#                                 'PC1': pca_result_without_id[:, 0],
#                                 'PC2': pca_result_without_id[:, 1],
#                                 'Label': label_without_id,
#                                 'Color': color_labels_wo_id
#                             })
#                             fig = px.scatter(pca_df1, x='PC1', y='PC2', color='Label',
#                                              title=f'PCA Plot for All Numeric Features without CELL ID - {file_name}',
#                                              color_continuous_scale='viridis')
#                             fig.update_layout(legend_title='Cell Type, Compound, Concentration')
#                             fig.add_annotation(text=f'Explained Variance Ratio: {pca.explained_variance_ratio_}', xref='paper',
#                                                yref='paper',
#                                                x=0.98, y=0.02, showarrow=False, font=dict(size=10))
#                             fig.write_html(os.path.join(general_results_dir, f"pca_plot_wo_id_{file_name[:-4]}.html"))
#                             #fig.write_image(os.path.join(general_results_dir, f'pca_plot_without_id2.pdf'), engine="kaleido")
#                             # PCA based on all numeric features with CELL ID-2 interactive
#                             pca_result_with_id = pca.fit_transform(numeric_features_imputed)
#                             # A label for each unique combination of Cell Type, Compound, Concentration, and CELL ID
#                             label_with_id = plate_df_2['cell_type'].astype(str) + '_' + \
#                                                plate_df_2['compound'].astype(str) + '_' + \
#                                                plate_df_2['concentration'].astype(str) + '_' + \
#                                                plate_df_2['cell_id'].astype(str)
#                             color_labels_with_id = pd.factorize(label_with_id)[0]
#                             # Create DataFrame for Plotly
#                             pca_df2 = pd.DataFrame({
#                                 'PC1': pca_result_with_id[:, 0],
#                                 'PC2': pca_result_with_id[:, 1],
#                                 'Label': label_with_id,
#                                 'Color': color_labels_with_id
#                             })
#                             fig = px.scatter(pca_df2, x='PC1', y='PC2', color='Label',
#                                              title=f'PCA Plot for All Numeric Features with CELL ID - {file_name}',
#                                              color_continuous_scale='viridis')
#                             fig.update_layout(legend_title='Cell Type, Compound, Concentration, ID')
#                             fig.add_annotation(text=f'Explained Variance Ratio: {pca.explained_variance_ratio_}', xref='paper',
#                                                yref='paper',
#                                                x=0.98, y=0.02, showarrow=False, font=dict(size=10))
#                             fig.write_html(os.path.join(general_results_dir, f"pca_plot_with_id_{file_name[:-4]}.html"))
#                         except Exception as e:
#                             logging.error(f"PCA failed: {e}")
#                         print("pca done!")
#                         #directory for outlier results
#                         try:
#                             outlier_files_dir = os.path.join(results_folder, "outlier_files")
#                             os.makedirs(outlier_files_dir, exist_ok=True)
#                             for col in plate_df_2.columns:
#                                 if col in index_columns:
#                                     plate_df_2[col] = plate_df_2[col].astype(str)
#                             perform_outlier_detection = input(
#                                 "Do you want to perform outlier detection? (yes/no): ").strip().lower()
#                             if perform_outlier_detection == 'yes':
#                                 print("Applying basic Outlier methods...")
#                                 # prompting index column for outliers detecetion
#                                 print("Columns available for outlier detection:")
#                                 for idx, column in enumerate(plate_df_2.columns):
#                                     print(f"{idx}: {column}")
#                                 index_column = input("Choose one column as index for outlier detection: ")
#                                 out_dat_init = plate_df_2.copy()
#                                 # final filtered DataFrame
#                                 init_file_path = os.path.join(outlier_files_dir,
#                                                                'out_dat_init.csv')
#                                 out_dat_init.to_csv(init_file_path, index=False)
#                                 #outlier detection functions options
#                                 ## perc
#                                 def compute_percentiles(df, output_file):
#                                     up_bounds = []
#                                     low_bounds = []
#                                     above_count = []
#                                     below_count = []
#                                     numeric_cols = df.select_dtypes(include='number').columns
#                                     for col in numeric_cols:
#                                         percentile_low = df[col].quantile(0.01)
#                                         percentile_high = df[col].quantile(0.99)
#                                         up_bound = percentile_high
#                                         low_bound = percentile_low
#                                         above = df[df[col] > percentile_high]
#                                         below = df[df[col] < percentile_low]
#                                         above_count.append(len(above))
#                                         below_count.append(len(below))
#                                         df = df[(df[col] >= percentile_low) & (df[col] <= percentile_high)]
#                                         up_bounds.append(up_bound)
#                                         low_bounds.append(low_bound)
#                                     output = pd.DataFrame(
#                                         {'Column': numeric_cols, 'Up Bound': up_bounds, 'Low Bound': low_bounds,
#                                          'Above Count': above_count,
#                                          'Below Count': below_count})
#                                     output.to_csv(output_file, index=False)
#                                     return df
#                                 ## iqr
#                                 def compute_iqr(df, output_file):
#                                     up_bounds = []
#                                     low_bounds = []
#                                     above_count = []
#                                     below_count = []
#                                     numeric_cols = df.select_dtypes(include='number').columns
#                                     for col in numeric_cols:
#                                         percentile25th = df[col].quantile(0.25)
#                                         percentile75th = df[col].quantile(0.75)
#                                         iqr = percentile75th - percentile25th
#                                         up_bound = percentile75th + 1.5 * iqr
#                                         low_bound = percentile25th - 1.5 * iqr
#                                         above = df[df[col] > up_bound]
#                                         below = df[df[col] < low_bound]
#                                         above_count.append(len(above))
#                                         below_count.append(len(below))
#                                         df = df[(df[col] >= low_bound) & (df[col] <= up_bound)]
#                                         up_bounds.append(up_bound)
#                                         low_bounds.append(low_bound)
#                                     output = pd.DataFrame(
#                                         {'Column': numeric_cols, 'Up Bound': up_bounds, 'Low Bound': low_bounds,
#                                          'Above Count': above_count,
#                                          'Below Count': below_count})
#                                     output.to_csv(output_file, index=False)
#                                     return df
#                                 ## zscore
#                                 def compute_scores(df, output_file):
#                                     up_scores = []
#                                     low_scores = []
#                                     above_count = []
#                                     below_count = []
#                                     numeric_cols = df.select_dtypes(include='number').columns
#                                     for col in numeric_cols:
#                                         up_score = df[col].mean() + 3 * df[col].std()
#                                         low_score = df[col].mean() - 3 * df[col].std()
#                                         above = df[df[col] > up_score]
#                                         below = df[df[col] < low_score]
#                                         above_count.append(len(above))
#                                         below_count.append(len(below))
#                                         df = df[(df[col] >= low_score) & (df[col] <= up_score)]
#                                         up_scores.append(up_score)
#                                         low_scores.append(low_score)
#                                     output = pd.DataFrame(
#                                         {'Column': numeric_cols, 'Up Score': up_scores, 'Low Score': low_scores,
#                                          'Above Count': above_count,
#                                          'Below Count': below_count})
#                                     output.to_csv(output_file, index=False)
#                                     return df
#                                 ## LOF
#                                 def LOF_outlier(df, output_file1, output_file2):
#                                     # outlier detection using LOF
#                                     types = df[index_column].unique()
#                                     print(types)
#                                     mask = []
#                                     features = df.select_dtypes(include='number').columns
#                                     detector_list = [
#                                         ("Local Outlier Factor 30", LocalOutlierFactor(n_neighbors=8))
#                                     ]
#                                     for name, algorithm in detector_list:
#                                         errors = np.full(len(df), fill_value=np.nan)
#                                         outliers = np.full(len(df), fill_value=np.nan)
#                                         for type in types:
#                                             x = df.loc[:, features].values
#                                             F = x.sum(1)
#                                             mask = np.zeros(x.shape[0])
#                                             mask[np.isfinite(F)] = 1
#                                             mask_type = mask * np.array(df[index_column] == type)
#                                             Curr_df = df.loc[mask_type == 1, features]
#                                             x = Curr_df.values
#                                             if name == "Local Outlier Factor 30":
#                                                 algorithm.fit(x)
#                                                 errors[mask_type == 1] = algorithm.negative_outlier_factor_
#                                                 outliers[mask_type == 1] = algorithm.fit_predict(x)
#                                         df[name] = errors
#                                         df[f'{name}_outliers'] = outliers
#                                         df.set_index(name, inplace=True,
#                                                      append=True, drop=False)
#                                         df.to_csv(os.path.join(output_file1))
#                                         # exclude rows that were defined as outliers
#                                         for col in df.columns:
#                                             if col.endswith("_outliers"):
#                                                 df = df[df[col] != -1]
#                                     df.to_csv(os.path.join(output_file2))
#                                     print("outliers removed!")
#                                     return df
#                                 #run the outlier detection functions
#                                 df_percentiles = compute_percentiles(out_dat_init,
#                                                                      os.path.join(outlier_files_dir,
#                                                                                   'percentiles_outliers.csv'))
#                                 perc_file_path = os.path.join(outlier_files_dir,
#                                                                'perc_dat.csv')
#                                 df_percentiles.to_csv(perc_file_path, index=False)
#                                 df_iqr = compute_iqr(out_dat_init, os.path.join(outlier_files_dir, 'iqr_outliers.csv'))
#                                 iqr_file_path = os.path.join(outlier_files_dir,
#                                                                'iqr_dat.csv')
#                                 df_iqr.to_csv(iqr_file_path, index=False)
#                                 df_scores = compute_scores(out_dat_init,
#                                                            os.path.join(outlier_files_dir, 'scores_outliers.csv'))
#                                 zscore_file_path = os.path.join(outlier_files_dir,
#                                                                'zscore_dat.csv')
#                                 df_scores.to_csv(zscore_file_path, index=False)
#                                 df_lof = LOF_outlier(out_dat_init,
#                                                      os.path.join(outlier_files_dir, 'processed_data_outliers_1.csv'),
#                                                      os.path.join(outlier_files_dir, 'processed_data_outliers_2.csv'))
#                                 print("Please review outlier detection methods in outlier results folder and choose one method to continue with:")
#                                 method_choice = input("Enter 1 for percentiles, 2 for IQR, 3 for scores, 4 for lof: ")
#                                 if method_choice == '1':
#                                     out_dat = df_percentiles
#                                 elif method_choice == '2':
#                                     out_dat = df_iqr
#                                 elif method_choice == '3':
#                                     out_dat = df_scores
#                                 elif method_choice == '4':
#                                     out_dat = df_lof
#                                     out_dat = out_dat.drop(columns=["Local Outlier Factor 30", "Local Outlier Factor 30_outliers", "Local Outlier Factor 30"])
#                                 else:
#                                     print("Invalid choice. Defaulting to original DataFrame.")
#                                     out_dat = out_dat_init
#                                     out_dat = out_dat.drop(
#                                         columns=["Local Outlier Factor 30", "Local Outlier Factor 30_outliers",
#                                                  "Local Outlier Factor 30"])
#                             else:
#                                 print("No outlier detection ran.")
#                                 out_dat = plate_df_2.copy()
#                             #final filtered DataFrame
#                             final_file_path = os.path.join(results_folder, 'final_filtered_data(post outlier detection).csv')
#                             out_dat.to_csv(final_file_path, index=False)
#                             print("Final filtered DataFrame saved to", final_file_path)
#                         except Exception as e:
#                             logging.error(f"Failed to perform outlier detection{e}")
#                         #prompting index column for final proc- "PC" is must!
#                         for idx, column in enumerate(out_dat.columns):
#                             print(f"{idx}: {column}")
#                         #columns to set as indexes
#                         try:
#                             columns_to_set_index = input("Enter ALL needed columns to set as indexes -PC IS MUST! (comma-separated): ").split(',')
#                             columns_to_set_index = [col.strip() for col in columns_to_set_index if col.strip()]
#                             logging.info(f"Columns to set as index: {columns_to_set_index}")
#                             if columns_to_set_index:
#                                 missing_index_cols = [col for col in columns_to_set_index if col not in out_dat.columns]
#                                 if missing_index_cols:
#                                     logging.warning(f"Index columns not found: {missing_index_cols}")
#                                 out_dat = out_dat.set_index([col for col in columns_to_set_index if col in out_dat.columns], drop=False)
#                                 logging.info(f"Set columns as index: {columns_to_set_index}")
#                         except Exception as e:
#                             logging.error(f"Error setting index columns: {e}")
#                         out_dat = out_dat.select_dtypes(include=np.number)
#                         #target variable selection for feature qc
#                         try:
#                             #user to select the correct index title as the target variable
#                             if out_dat.index.names:
#                                 logging.info("Select the correct index title to be the target variable for feature quality control:")
#                                 for idx, index_title in enumerate(out_dat.index.names):
#                                     logging.info(f"{idx + 1}: {index_title}")
#                                 while True:
#                                     try:
#                                         selected_index = int(
#                                             input("Enter the number corresponding to the correct index title: ")) - 1
#                                         if selected_index < 0 or selected_index >= len(out_dat.index.names):
#                                             raise ValueError("Selection out of range.")
#                                         break
#                                     except ValueError as e:
#                                         logging.error(f"Invalid selection. Please try again. Error: {e}")
#                                 #selected index title as the target variable
#                                 target_variable = out_dat.index.names[selected_index]
#                                 logging.info(f"Target variable selected: {target_variable}")
#                             else:
#                                 raise ValueError("No index names found in the data.")
#                             #factor the target variable
#                             target_variable_labels, _ = pd.factorize(out_dat.index.get_level_values(target_variable))
#                             logging.info(f"Target variable labels: {target_variable_labels}")
#                             #directory for feature qc results
#                             feature_qc_dir = os.path.join(results_folder, "feature_qc")
#                             os.makedirs(feature_qc_dir, exist_ok=True)
#                             #statistical summary for features
#                             statistical_summary = out_dat.describe()
#                             statistical_summary.to_csv(os.path.join(feature_qc_dir, "statistical_summary.csv"))
#                             logging.info("Statistical summary saved to 'statistical_summary.csv'.")
#                             #feature stability scores -CV - for normalizied features
#                             normalized_data = (out_dat - out_dat.min()) / (out_dat.max() - out_dat.min())
#                             if normalized_data.isna().any().any():
#                                 logging.warning("Normalization produced NaN values. Check your data for zero variance features.")
#                             stability_scores = normalized_data.var() / normalized_data.mean()
#                             logging.info("Stability scores calculated.")
#                             #correlation
#                             correlation_matrix = out_dat.corr(method='spearman')
#                             correlation_matrix.to_csv(os.path.join(feature_qc_dir, "correlation_matrix.csv"))
#                             logging.info("Correlation matrix saved to 'correlation_matrix.csv'.")
#                             #plot the correlation matrix
#                             lan = plt.rcParams['font.family'] = ['Arial']
#                             plt.figure(figsize=(32, 18))
#                             font = {'family': 'Arial',
#                                     'size': 10}
#                             plt.rc('font', **font)
#                             plt.yticks(family='Arial', rotation="horizontal", size=10)
#                             plt.xticks(family='Arial', rotation="vertical", size=8)
#                             plt.margins(0.1)
#                             plt.subplots_adjust(bottom=0.3)
#                             sns.heatmap(correlation_matrix, annot=True, cmap='coolwarm', fmt=".2f")
#                             plt.title("Correlation Matrix")
#                             plt.tight_layout()
#                             correlation_plot_path = os.path.join(feature_qc_dir, "correlation_matrix_plot.png")
#                             plt.savefig(correlation_plot_path)
#                             plt.close()
#                             logging.info(f"Correlation matrix plot saved to: {correlation_plot_path}")
#                             #CV plot
#                             if stability_scores is not None and not stability_scores.isna().all():
#                                 plt.figure(figsize=(32, 16))
#                                 plt.bar(stability_scores.index, stability_scores.values, color='skyblue')
#                                 plt.title("Feature Stability Scores (Coefficient of Variation)")
#                                 plt.xlabel("Feature")
#                                 plt.ylabel("Stability Score")
#                                 plt.xticks(rotation=45, ha='right')
#                                 plt.margins(0.1)
#                                 plt.tight_layout()
#                                 cv_plot_path = os.path.join(feature_qc_dir, "stability_scores_plot.png")
#                                 plt.savefig(cv_plot_path)
#                                 plt.close()
#                                 logging.info(f"CV plot saved to: {cv_plot_path}")
#                             else:
#                                 logging.warning("Stability scores not available or invalid. Skipping CV plot.")
#                             #random forest score computation per features
#                             rf_scores = {}
#                             try:
#                                 for feature in out_dat.columns:
#                                     model = RandomForestRegressor()
#                                     feature_values = out_dat[feature].values.reshape(-1, 1)
#                                     model.fit(feature_values, target_variable_labels)
#                                     rf_scores[feature] = model.score(feature_values, target_variable_labels)
#                                     logging.debug(f"Random Forest score for feature '{feature}': {rf_scores[feature]}")
#                             except Exception as e:
#                                 logging.error(f"Error occurred during Random Forest model training: {e}")
#                                 rf_scores[feature] = np.nan
#                             #RF scores to csv
#                             rf_scores_df = pd.DataFrame(rf_scores, index=['RF_Score'])
#                             rf_scores_df.to_csv(os.path.join(feature_qc_dir, "random_forest_scores.csv"))
#                             logging.info("Random Forest scores saved to 'random_forest_scores.csv'.")
#                             #plot RF scores
#                             plt.figure(figsize=(32, 16))
#                             plt.bar(rf_scores.keys(), rf_scores.values(), color='lightgreen')
#                             plt.title("Random Forest Scores for each Feature")
#                             plt.xlabel("Feature")
#                             plt.ylabel("Random Forest Score")
#                             plt.xticks(rotation=45, ha='right')
#                             plt.margins(0.1)
#                             plt.tight_layout()
#                             rf_plot_path = os.path.join(feature_qc_dir, "random_forest_scores_plot.png")
#                             plt.savefig(rf_plot_path)
#                             plt.close()
#                             logging.info(f"Random Forest scores plot saved to: {rf_plot_path}")
#                             print(f"Feature QC finsihed and saved in{feature_qc_dir}")
#                             #want to change feature names?
#                             rename_columns = input("Do you want to rename any columns (recommended for long feature names) ? (yes/no): ").strip().lower()
#                             if rename_columns == 'yes':
#                                 logging.info("Current Columns:")
#                                 for idx, column in enumerate(out_dat.columns):
#                                     logging.info(f"{idx}: {column}")
#                                 renaming_input = input(
#                                     "Enter the column name and new name separated by ':', each pair separated by commas-no space! (e.g., col1:new_col1,col2:new_col2): ").strip()
#                                 try:
#                                     rename_dict = dict(pair.split(':') for pair in renaming_input.split(',') if ':' in pair)
#                                     if rename_dict:
#                                         out_dat.rename(columns=rename_dict, inplace=True)
#                                         logging.info("Names changed successfully.")
#                                     else:
#                                         logging.warning("No valid renaming pairs provided.")
#                                 except Exception as e:
#                                     logging.error(f"Error renaming columns: {e}")
#                             out_dat.to_csv(os.path.join(results_folder, "data_with_new_names.csv"))
#                             logging.info(
#                                 f"Data with new names saved to: {os.path.join(results_folder, 'data_with_new_names.csv')}")
#                         except Exception as e:
#                             logging.error(f"An error occurred during renaming and feature qc process: {e}")
#                             return
#                         print("ready_for_norm!")
#                         #imputation and mormalization segment
#                         try:
#                             print("imputing data for clean normalization")
#                             #initialize simpleimputer for missing values (NaN) - median!
#                             imputer = SimpleImputer(strategy='median')
#                             #impute missing values (NaN)
#                             df_imputed = pd.DataFrame(imputer.fit_transform(out_dat), columns=out_dat.columns)
#                             logging.info("Imputation of missing values completed.")
#                             #replace zero values with NaN to use the same imputer
#                             df_imputed.replace(0, np.nan, inplace=True)
#                             logging.debug("Replaced zero values with NaN for further imputation.")
#                             #impute zero values (now NaNs) using the same strategy
#                             df_final_imputed = pd.DataFrame(imputer.fit_transform(df_imputed), columns=out_dat.columns)
#                             df_final_imputed.index = out_dat.index
#                             logging.info("Imputation of zero values completed.")
#                             #save imputed data
#                             final_imputed_path = os.path.join(results_folder, "df_final_imputed.csv")
#                             df_final_imputed.to_csv(final_imputed_path)
#                             logging.info(f"Imputed data saved to {final_imputed_path}")
#                             #normalization
#                             logging.info("Proceeding to normalization...")
#                             selected_data = df_final_imputed.copy()
#                             # Create a directory for normalization results
#                             normalization_dir = os.path.join(results_folder, "normalization")
#                             os.makedirs(normalization_dir, exist_ok=True)
#                             logging.info(f"Normalization results will be saved to: {normalization_dir}")
#                             #split selected_data based on PC column - error handling for missing 'PC' index
#                             if 'PC' not in selected_data.index.names:
#                                 raise ValueError("'PC' index not found in selected data. Ensure 'PC' is part of the index.")
#                             pc_groups = selected_data.groupby(level='PC', group_keys=False)
#                             df_reverted = pc_groups.apply(lambda x: x)
#                             data_3_path = os.path.join(results_folder, "df_proc_3.csv")
#                             df_reverted.to_csv(data_3_path)
#                             logging.info(f"Final processed data saved to {data_3_path}")
#                             #running over each group and performing normalization
#                             normalization_methods = {
#                                 "min-max_normalization": lambda x: (x - x.min()) / (x.max() - x.min()),
#                                 "central_log_normalization": lambda x: np.log(x + np.sqrt(x ** 2 + 1)),
#                                 "z-score_normalization": lambda x: (x - x.mean()) / x.std()
#                             }
#                             for pc_value, pc_group in pc_groups:
#                                 pc_dir = os.path.join(normalization_dir, f"df_{pc_value}")
#                                 os.makedirs(pc_dir, exist_ok=True)
#                                 logging.info(f"Normalization directory created at: {pc_dir}")
#                                 #original data
#                                 original_data_path = os.path.join(pc_dir, "original_data.csv")
#                                 pc_group.to_csv(original_data_path)
#                                 logging.info(f"Original data saved to {original_data_path}")
#                                 #normalization methods and saving results
#                                 normalized_data_dict = {}
#                                 for method_name, method_func in normalization_methods.items():
#                                     try:
#                                         normalized_data = method_func(pc_group)
#                                         normalized_data.to_csv(os.path.join(pc_dir, f"{method_name.lower().replace(' ', '_')}.csv"))
#                                         normalized_data_dict[method_name] = normalized_data
#                                         logging.info(f"{method_name} applied and saved to {method_name.lower().replace(' ', '_')}.csv")
#                                     except Exception as e:
#                                         logging.error(f"Error applying {method_name} for PC value {pc_value}: {e}")
#                                 #box-cox transformation - data must be positive
#                                 try:
#                                     transformed_data = pd.DataFrame()
#                                     for column in pc_group.columns:
#                                         if (pc_group[column] <= 0).any():
#                                             raise ValueError(
#                                                 f"Box-Cox transformation requires positive data, found non-positive values in {column}.")
#                                         transformed_column, _ = boxcox(pc_group[column])
#                                         transformed_data[column] = transformed_column
#                                     transformed_data.index = pc_group.index
#                                     transformed_data.to_csv(os.path.join(pc_dir, "box_cox_normalized_data.csv"))
#                                     logging.info(f"Box-Cox normalized data saved to box_cox_normalized_data.csv")
#                                 except Exception as e:
#                                     logging.error(f"Error applying Box-Cox transformation for PC value {pc_value}: {e}")
#                                 #log normalization
#                                 try:
#                                     log_data = pd.DataFrame()
#                                     for column in pc_group.columns:
#                                         #handle negative or zero values for log normalization
#                                         if (pc_group[column] <= 0).any():
#                                             raise ValueError(
#                                                 f"Log normalization requires positive data, found non-positive values in {column}.")
#                                         log_column = np.log(pc_group[column])
#                                         log_data[column] = log_column
#                                     log_data.index = pc_group.index
#                                     log_data.to_csv(os.path.join(pc_dir, "log_normalized_data.csv"))
#                                     logging.info(f"Log normalized data saved to log_normalized_data.csv")
#                                 except Exception as e:
#                                     logging.error(f"Error applying log normalization for PC value {pc_value}: {e}")
#                                 #QQ plots and histograms for each feature in each normalized data frame
#                                 normalized_dfs = [
#                                     ("Original Data", pc_group),
#                                     ("min-max_normalization", normalized_data_dict.get("min-max_normalization")),
#                                     ("central_log_normalization", normalized_data_dict.get("central_log_normalization")),
#                                     ("z-score_normalization", normalized_data_dict.get("z-score_normalization")),
#                                     ("Box_Cox Normalized Data", transformed_data),
#                                     ("Log Normalized Data", log_data)
#
#                                 ]
#                                 for df_name, df in normalized_dfs:
#                                     if df is None:
#                                         logging.warning(f"{df_name} is not available for PC value {pc_value}. Skipping plots.")
#                                         continue
#                                     df_dir = os.path.join(pc_dir, df_name.replace(' ', '_').lower())
#                                     os.makedirs(df_dir, exist_ok=True)
#                                     logging.info(f"Plot directory created at: {df_dir}")
#                                     for feature in df.columns:
#                                         try:
#                                             ## qq
#                                             plt.figure()
#                                             scipy.stats.probplot(df[feature], dist="norm", plot=plt)
#                                             plt.title(f"Q-Q-{feature}")
#                                             qq_plot_path_2 = os.path.join(df_dir, f"qq_2_plot_{feature}.png")
#                                             plt.savefig(qq_plot_path_2)
#                                             plt.close()
#                                             plt.clf()
#                                         except Exception as e:
#                                             logging.error(f"Error plotting Q-Q for {feature} in {df_name}: {e}")
#                                         # histogram
#                                         try:
#                                             plt.figure(figsize=(18, 12))
#                                             sns.histplot(df[feature], bins=20, kde=True, color='skyblue')
#                                             plt.title(f"Histogram for {feature} ({df_name})")
#                                             plt.xlabel("Value")
#                                             plt.ylabel("Frequency")
#                                             plt.tight_layout()
#                                             histogram_path = os.path.join(df_dir, f"histogram_{feature}.png")
#                                             plt.savefig(histogram_path)
#                                             plt.close()
#                                             logging.info(f"Histogram saved to: {histogram_path}")
#                                         except Exception as e:
#                                             logging.error(f"Error plotting histogram for {feature} in {df_name}: {e}")
#                             #user to select the normalized data frame to continue with
#                             try:
#                                 normalized_options = ["min-max_normalization", "central_log_normalization",
#                                                       "z-score_normalization",
#                                                       "Box_Cox Normalized Data", "Log Normalized Data"]
#                                 logging.info("Available normalization methods:")
#                                 for idx, option in enumerate(normalized_options):
#                                     logging.info(f"{idx + 1}: {option}")
#                                 selected_normalization = int(
#                                     input("Check normalized data and enter the number corresponding to the desired normalized data frame: ")) - 1
#                                 if selected_normalization < 0 or selected_normalization >= len(normalized_options):
#                                     raise ValueError("Selection out of range.")
#                                 selected_normalization_name = normalized_options[selected_normalization]
#                                 logging.info(f"Selected normalization method: {selected_normalization_name}")
#                                 #combine all PC values for the selected normalization into one final normalized data frame
#                                 final_normalized_data = pd.DataFrame()
#                                 for pc_value, pc_group in pc_groups:
#                                     pc_dir = os.path.join(normalization_dir, f"df_{pc_value}")
#                                     file_path = os.path.join(pc_dir, f"{selected_normalization_name.lower().replace(' ', '_')}.csv")
#                                     if not os.path.exists(file_path):
#                                         logging.warning(f"Normalized file not found: {file_path}. Skipping.")
#                                         continue
#                                     normalized_df = pd.read_csv(file_path, index_col=0)
#                                     final_normalized_data = pd.concat([final_normalized_data, normalized_df])
#                                     logging.debug(f"Data from {file_path} added to final normalized data.")
#                                 #save final normalized data
#                                 final_normalized_data_path = os.path.join(results_folder, "final_normalized_data.csv")
#                                 final_normalized_data.to_csv(final_normalized_data_path)
#                                 logging.info(f"Final normalized data saved to: {final_normalized_data_path}")
#                             except Exception as e:
#                                 logging.error(f"An error occurred during the imputation or normalization process: {e}")
#                                 return
#                         except Exception as e:
#                             logging.error(f"An unexpected error occurred: {e}")
#                             return
#                         finally:
#                             logging.info('All processing steps completed.')
#                         #dashboard HTML generation
#                         #use relative path for links
#                         print("creating dashboard for results summary")
#                         dashboard_page = os.path.join(main_results_dir, f'{item_name}_dashboard.html')
#                         dashboard_content = f"""
#                         <html>
#                         <head><title>Dashboard for Folder {item_name}</title></head>
#                         <body>
#                             <h1>Results for Folder: {item_name}</h1>
#                             <p>Results stored in: {results_folder}</p>
#                         """
#                         #add a section for each  file in the folder
#                         dashboard_content += "<h2>Processed Plates</h2>"
#                         for file_name in os.listdir(item_path):
#                             if file_name.endswith(".csv"):
#                                 results_folder = os.path.join(item_path, f"results_{file_name}_res")
#                                 barplot_dir = os.path.join(results_folder, "box_plots")
#                                 summary_dir = os.path.join(results_folder, "general_results")
#                                 barplot_relative_path = os.path.relpath(barplot_dir, main_results_dir)
#                                 summary_relative_path = os.path.relpath(summary_dir, main_results_dir)
#                                 dashboard_content += f"<h3>Results for Plate: {file_name}</h3>"
#                                 dashboard_content += "<ul>"
#                                 #bar plots for important features
#                                 for feature in important_features:
#                                     bar_plot_html = os.path.join(barplot_relative_path, f'{feature}_bar_plot.html')
#                                     dashboard_content += f'<li><a href="{bar_plot_html}">Bar Plot for {feature}</a></li>'
#
#                                 #PCA plot
#                                 pca_plot_html_id = os.path.join(summary_relative_path, f'pca_plot_with_id_{file_name[:-4]}.html')
#                                 dashboard_content += f'<li><a href="{pca_plot_html_id}">PCA Plot for {file_name}</a></li>'
#                                 #PCA plot2
#                                 pca_plot_html_reg = os.path.join(summary_relative_path, f'pca_plot_wo_id_{file_name[:-4]}.html')
#                                 dashboard_content += f'<li><a href="{pca_plot_html_reg}">PCA Plot for {file_name}</a></li>'
#                                 #number of cells plot for the current plate
#                                 count_plot = os.path.join(summary_relative_path,
#                                                                 f'number_of_cells_bar_plot.pdf')
#                                 dashboard_content += f'<li><a href="{count_plot}">number of cells Plot for {file_name}</a></li>'
#                                 #sum file
#                                 sum_file = os.path.join(summary_relative_path,
#                                                                 f'init_statistical_summary.xlsx')
#                                 dashboard_content += f'<li><a href="{sum_file}">stat_summary for {file_name}</a></li>'
#                                 #heatmap
#                                 heat_plot_pdf_reg = os.path.join(summary_relative_path,
#                                                                  f'heatmap.pdf')
#                                 dashboard_content += f'<li><a href="{heat_plot_pdf_reg}">heatmap for {file_name}</a></li>'
#
#                                 dashboard_content += "</ul>"
#
#                         #closing the HTML content
#                         dashboard_content += "</body></html>"
#
#                         #updated dashboard content to the HTML file
#                         with open(dashboard_page, 'w') as f:
#                             f.write(dashboard_content)
#                         logging.info(f"Dashboard updated with  plots and saved for {item_name} at {dashboard_page}")
#                     logging.info(f"{file_name} proccessed succesfully")
#     except Exception as e:
#         logging.error(f"failed proccessing {plate_df}")
#
# ########
# ##this should be the main folder /results
# # main_folder_path = "noam_ghf/results"
# # ##define redundant columns - any columns that are un-needed upfront
# # unwanted_columns=[Number of Analyzed Fields, Time [s], Plane, Timepoint, Cell Count]
# # ## cell count feat
# # cell_count_feat = "selected_cells - Number of Objects"
# # ## which columns define the groups tested? select all relvant.
# # index_columns = [Compound, Concentration, Severity, CELL ID, Cell Type, PC]
# # #define important columns - all columns that should be tested for the analysis, in addition to technical columns.
# # important_features = [spots_chanel_3_final - Corrected Spot Intensity - Mean per Well, spots_chanel_3_final - Intensity_Spot chanel_3_final Mean - Mean per Well,
# #                       "modified_ir_chanel_3_at - Intensity modified_ir_chanel_3_at Mean - Mean per Well", "selected_cells - chanel_3_Intensity Mean - Mean per Well",
# #                       "total_ir_chanel_3_area_normalized", "total_ir_chanel_3_intensity_normalized", "selected_cells - chanel_3 Gabor Max 2 px w2 - Mean per Well",
# #                       "selected_cells - Cell Area [ÂµmÂ²] - Mean per Well", "selected_cells - chanel_1 Area"]
#
# ## run constant
# #process_plate_data_on_folders(main_folder_path, unwanted_columns, important_features, index_columns, cell_count_feat)
#
#
# def gather_user_lists():
#
#     print("Choose appropriate columns for your analysis!")
#     #options based on loaded data
#     print("usual unwanted columns: Number of Analyzed Fields, Height, Plane, Time, Timepoint, Cell Count")
#     print("enter unwanted columns:")
#     list1 = input("enter unwanted columns:(comma-separated): ").split(',')
#
#     print("usual important columns: modified_spots_chanel_3 - Total Spot Area - Mean per Well, all_cells - chanel_3 SER Spot 1 px - Mean per Well, spots_chanel_3_final - Spot to Region Intensity - Mean per Well, spots_chanel_3_final - Corrected Spot Intensity - Mean per Well, spots_chanel_3_final - Region Intensity - Mean per Well, ir_chanel_3_at - Intensity ir_chanel_3_at chanel_3 Mean - Mean per Well,  all_cells - chanel_3 Gabor Max 2 px w2 - Mean per Well, all_cells - chanel_3 SER Valley 1 px - Mean per Well, all_cells - chanel_3_intensity Mean - Mean per Well, all_cells - chanel_2_intensity Mean - Mean per Well, all_cells - chanel_1_intensity Mean - Mean per Well, chan3_intensitymean_div_chan2area, spotschan3_intensitymean_div_chan2area")
#     print("enter important_features:")
#     list2 = input("enter important_features:(comma-separated): ").split(',')
#
#     print("must index columns!!!: PC, cell_type, cell_id, compound, concentration")
#     print("enter index_columns:")
#     list3 = input("enter index_columns: (comma-separated): ").split(',')
#
#     #list of strings
#     list1 = [item.strip() for item in list1]
#     list2 = [item.strip() for item in list2]
#     list3 = [item.strip() for item in list3]
#
#     return list1, list2, list3
#
# #RUN MAIN
# def main(main_folder, main_folder_path, final_folder_path):
#     print("Welcome to Mitochondrial Activity analysis (TMRE section)!")
#     process_folders(main_folder)
#     prompt_for_review()
#     process_data_combine(main_folder_path)
#     prompt_for_review()
#     list1, list2, list3 = gather_user_lists()
#     process_plate_data_on_folders(final_folder_path, list1, list2, list3)
#     print("Completed Mitochondrial Activity analysis (TMRE section), good luck with statistics:)")
# ## usual cell_count_feat names
# # selected_cells - Number of Objects
# ## all cells
# ##none
# ### usual unwanted columns
# # unwanted_columns=Number of Analyzed Fields, Time [s], Plane, Timepoint, Cell Count
# ## usual indexes
# ## must index - PC, cell_type, cell_id, compound, concentration
# # index_columns = PC, cell_type, cell_id, compound, concentration
# ## all_cells - Number of Objects
# ## usual important features
# # important_features = modified_spots_chanel_3 - Total Spot Area - Mean per Well, all_cells - chanel_3 SER Spot 1 px - Mean per Well, spots_chanel_3_final - Spot to Region Intensity - Mean per Well, spots_chanel_3_final - Corrected Spot Intensity - Mean per Well, spots_chanel_3_final - Region Intensity - Mean per Well, ir_chanel_3_at - Intensity ir_chanel_3_at chanel_3 Mean - Mean per Well,  all_cells - chanel_3 Gabor Max 2 px w2 - Mean per Well, all_cells - chanel_3 SER Valley 1 px - Mean per Well, all_cells - chanel_3_intensity Mean - Mean per Well, all_cells - chanel_2_intensity Mean - Mean per Well, all_cells - chanel_1_intensity Mean - Mean per Well, chan3_intensitymean_div_chan2area, spotschan3_intensitymean_div_chan2area
#
# if __name__ == "__main__":
#     main_folder = input("Enter folder name(containing a folder with .txt files): ")
#     main_folder_path = input("Enter: the main folder name/results ")
#     final_folder_path = input("Enter: the main folder name/MA-TMRE_result_combined ")
#     main(main_folder, main_folder_path, final_folder_path)



##### check CLR save ?
### validate qq + hist forl ong names
## add proteceion for outlier
## printing improve

### MA -Combined Plates
#required libraries :
import os
import logging
import pandas as pd
import numpy as np
import os
import glob
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
import itertools
import statsmodels.api as sm
from scipy.stats import boxcox, zscore
from sklearn.neighbors import LocalOutlierFactor
from sklearn.ensemble import RandomForestRegressor
import kaleido
from openpyxl import Workbook
from openpyxl.utils.dataframe import dataframe_to_rows
from scipy.stats import zscore
from sklearn.impute import SimpleImputer
import scipy
from sklearn.decomposition import PCA
import matplotlib.patches as mpatches
from colorama import Fore, Style, init
init(autoreset=True)
####
# IF combine plates vs
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
def process_folders(main_folder):
    try:
        #check main folder
        if not os.path.isdir(main_folder):
            print(f"Error: The folder '{main_folder}' does not exist.")
            return
        #run through subfolders in main folder
        for subfolder_name in os.listdir(main_folder):
            subfolder_path = os.path.join(main_folder, subfolder_name)
            if os.path.isdir(subfolder_path):
                print(f"\nProcessing subfolder: {subfolder_name}")
                #user to choose value for the 'PC' column- plate condition- must.
                pc_value = input("Enter the value for the 'PC' column(e.g- plate condition ), must add (numerical or string): ").strip()
                if not pc_value:
                    print("Error: 'PC' column value cannot be empty. skipping this subfolder.")
                    continue
                #run through text files in subfolder
                for txt_file in glob.glob(os.path.join(subfolder_path, '*.txt')):
                    print(f"Reading file: {txt_file}")
                    try:
                        #Load file and process
                        df = pd.read_csv(txt_file, skiprows=8, engine='python', sep='\t')
                    except Exception as e:
                        print(f"Error reading file '{txt_file}': {e}")
                        continue
                    #adding 'PC' column with the value
                    df['PC'] = pc_value
                    #extract plate name from file name
                    plate_name = os.path.splitext(os.path.basename(txt_file))[0]
                    print(f"Processing plate: {plate_name}")
                    #output folder
                    output_folder = os.path.join(main_folder, "results", f"{subfolder_name}_csv_files")
                    os.makedirs(output_folder, exist_ok=True)
                    #save CSV
                    csv_filename = f"{plate_name}_csv.csv"
                    #prompt user for flagged wells to remove
                    flagged_wells_input = input(
                        f"Enter flagged wells(bad wells) to remove for plate '{plate_name}' (e.g., row_column, ...) or 'none' if none: ").strip()
                    #process flagged wells
                    flagged_wells = []
                    if flagged_wells_input.lower() != 'none':
                        flagged_wells = [fw.strip() for fw in flagged_wells_input.split(',') if fw.strip()]
                    # remove flagged wells
                    for flagged_well in flagged_wells:
                        try:
                            row, column = flagged_well.split('_')
                            row = int(row)
                            column = int(column)
                            df = df.loc[~((df['Row'] == row) & (df['Column'] == column))]
                            print(f"Flagged well removed: Row {row}, Column {column}")
                        except Exception as e:
                            print(f"Error processing flagged well '{flagged_well}': {e}")
                            continue
            #editing column names
                    try:
                        df.columns = [col.split("[")[0].strip() if "[" in col else col for col in df.columns]
                        print("Column names cleaned.")
                    except Exception as e:
                        logging.error(f"Error cleaning column names: {e}")
                        return
                    #save modified DF
                    csv_path = os.path.join(output_folder, csv_filename)
                    try:
                        df.to_csv(csv_path, index=False)
                        print(f"File '{csv_filename}' saved in folder '{output_folder}'")
                    except Exception as e:
                        print(f"Error saving file '{csv_filename}': {e}")
                        continue
        #done processing all subfolders
        print("\nAll subfolders processed successfully.")
        print(f"review the files saved in {main_folder}/results.")

    except Exception as e:
        print(f"Unexpected error occurred: {e}")
##################
def prompt_for_review():
    input("review the csv files. when ready, press enter to continue...")
######################
def process_data_combine(folder_path):
    try:
        #output result folder
        output_folder = os.path.join(main_folder, "MA-TMRE_result_combined")
        os.makedirs(output_folder, exist_ok=True)
        result_combine_dir = os.path.join(main_folder, "MA-TMRE_result_combined/combined_csv_file")
        if not os.path.exists(result_combine_dir):
            os.makedirs(result_combine_dir, exist_ok=True)
        #subdirectories
        subdirectories = [subdir for subdir in os.listdir(folder_path) if os.path.isdir(os.path.join(folder_path, subdir))]
        combined_dfs = []
        #run on each subdirectory
        for subdir in subdirectories:
            sub_dir_path = os.path.join(folder_path, subdir)
            print("Processing files in:", sub_dir_path)
            dfs = []
            for file in os.listdir(sub_dir_path):
                if file.endswith('.csv'):
                    file_path = os.path.join(sub_dir_path, file)
                    df = pd.read_csv(file_path, encoding='unicode_escape')
                    df['plate'] = os.path.splitext(file)[0]  #add plate column- important
                    dfs.append(df)
        print("plate column added to data frame")
        combined_df = pd.concat(dfs, ignore_index=True)
        combined_dfs.append(combined_df)
        combined_df = pd.concat(combined_dfs, ignore_index=True)
        print("Data frames combined!")
        combined_df['well'] = combined_df['Row'].astype('str') + combined_df['Column'].astype('str')
        combined_df['well'] = combined_df['well'].astype('str')
        print("well column added to data frame")
        combined_df.to_csv(os.path.join(result_combine_dir, "combined_df_data.csv"), index=False)
        print(f"review the files saved in {main_folder}/MA-TMRE_result_combined/combined_csv_file.")
    except Exception as e:
        logging.error(f"failed to combine plates...{e}")
####################
def process_plate_data_on_folders(main_folder, unwanted_columns, important_features, index_columns):
    try:
    ## open main folder for dash
        main_results_dir = os.path.join(main_folder, "main_results_folder")
        os.makedirs(main_results_dir, exist_ok=True)
        unwanted_columns = list(set(unwanted_columns))
        important_features = list(set(important_features))
        index_columns = list(set(index_columns))
        #running on each item in the main input folder
        for item_name in os.listdir(main_folder):
            item_path = os.path.join(main_folder, item_name)
            if os.path.isdir(item_path):
                #running on each csv in each subfolder
                for file_name in os.listdir(item_path):
                    if file_name.endswith(".csv"):
                        file_path = os.path.join(item_path, file_name)
                        print(f"Processing file: {file_path}")
                        #load each plate
                        plate_df = pd.read_csv(file_path)
                        #results folder with the name of the csv file
                        results_folder = os.path.join(item_path, f"results_{file_name}_res")
                        os.makedirs(results_folder, exist_ok=True)
                        #drop unwanted columns & start proc
                        try:
                            plate_df = plate_df.drop(columns=unwanted_columns)
                            print(f"Dropped unwanted columns: {unwanted_columns}")
                        except KeyError as e:
                            logging.warning(f"Error dropping unwanted columns: {e}")
                        #checking important features
                        print("These are the important columns chosen for the analysis:")
                        print(important_features)
                        missing_features = [feature for feature in important_features if feature not in plate_df.columns]
                        if missing_features:
                            logging.warning(f"Missing important features: {missing_features}")
                            important_features = [f for f in important_features if f in plate_df.columns]
                        else:
                            print("all selected features found!")
                        if not important_features:
                            logging.error("No valid important features available. Skipping.")
                            continue
                        try:
                            print("These are the index columns chosen for the analysis:")
                            print(index_columns)
                            index_columns.append('well')
                            index_columns.append('plate')
                            print("'well' & 'plate' columns added to index_columns.")
                        except Exception as e:
                            logging.error(f"Error adding 'well' and 'plate' columns to index columns: {e}")
                        #generate combinations of index_columns for new indexes
                        try:
                            new_index_columns = []
                            for r in range(1, len(index_columns) + 1):
                                for comb in itertools.combinations(index_columns, r):
                                    new_index_columns.append(comb)
                            print(f"Generated {len(new_index_columns)} new index column combinations.")
                        except Exception as e:
                            logging.error(f"Error generating new index column combinations: {e}")
                        #create new index columns based on combinations
                        try:
                            new_cols = {}
                            for cols in new_index_columns:
                                col_name = '_'.join(cols)
                                if col_name not in plate_df.columns:
                                    new_cols[col_name] = plate_df[list(cols)].apply(lambda x: '_'.join(map(str, x)), axis=1)
                                else:
                                    print(f"{col_name} already exists and will not be added")
                            plate_df = pd.concat([plate_df, pd.DataFrame(new_cols)], axis=1)
                            print("Created new index columns based on combinations.")
                        except Exception as e:
                            logging.error(f"Error creating new index columns: {e}")
                        ## keep only new_index and imp features
                        try:
                            print("subsetting data to focus on important features & index columns...")
                            columns1 = plate_df[index_columns]
                            columns2 = plate_df[important_features]
                            plate_df_int = pd.concat([columns1, columns2, pd.DataFrame(new_cols)], axis=1)
                        except Exception as e:
                            logging.error(f"Failed to subset{e}")
                            plate_df_int = plate_df.copy()
                        try:
                            # remove rows and columns with ≥90% zero or NA values
                            filtered_df = plate_df_int.drop(columns=index_columns, errors='ignore')
                            filtered_df = filtered_df.dropna(thresh=0.1 * len(filtered_df), axis=0)
                            filtered_df = filtered_df.replace(0, np.nan)
                            filtered_df = filtered_df.dropna(thresh=0.1 * len(filtered_df.columns), axis=1)
                            final_df = pd.concat([columns1, filtered_df], axis=1)
                            print("Dropped rows and columns with ≥90% zero or NA values.")
                            plate_df_2 = final_df.copy()
                        except Exception as e:
                            logging.error(f"Error dropping rows/columns with high NA/zero values: {e}")
                        ## add cell count feat
                        try:
                            print("usual cell_count_feat names: all_cells - Number of Objects")
                            cell_count_feat = input("Enter cell_count_feat(check in csv if not sure): ")
                            plate_df_2[cell_count_feat] = plate_df[cell_count_feat]
                            plate_df_2[cell_count_feat] = pd.to_numeric(plate_df_2[cell_count_feat])
                        except Exception as e:
                            logging.error(f"failed to locate cell count feat {e}")
                        ##user to remove rows based on any avaialble columns
                        try:
                            print("Section to remove rows based on chosen column, choose column following by valuesto remove, otherwise, type no")
                            while True:
                                print("Columns in the dataset:")
                                for idx, column in enumerate(plate_df_2.columns):
                                    print(f"{idx}: {column}")
                                column_to_filter = input("Enter the column name to filter by (or 'no' to skip): ").strip()
                                if column_to_filter.lower() == 'no':
                                    break
                                if column_to_filter not in plate_df_2.columns:
                                    print("Invalid column name. Please try again.")
                                    continue
                                print(f"Values in the selected column '{column_to_filter}':")
                                print(plate_df_2[column_to_filter].unique())
                                values_to_remove = input("Enter the values to remove (comma-separated): ").split(',')
                                values_to_remove = [val.strip() for val in values_to_remove]
                                plate_df_2 = plate_df_2[~plate_df_2[column_to_filter].isin(values_to_remove)]
                                print(f"Rows with values {values_to_remove} in column '{column_to_filter}' have been removed.")
                        except Exception as e:
                            logging.error(f"Error removing custom column values: {e}")
                            return
                        #remove additional columns if needed
                        try:
                            print("Section to remove columns from the data, choose column to remove, otherwise, type none")
                            for idx, column in enumerate(plate_df_2.columns):
                                print(f"{idx + 1}: {column}")
                            features_to_remove = input("Enter comma-separated column names to remove (or 'none'): ").strip()
                            if features_to_remove.lower() != 'none':
                                features_to_remove = [col.strip() for col in features_to_remove.split(',')]
                                plate_df_2.drop(columns=features_to_remove, inplace=True)
                                print(f"Features removed: {features_to_remove}")
                        except KeyError as e:
                            logging.warning(f"Some features to remove were not found: {e}")
                        except Exception as e:
                            logging.error(f"Error removing features: {e}")
                        #IDs to remove ?
                        try:
                            print("Section to remove cell id values from the data, choose id to remove, otherwise, type none")
                            print(plate_df_2["cell_id"].unique())
                            remove_id = input("Enter ID value to remove (or 'none'): ").strip()
                            if remove_id.lower() != 'none':
                                plate_df_2 = plate_df_2[~plate_df_2['cell_id'].astype(str).str.contains(remove_id)]
                                logging.info(f"Rows with ID '{remove_id}' removed.")
                        except KeyError:
                            logging.warning("Column 'ID' not found. Skipping ID removal.")
                        except Exception as e:
                            logging.error(f"Error removing rows based on ID: {e}")
                        #save first df
                        plate_df_2.to_csv(os.path.join(results_folder, "df_proc_1.csv"), index=False)
                        print("first processed data frame saved to result directory")
                        #directory for main results each subfolder
                        general_results_dir = os.path.join(results_folder, "general_results")
                        os.makedirs(general_results_dir, exist_ok=True)
                        #index column for statistical summary
                        try:
                            print(new_index_columns)
                            index_col = input("Enter the index column for statistical summary: seperated by '_' ").strip()
                            summary_df = plate_df_2.groupby(index_col).describe()
                            print(f"Generated statistical summary grouped by '{index_col}'.")
                        except KeyError:
                            logging.error(f"Index column '{index_col}' not found in DataFrame.")
                            summary_df = pd.DataFrame()
                        except Exception as e:
                            logging.error(f"Error generating statistical summary: {e}")
                            summary_df = pd.DataFrame()
                        #save statistical summary
                        try:
                            summary_excel_path = os.path.join(general_results_dir, "init_statistical_summary.xlsx")
                            summary_df.to_excel(summary_excel_path)
                            print(f"Statistical summary saved to '{summary_excel_path}'")
                        except Exception as e:
                            logging.error(f"Error saving statistical summary to Excel: {e}")
                        #directory for cell count results
                        print("cell count process initiated")
                        cell_count_results_dir = os.path.join(results_folder, "cell_count_results")
                        os.makedirs(cell_count_results_dir, exist_ok=True)
                        #histogram for cell count
                        try:
                            if cell_count_feat in plate_df_2.columns:
                                fig = px.histogram(plate_df_2, x=cell_count_feat, title="Histogram of Cell Count")
                                histogram_file_path = os.path.join(cell_count_results_dir, "histogram.html")
                                fig.write_html(histogram_file_path)
                                print(f"Cell count histogram saved to {histogram_file_path}")
                            else:
                                logging.warning(f"{cell_count_feat} column not found. histogram not generated.")
                        except Exception as e:
                            logging.error(f"Error generating histogram: {e}")
                        #statistics for cell count
                        try:
                            if cell_count_feat in plate_df_2.columns:
                                stats = plate_df_2[cell_count_feat].describe()
                                stats_file_path = os.path.join(cell_count_results_dir, "cell_count-statistics.txt")
                                with open(stats_file_path, "w") as f:
                                    f.write(stats.to_string())
                                print(f"Cell count statistics saved to {stats_file_path}")
                            else:
                                logging.warning(f"{cell_count_feat} column not found. statistics not generated.")
                        except Exception as e:
                            logging.error(f"Error creating cell count statistics: {e}")
                        #input for minimum and maximum cell count values
                        try:
                            min_input = input("Enter the minimum value for cell count (leave blank to skip): ").strip()
                            max_input = input("Enter the maximum value for cell count (leave blank to skip): ").strip()
                            min_value = float(min_input) if min_input else None
                            max_value = float(max_input) if max_input else None
                            if min_value is not None and max_value is not None:
                                plate_df_2 = plate_df_2[(plate_df_2[cell_count_feat] >= min_value) & (
                                        plate_df_2[cell_count_feat] <= max_value)]
                                print(f"Filtered data with cell count between {min_value} and {max_value}")
                            elif min_value is not None:
                                plate_df_2 = plate_df_2[plate_df_2[cell_count_feat] >= min_value]
                                print(f"Filtered data with cell count >= {min_value}")
                            elif max_value is not None:
                                plate_df_2 = plate_df_2[plate_df_2[cell_count_feat] <= max_value]
                                print(f"Filtered data with cell count <= {max_value}")
                            else:
                                print("No filtering applied to cell count.")
                        except ValueError:
                            logging.error("Invalid input for minimum or maximum cell count. Skipping filtering.")
                        except Exception as e:
                            logging.error(f"Error filtering cell count: {e}")
                        plate_df_2.to_csv(os.path.join(results_folder, "df_proc_2.csv"))
                        print("second processed data frame saved to result directory")
                        print(new_index_columns)
                        #directory for box plots results
                        box_plot_dir = os.path.join(results_folder, "box_plots")
                        os.makedirs(box_plot_dir, exist_ok=True)
                        #prompting index column for box plots
                        index_col = input("Enter the index column for box plot vis: ")
                        bar_df = plate_df_2.copy()
                        # box plots for selected index- important features only
                        try:
                            print("box plotting selected features!")
                            for col in bar_df:
                                if col in important_features:
                                    bp = px.box(bar_df, x=index_col, y=bar_df[col], points="all", color=index_col, notched=True)
                                    bp.update_traces(quartilemethod="inclusive")  # or "inclusive", or "linear" by default
                                    bp.update_layout(
                                        font_family="Arial",
                                        font_color="Black",
                                        font_size=20,
                                        font=dict(
                                            family="Arial",
                                            size=20,
                                            color="Black"
                                        )
                                    )
                                    print("one_down")
                                    bp.write_html(os.path.join(box_plot_dir, f'{col}_bar_plot.html'))
                                    #bp.write_image(os.path.join(box_plot_dir, f"box_plot{col}.pdf"), engine="kaleido")
                            print("box plots done!")
                        except Exception as e:
                            logging.error(f"failed box plotting{e}")
                        #directory for summary files
                        summary_files_dir = os.path.join(results_folder, "summary_files")
                        os.makedirs(summary_files_dir, exist_ok=True)
                        #creating excel  with sheets with split data and summary stats
                        try:
                            print("creating excel files and statistical summary & ppt for combination of groups!")
                            group_columns = ['cell_type', 'compound', 'concentration']
                            excel_file_path = os.path.join(summary_files_dir, f'summary_file.xlsx')
                            with pd.ExcelWriter(excel_file_path, engine='openpyxl') as writer:
                                for name, group in plate_df_2.groupby(group_columns):
                                    try:
                                        sheet_name = f"{name[0]}_{name[1]}_{name[2]}"
                                        group.to_excel(writer, sheet_name=sheet_name, index=False)
                                        ####summary statistics to another sheet
                                        summary_df = group.describe(include='all').T[
                                            ['mean', 'std', 'min', '25%', '50%', '75%', 'max']]
                                        summary_df.to_excel(writer, sheet_name=f"{sheet_name}_Summary")
                                    except ValueError:
                                        pass
                                if not writer.sheets:
                                    pd.DataFrame({'Dummy': [1]}).to_excel(writer, sheet_name='DummySheet', index=False)
                            #ppt with slides for each unique combination
                            print("ppt file generated!")
                            ppt_file_path = os.path.join(summary_files_dir, f'slides.pptx')
                            prs = Presentation()
                            for name, _ in plate_df_2.groupby(group_columns):
                                slide = prs.slides.add_slide(prs.slide_layouts[5])
                                title = slide.shapes.title
                                title.text = f"Combination: {name[0]}, {name[1]}, {name[2]}"
                            prs.save(ppt_file_path)
                        except Exception as e:
                            logging.error(f"failed creating summary{e}")
                        #bar plot for Number of Cells
                        try:
                            print("working on final cell count vis!")
                            group_columns = ['cell_type', 'compound', 'concentration', 'cell_id']
                            plt.figure(figsize=(20, 18))
                            for name, group in plate_df_2.groupby(group_columns):
                                try:
                                    plt.bar(str(name), group[cell_count_feat].mean(),
                                            yerr=group[cell_count_feat].std(), capsize=5,
                                            label=str(name))
                                except ValueError:
                                    pass
                            plt.title(f'Bar Plot for Number of Cells - {file_name}')
                            plt.xlabel('cell_type, compound, concentration, cell_id')
                            plt.ylabel('Number of Cells')
                            plt.xticks(rotation=45, ha='right')
                            plt.legend()
                            plt.savefig(os.path.join(general_results_dir, f'number_of_cells_bar_plot.pdf'))
                            plt.close()
                        except Exception as e:
                            logging.error(f"failed bar plotting cell number{e}")
                        ## remove cell feat
                        try:
                            plate_df_2 = plate_df_2.drop(columns=[cell_count_feat])
                        except Exception as e:
                            logging.error(f"failed to remove {cell_count_feat}...{e}")
                        #heatmap for selected features (mean values & normalized)
                        try:
                            heat_df = plate_df_2.copy()
                            index_col = input("Enter the index column for heatmap: ")
                            print("heatmap prep!")
                            heat_df = heat_df.set_index(index_col, drop=True)
                            numeric_features = heat_df[important_features].select_dtypes(include=np.number)
                            numeric_features_normalized = (numeric_features - numeric_features.mean()) / numeric_features.std()
                            numeric_features_normalized = numeric_features_normalized.groupby(index_col).mean()
                            lan = plt.rcParams['font.family'] = ['Arial']
                            plt.figure(figsize=(20, 12))
                            font = {'family': 'Arial',
                                    'size': 22}
                            plt.rc('font', **font)
                            plt.yticks(family='Arial', rotation="horizontal", size=10)
                            plt.xticks(family='Arial', rotation="vertical", size=8)
                            plt.margins(0.1)
                            plt.subplots_adjust(bottom=0.3)
                            sns.heatmap(numeric_features_normalized, annot=True, cmap='coolwarm')
                            plt.title(f'Heatmap for Selected Numeric Features - {file_name}')
                            plt.tight_layout()
                            plt.savefig(os.path.join(general_results_dir, f'heatmap.pdf'))
                            plt.close()
                            print("heatmap done!")
                        except Exception as e:
                            logging.error(f"Heatmap failed: {e}")
                        #PCA based on all numeric features without ID
                        try:
                            print("creating pca for group with no ID and with id")
                            pca = PCA(n_components=2)
                            # Impute missing values in numeric features
                            numeric_features = plate_df_2.select_dtypes(include=np.number)
                            imputer = SimpleImputer(strategy='mean')
                            numeric_features_imputed = imputer.fit_transform(numeric_features)
                            #interactive pcas
                            pca_result_without_id = pca.fit_transform(numeric_features_imputed)
                            label_without_id = plate_df_2['cell_type'].astype(str) + '_' + \
                                            plate_df_2['compound'].astype(str) + '_' + \
                                            plate_df_2['concentration'].astype(str)
                            color_labels_wo_id = pd.factorize(label_without_id)[0]
                            pca_df1 = pd.DataFrame({
                                'PC1': pca_result_without_id[:, 0],
                                'PC2': pca_result_without_id[:, 1],
                                'Label': label_without_id,
                                'Color': color_labels_wo_id,
                                'cell_id': plate_df_2['cell_id']
                            })
                            fig = px.scatter(pca_df1, x='PC1', y='PC2', color='Label',
                                             title=f'PCA Plot for All Numeric Features without CELL ID - {file_name}', hover_data=['cell_id'],
                                             color_continuous_scale='viridis')
                            fig.update_layout(legend_title='Cell Type, Compound, Concentration')
                            fig.add_annotation(text=f'Explained Variance Ratio: {pca.explained_variance_ratio_}', xref='paper',
                                               yref='paper',
                                               x=0.98, y=0.02, showarrow=False, font=dict(size=10))
                            fig.write_html(os.path.join(general_results_dir, f"pca_plot_wo_id_{file_name[:-4]}.html"))
                            #fig.write_image(os.path.join(general_results_dir, f'pca_plot_without_id2.pdf'), engine="kaleido")
                            # PCA based on all numeric features with CELL ID-2 interactive
                            pca_result_with_id = pca.fit_transform(numeric_features_imputed)
                            # A label for each unique combination of Cell Type, Compound, Concentration, and CELL ID
                            label_with_id = plate_df_2['cell_type'].astype(str) + '_' + \
                                               plate_df_2['compound'].astype(str) + '_' + \
                                               plate_df_2['concentration'].astype(str) + '_' + \
                                               plate_df_2['cell_id'].astype(str)
                            color_labels_with_id = pd.factorize(label_with_id)[0]
                            # Create DataFrame for Plotly
                            pca_df2 = pd.DataFrame({
                                'PC1': pca_result_with_id[:, 0],
                                'PC2': pca_result_with_id[:, 1],
                                'Label': label_with_id,
                                'Color': color_labels_with_id,
                                'cell_id': plate_df_2['cell_id']
                            })
                            fig = px.scatter(pca_df2, x='PC1', y='PC2', color='Label',
                                             title=f'PCA Plot for All Numeric Features with CELL ID - {file_name}', hover_data=['cell_id'],
                                             color_continuous_scale='viridis')
                            fig.update_layout(legend_title='Cell Type, Compound, Concentration, ID')
                            fig.add_annotation(text=f'Explained Variance Ratio: {pca.explained_variance_ratio_}', xref='paper',
                                               yref='paper',
                                               x=0.98, y=0.02, showarrow=False, font=dict(size=10))
                            fig.write_html(os.path.join(general_results_dir, f"pca_plot_with_id_{file_name[:-4]}.html"))
                        except Exception as e:
                            logging.error(f"PCA failed: {e}")
                        print("pca done!")
                        #directory for outlier results
                        try:
                            outlier_files_dir = os.path.join(results_folder, "outlier_files")
                            os.makedirs(outlier_files_dir, exist_ok=True)
                            for col in plate_df_2.columns:
                                if col in index_columns:
                                    plate_df_2[col] = plate_df_2[col].astype(str)
                            perform_outlier_detection = input(
                                "Do you want to perform outlier detection? (yes/no): ").strip().lower()
                            if perform_outlier_detection == 'yes':
                                print("Applying basic Outlier methods...")
                                # prompting index column for outliers detection
                                print("Columns available for outlier detection:")
                                for idx, column in enumerate(plate_df_2.columns):
                                    print(f"{idx}: {column}")
                                index_column = input("Choose one column as index for outlier detection: ")
                                out_dat_init = plate_df_2.copy()
                                # final filtered DF
                                init_file_path = os.path.join(outlier_files_dir,
                                                               'out_dat_init.csv')
                                out_dat_init.to_csv(init_file_path, index=False)
                                #outlier detection functions options
                                ## perc
                                def compute_percentiles(df, output_file):
                                    up_bounds = []
                                    low_bounds = []
                                    above_count = []
                                    below_count = []
                                    numeric_cols = df.select_dtypes(include='number').columns
                                    for col in numeric_cols:
                                        percentile_low = df[col].quantile(0.01)
                                        percentile_high = df[col].quantile(0.99)
                                        up_bound = percentile_high
                                        low_bound = percentile_low
                                        above = df[df[col] > percentile_high]
                                        below = df[df[col] < percentile_low]
                                        above_count.append(len(above))
                                        below_count.append(len(below))
                                        df = df[(df[col] >= percentile_low) & (df[col] <= percentile_high)]
                                        up_bounds.append(up_bound)
                                        low_bounds.append(low_bound)
                                    output = pd.DataFrame(
                                        {'Column': numeric_cols, 'Up Bound': up_bounds, 'Low Bound': low_bounds,
                                         'Above Count': above_count,
                                         'Below Count': below_count})
                                    output.to_csv(output_file, index=False)
                                    return df
                                ## iqr
                                def compute_iqr(df, output_file):
                                    up_bounds = []
                                    low_bounds = []
                                    above_count = []
                                    below_count = []
                                    numeric_cols = df.select_dtypes(include='number').columns
                                    for col in numeric_cols:
                                        percentile25th = df[col].quantile(0.25)
                                        percentile75th = df[col].quantile(0.75)
                                        iqr = percentile75th - percentile25th
                                        up_bound = percentile75th + 1.5 * iqr
                                        low_bound = percentile25th - 1.5 * iqr
                                        above = df[df[col] > up_bound]
                                        below = df[df[col] < low_bound]
                                        above_count.append(len(above))
                                        below_count.append(len(below))
                                        df = df[(df[col] >= low_bound) & (df[col] <= up_bound)]
                                        up_bounds.append(up_bound)
                                        low_bounds.append(low_bound)
                                    output = pd.DataFrame(
                                        {'Column': numeric_cols, 'Up Bound': up_bounds, 'Low Bound': low_bounds,
                                         'Above Count': above_count,
                                         'Below Count': below_count})
                                    output.to_csv(output_file, index=False)
                                    return df
                                ## zscore
                                def compute_scores(df, output_file):
                                    up_scores = []
                                    low_scores = []
                                    above_count = []
                                    below_count = []
                                    numeric_cols = df.select_dtypes(include='number').columns
                                    for col in numeric_cols:
                                        up_score = df[col].mean() + 3 * df[col].std()
                                        low_score = df[col].mean() - 3 * df[col].std()
                                        above = df[df[col] > up_score]
                                        below = df[df[col] < low_score]
                                        above_count.append(len(above))
                                        below_count.append(len(below))
                                        df = df[(df[col] >= low_score) & (df[col] <= up_score)]
                                        up_scores.append(up_score)
                                        low_scores.append(low_score)
                                    output = pd.DataFrame(
                                        {'Column': numeric_cols, 'Up Score': up_scores, 'Low Score': low_scores,
                                         'Above Count': above_count,
                                         'Below Count': below_count})
                                    output.to_csv(output_file, index=False)
                                    return df
                                ## LOF
                                def LOF_outlier(df, output_file1, output_file2):
                                    # outlier detection using LOF
                                    types = df[index_column].unique()
                                    print(types)
                                    mask = []
                                    features = df.select_dtypes(include='number').columns
                                    detector_list = [
                                        ("Local Outlier Factor 30", LocalOutlierFactor(n_neighbors=8))
                                    ]
                                    for name, algorithm in detector_list:
                                        errors = np.full(len(df), fill_value=np.nan)
                                        outliers = np.full(len(df), fill_value=np.nan)
                                        for type in types:
                                            x = df.loc[:, features].values
                                            F = x.sum(1)
                                            mask = np.zeros(x.shape[0])
                                            mask[np.isfinite(F)] = 1
                                            mask_type = mask * np.array(df[index_column] == type)
                                            Curr_df = df.loc[mask_type == 1, features]
                                            x = Curr_df.values
                                            if name == "Local Outlier Factor 30":
                                                algorithm.fit(x)
                                                errors[mask_type == 1] = algorithm.negative_outlier_factor_
                                                outliers[mask_type == 1] = algorithm.fit_predict(x)
                                        df[name] = errors
                                        df[f'{name}_outliers'] = outliers
                                        df.set_index(name, inplace=True,
                                                     append=True, drop=False)
                                        df.to_csv(os.path.join(output_file1))
                                        # exclude rows that were defined as outliers
                                        for col in df.columns:
                                            if col.endswith("_outliers"):
                                                df = df[df[col] != -1]
                                    df.to_csv(os.path.join(output_file2))
                                    print("outliers removed!")
                                    return df
                                #run the outlier detection functions
                                df_percentiles = compute_percentiles(out_dat_init,
                                                                     os.path.join(outlier_files_dir,
                                                                                  'percentiles_outliers.csv'))
                                perc_file_path = os.path.join(outlier_files_dir,
                                                               'perc_dat.csv')
                                df_percentiles.to_csv(perc_file_path, index=False)
                                df_iqr = compute_iqr(out_dat_init, os.path.join(outlier_files_dir, 'iqr_outliers.csv'))
                                iqr_file_path = os.path.join(outlier_files_dir,
                                                               'iqr_dat.csv')
                                df_iqr.to_csv(iqr_file_path, index=False)
                                df_scores = compute_scores(out_dat_init,
                                                           os.path.join(outlier_files_dir, 'scores_outliers.csv'))
                                zscore_file_path = os.path.join(outlier_files_dir,
                                                               'zscore_dat.csv')
                                df_scores.to_csv(zscore_file_path, index=False)
                                df_lof = LOF_outlier(out_dat_init,
                                                     os.path.join(outlier_files_dir, 'processed_data_outliers_1.csv'),
                                                     os.path.join(outlier_files_dir, 'processed_data_outliers_2.csv'))
                                print("Please review outlier detection methods in outlier results folder and choose one method to continue with:")
                                method_choice = input("Enter 1 for percentiles, 2 for IQR, 3 for scores, 4 for lof: ")
                                if method_choice == '1':
                                    out_dat = df_percentiles
                                elif method_choice == '2':
                                    out_dat = df_iqr
                                elif method_choice == '3':
                                    out_dat = df_scores
                                elif method_choice == '4':
                                    out_dat = df_lof
                                    out_dat = out_dat.drop(columns=["Local Outlier Factor 30", "Local Outlier Factor 30_outliers", "Local Outlier Factor 30"])
                                else:
                                    print("Invalid choice. Defaulting to original DataFrame.")
                                    out_dat = out_dat_init
                                    out_dat = out_dat.drop(
                                        columns=["Local Outlier Factor 30", "Local Outlier Factor 30_outliers",
                                                 "Local Outlier Factor 30"])
                            else:
                                print("No outlier detection ran.")
                                out_dat = plate_df_2.copy()
                            #final filtered DataFrame
                            final_file_path = os.path.join(results_folder, 'final_filtered_data(post outlier detection).csv')
                            out_dat.to_csv(final_file_path, index=False)
                            print("Final filtered DataFrame saved to", final_file_path)
                        except Exception as e:
                            logging.error(f"Failed to perform outlier detection{e}")
                        #prompting index column for final proc- "PC" is must!
                        for idx, column in enumerate(out_dat.columns):
                            print(f"{idx}: {column}")
                        #columns to set as indexes
                        try:
                            print(Fore.RED + "==== Please make sure you choose all needed index columns - PC is must! ====")
                            columns_to_set_index = input("Enter ALL needed columns to set as indexes -PC IS MUST! (comma-separated): ").split(',')
                            columns_to_set_index = [col.strip() for col in columns_to_set_index if col.strip()]
                            print(f"Columns to set as index: {columns_to_set_index}")
                            if columns_to_set_index:
                                missing_index_cols = [col for col in columns_to_set_index if col not in out_dat.columns]
                                if missing_index_cols:
                                    logging.warning(f"Index columns not found: {missing_index_cols}")
                                out_dat = out_dat.set_index([col for col in columns_to_set_index if col in out_dat.columns], drop=False)
                                print(f"Set columns as index: {columns_to_set_index}")
                        except Exception as e:
                            logging.error(f"Error setting index columns: {e}")
                        out_dat = out_dat.select_dtypes(include=np.number)
                        #target variable selection for feature qc
                        try:
                            #user to select the correct index title as the target variable
                            if out_dat.index.names:
                                print("Select the correct index title to be the target variable for feature quality control:")
                                for idx, index_title in enumerate(out_dat.index.names):
                                    print(f"{idx + 1}: {index_title}")
                                while True:
                                    try:
                                        selected_index = int(
                                            input("Enter the number corresponding to the correct index title: ")) - 1
                                        if selected_index < 0 or selected_index >= len(out_dat.index.names):
                                            raise ValueError("Selection out of range.")
                                        break
                                    except ValueError as e:
                                        logging.error(f"Invalid selection. Please try again. Error: {e}")
                                #selected index title as the target variable
                                target_variable = out_dat.index.names[selected_index]
                                print(f"Target variable selected: {target_variable}")
                            else:
                                raise ValueError("No index names found in the data.")
                            #factor the target variable
                            target_variable_labels, _ = pd.factorize(out_dat.index.get_level_values(target_variable))
                            print(f"Target variable labels: {target_variable_labels}")
                            #directory for feature qc results
                            feature_qc_dir = os.path.join(results_folder, "feature_qc")
                            os.makedirs(feature_qc_dir, exist_ok=True)
                            #statistical summary for features
                            statistical_summary = out_dat.describe()
                            statistical_summary.to_csv(os.path.join(feature_qc_dir, "statistical_summary.csv"))
                            print("Statistical summary saved to 'statistical_summary.csv'.")
                            #feature stability scores -CV - for normalizied features
                            normalized_data = (out_dat - out_dat.min()) / (out_dat.max() - out_dat.min())
                            if normalized_data.isna().any().any():
                                logging.warning("Normalization produced NaN values. Check your data for zero variance features.")
                            stability_scores = normalized_data.var() / normalized_data.mean()
                            print("Stability scores calculated.")
                            #correlation
                            correlation_matrix = out_dat.corr(method='spearman')
                            correlation_matrix.to_csv(os.path.join(feature_qc_dir, "correlation_matrix.csv"))
                            print("Correlation matrix saved to 'correlation_matrix.csv'.")
                            #plot the correlation matrix
                            lan = plt.rcParams['font.family'] = ['Arial']
                            plt.figure(figsize=(32, 18))
                            font = {'family': 'Arial',
                                    'size': 10}
                            plt.rc('font', **font)
                            plt.yticks(family='Arial', rotation="horizontal", size=10)
                            plt.xticks(family='Arial', rotation="vertical", size=8)
                            plt.margins(0.1)
                            plt.subplots_adjust(bottom=0.3)
                            sns.heatmap(correlation_matrix, annot=True, cmap='coolwarm', fmt=".2f")
                            plt.title("Correlation Matrix")
                            plt.tight_layout()
                            correlation_plot_path = os.path.join(feature_qc_dir, "correlation_matrix_plot.png")
                            plt.savefig(correlation_plot_path)
                            plt.close()
                            print(f"Correlation matrix plot saved to: {correlation_plot_path}")
                            #CV plot
                            if stability_scores is not None and not stability_scores.isna().all():
                                plt.figure(figsize=(32, 16))
                                plt.bar(stability_scores.index, stability_scores.values, color='skyblue')
                                plt.title("Feature Stability Scores (Coefficient of Variation)")
                                plt.xlabel("Feature")
                                plt.ylabel("Stability Score")
                                plt.xticks(rotation=45, ha='right')
                                plt.margins(0.1)
                                plt.tight_layout()
                                cv_plot_path = os.path.join(feature_qc_dir, "stability_scores_plot.png")
                                plt.savefig(cv_plot_path)
                                plt.close()
                                print(f"CV plot saved to: {cv_plot_path}")
                            else:
                                logging.warning("Stability scores not available or invalid. Skipping CV plot.")
                            #random forest score computation per features
                            rf_scores = {}
                            try:
                                for feature in out_dat.columns:
                                    model = RandomForestRegressor()
                                    feature_values = out_dat[feature].values.reshape(-1, 1)
                                    model.fit(feature_values, target_variable_labels)
                                    rf_scores[feature] = model.score(feature_values, target_variable_labels)
                                    logging.debug(f"Random Forest score for feature '{feature}': {rf_scores[feature]}")
                            except Exception as e:
                                logging.error(f"Error occurred during Random Forest model training: {e}")
                                rf_scores[feature] = np.nan
                            #RF scores to csv
                            rf_scores_df = pd.DataFrame(rf_scores, index=['RF_Score'])
                            rf_scores_df.to_csv(os.path.join(feature_qc_dir, "random_forest_scores.csv"))
                            print("Random Forest scores saved to 'random_forest_scores.csv'.")
                            #plot RF scores
                            plt.figure(figsize=(32, 16))
                            plt.bar(rf_scores.keys(), rf_scores.values(), color='lightgreen')
                            plt.title("Random Forest Scores for each Feature")
                            plt.xlabel("Feature")
                            plt.ylabel("Random Forest Score")
                            plt.xticks(rotation=45, ha='right')
                            plt.margins(0.1)
                            plt.tight_layout()
                            rf_plot_path = os.path.join(feature_qc_dir, "random_forest_scores_plot.png")
                            plt.savefig(rf_plot_path)
                            plt.close()
                            print(f"Random Forest scores plot saved to: {rf_plot_path}")
                            print(f"Feature QC finsihed and saved in{feature_qc_dir}")
                            #want to change feature names?
                            rename_columns = input("Do you want to rename any columns (recommended for long feature names) ? (yes/no): ").strip().lower()
                            if rename_columns == 'yes':
                                print("Current Columns:")
                                for idx, column in enumerate(out_dat.columns):
                                    print(f"{idx}: {column}")
                                renaming_input = input(
                                    "Enter the column name and new name separated by ':', each pair separated by commas-no space! (e.g., col1:new_col1,col2:new_col2): ").strip()
                                try:
                                    rename_dict = dict(pair.split(':') for pair in renaming_input.split(',') if ':' in pair)
                                    if rename_dict:
                                        out_dat.rename(columns=rename_dict, inplace=True)
                                        print("Names changed successfully.")
                                    else:
                                        logging.warning("No valid renaming pairs provided.")
                                except Exception as e:
                                    logging.error(f"Error renaming columns: {e}")
                            out_dat.to_csv(os.path.join(results_folder, "data_with_new_names.csv"))
                            print(
                                f"Data with new names saved to: {os.path.join(results_folder, 'data_with_new_names.csv')}")
                        except Exception as e:
                            logging.error(f"An error occurred during renaming and feature qc process: {e}")
                            return
                        print("ready_for_norm!")
                        #imputation and mormalization segment
                        try:
                            print("imputing data for clean normalization")
                            #initialize simpleimputer for missing values (NaN) - median!
                            imputer = SimpleImputer(strategy='median')
                            #impute missing values (NaN)
                            df_imputed = pd.DataFrame(imputer.fit_transform(out_dat), columns=out_dat.columns)
                            print("Imputation of missing values completed.")
                            #replace zero values with NaN to use the same imputer
                            df_imputed.replace(0, np.nan, inplace=True)
                            logging.debug("Replaced zero values with NaN for further imputation.")
                            #impute zero values (now NaNs) using the same strategy
                            df_final_imputed = pd.DataFrame(imputer.fit_transform(df_imputed), columns=out_dat.columns)
                            df_final_imputed.index = out_dat.index
                            print("Imputation of zero values completed.")
                            #save imputed data
                            final_imputed_path = os.path.join(results_folder, "df_final_imputed.csv")
                            df_final_imputed.to_csv(final_imputed_path)
                            print(f"Imputed data saved to {final_imputed_path}")
                            #normalization
                            print("Proceeding to normalization...")
                            selected_data = df_final_imputed.copy()
                            #directory for normalization results
                            normalization_dir = os.path.join(results_folder, "normalization")
                            os.makedirs(normalization_dir, exist_ok=True)
                            print(f"Normalization results will be saved to: {normalization_dir}")
                            #split selected_data based on PC column - error handling for missing 'PC' index
                            if 'PC' not in selected_data.index.names:
                                raise ValueError("'PC' index not found in selected data. Ensure 'PC' is part of the index.")
                            pc_groups = selected_data.groupby(level='PC', group_keys=False)
                            df_reverted = pc_groups.apply(lambda x: x)
                            data_3_path = os.path.join(results_folder, "df_proc_3.csv")
                            df_reverted.to_csv(data_3_path)
                            print(f"Final processed data saved to {data_3_path}")
                            #running over each group and performing normalization
                            normalization_methods = {
                                "min-max_normalization": lambda x: (x - x.min()) / (x.max() - x.min()),
                                "central_log_normalization": lambda x: np.log(x + np.sqrt(x ** 2 + 1)),
                                "z-score_normalization": lambda x: (x - x.mean()) / x.std()
                            }
                            for pc_value, pc_group in pc_groups:
                                pc_dir = os.path.join(normalization_dir, f"df_{pc_value}")
                                os.makedirs(pc_dir, exist_ok=True)
                                print(f"Normalization directory created at: {pc_dir}")
                                #original data
                                original_data_path = os.path.join(pc_dir, "original_data.csv")
                                pc_group.to_csv(original_data_path)
                                print(f"Original data saved to {original_data_path}")
                                #normalization methods and saving results
                                normalized_data_dict = {}
                                for method_name, method_func in normalization_methods.items():
                                    try:
                                        normalized_data = method_func(pc_group)
                                        normalized_data.to_csv(os.path.join(pc_dir, f"{method_name.lower().replace(' ', '_')}.csv"))
                                        normalized_data_dict[method_name] = normalized_data
                                        print(f"{method_name} applied and saved to {method_name.lower().replace(' ', '_')}.csv")
                                    except Exception as e:
                                        logging.error(f"Error applying {method_name} for PC value {pc_value}: {e}")
                                #box-cox transformation - data must be positive
                                try:
                                    transformed_data = pd.DataFrame()
                                    for column in pc_group.columns:
                                        if (pc_group[column] <= 0).any():
                                            raise ValueError(
                                                f"Box-Cox transformation requires positive data, found non-positive values in {column}.")
                                        transformed_column, _ = boxcox(pc_group[column])
                                        transformed_data[column] = transformed_column
                                    transformed_data.index = pc_group.index
                                    transformed_data.to_csv(os.path.join(pc_dir, "box_cox_normalized_data.csv"))
                                    print(f"Box-Cox normalized data saved to box_cox_normalized_data.csv")
                                except Exception as e:
                                    logging.error(f"Error applying Box-Cox transformation for PC value {pc_value}: {e}")
                                #log normalization
                                try:
                                    log_data = pd.DataFrame()
                                    for column in pc_group.columns:
                                        #handle negative or zero values for log normalization
                                        if (pc_group[column] <= 0).any():
                                            raise ValueError(
                                                f"Log normalization requires positive data, found non-positive values in {column}.")
                                        log_column = np.log(pc_group[column])
                                        log_data[column] = log_column
                                    log_data.index = pc_group.index
                                    log_data.to_csv(os.path.join(pc_dir, "log_normalized_data.csv"))
                                    print(f"Log normalized data saved to log_normalized_data.csv")
                                except Exception as e:
                                    logging.error(f"Error applying log normalization for PC value {pc_value}: {e}")
                                #QQ plots and histograms for each feature in each normalized data frame
                                normalized_dfs = [
                                    ("Original Data", pc_group),
                                    ("min-max_normalization", normalized_data_dict.get("min-max_normalization")),
                                    ("central_log_normalization", normalized_data_dict.get("central_log_normalization")),
                                    ("z-score_normalization", normalized_data_dict.get("z-score_normalization")),
                                    ("Box_Cox Normalized Data", transformed_data),
                                    ("Log Normalized Data", log_data)

                                ]
                                for df_name, df in normalized_dfs:
                                    if df is None:
                                        logging.warning(f"{df_name} is not available for PC value {pc_value}. Skipping plots.")
                                        continue
                                    df_dir = os.path.join(pc_dir, df_name.replace(' ', '_').lower())
                                    os.makedirs(df_dir, exist_ok=True)
                                    print(f"Plot directory created at: {df_dir}")
                                    for feature in df.columns:
                                        try:
                                            ## qq
                                            plt.figure()
                                            scipy.stats.probplot(df[feature], dist="norm", plot=plt)
                                            plt.title(f"Q-Q-{feature}")
                                            qq_plot_path_2 = os.path.join(df_dir, f"qq_2_plot_{feature}.png")
                                            plt.savefig(qq_plot_path_2)
                                            plt.close()
                                            plt.clf()
                                        except Exception as e:
                                            logging.error(f"Error plotting Q-Q for {feature} in {df_name}: {e}")
                                        # histogram
                                        try:
                                            plt.figure(figsize=(18, 12))
                                            sns.histplot(df[feature], bins=20, kde=True, color='skyblue')
                                            plt.title(f"Histogram for {feature} ({df_name})")
                                            plt.xlabel("Value")
                                            plt.ylabel("Frequency")
                                            plt.tight_layout()
                                            histogram_path = os.path.join(df_dir, f"histogram_{feature}.png")
                                            plt.savefig(histogram_path)
                                            plt.close()
                                            print(f"Histogram saved to: {histogram_path}")
                                        except Exception as e:
                                            logging.error(f"Error plotting histogram for {feature} in {df_name}: {e}")
                            #user to select the normalized data frame to continue with
                            try:
                                normalized_options = ["min-max_normalization", "central_log_normalization",
                                                      "z-score_normalization",
                                                      "Box_Cox Normalized Data", "Log Normalized Data"]
                                print("Available normalization methods:")
                                for idx, option in enumerate(normalized_options):
                                    print(f"{idx + 1}: {option}")
                                selected_normalization = int(
                                    input("Check normalized data and enter the number corresponding to the desired normalized data frame: ")) - 1
                                if selected_normalization < 0 or selected_normalization >= len(normalized_options):
                                    raise ValueError("Selection out of range.")
                                selected_normalization_name = normalized_options[selected_normalization]
                                print(f"Selected normalization method: {selected_normalization_name}")
                                #combine all PC values for the selected normalization into one final normalized data frame
                                final_normalized_data = pd.DataFrame()
                                for pc_value, pc_group in pc_groups:
                                    pc_dir = os.path.join(normalization_dir, f"df_{pc_value}")
                                    file_path = os.path.join(pc_dir, f"{selected_normalization_name.lower().replace(' ', '_')}.csv")
                                    if not os.path.exists(file_path):
                                        logging.warning(f"Normalized file not found: {file_path}. Skipping.")
                                        continue
                                    normalized_df = pd.read_csv(file_path, index_col=0)
                                    final_normalized_data = pd.concat([final_normalized_data, normalized_df])
                                    logging.debug(f"Data from {file_path} added to final normalized data.")
                                #save final normalized data
                                final_normalized_data_path = os.path.join(results_folder, "final_normalized_data.csv")
                                final_normalized_data.to_csv(final_normalized_data_path)
                                logging.info(f"Final normalized data saved to: {final_normalized_data_path}")
                            except Exception as e:
                                logging.error(f"An error occurred during the imputation or normalization process: {e}")
                                return
                        except Exception as e:
                            logging.error(f"An unexpected error occurred: {e}")
                            return
                        finally:
                            logging.info('All processing steps completed.')
                        #dashboard HTML generation
                        #relative path for links
                        print("creating dashboard for results summary")
                        dashboard_page = os.path.join(main_results_dir, f'{item_name}_dashboard.html')
                        dashboard_content = f"""
                        <html>
                        <head><title>Dashboard for Folder {item_name}</title></head>
                        <body>
                            <h1>Results for Folder: {item_name}</h1>
                            <p>Results stored in: {results_folder}</p>
                        """
                        #add a section for each  file in the folder
                        dashboard_content += "<h2>Processed Plates</h2>"
                        for file_name in os.listdir(item_path):
                            if file_name.endswith(".csv"):
                                results_folder = os.path.join(item_path, f"results_{file_name}_res")
                                barplot_dir = os.path.join(results_folder, "box_plots")
                                summary_dir = os.path.join(results_folder, "general_results")
                                barplot_relative_path = os.path.relpath(barplot_dir, main_results_dir)
                                summary_relative_path = os.path.relpath(summary_dir, main_results_dir)
                                dashboard_content += f"<h3>Results for Plate: {file_name}</h3>"
                                dashboard_content += "<ul>"
                                #bar plots for important features
                                for feature in important_features:
                                    bar_plot_html = os.path.join(barplot_relative_path, f'{feature}_bar_plot.html')
                                    dashboard_content += f'<li><a href="{bar_plot_html}">Bar Plot for {feature}</a></li>'
                                #PCA plot
                                pca_plot_html_id = os.path.join(summary_relative_path, f'pca_plot_with_id_{file_name[:-4]}.html')
                                dashboard_content += f'<li><a href="{pca_plot_html_id}">PCA Plot for {file_name}</a></li>'
                                #PCA plot2
                                pca_plot_html_reg = os.path.join(summary_relative_path, f'pca_plot_wo_id_{file_name[:-4]}.html')
                                dashboard_content += f'<li><a href="{pca_plot_html_reg}">PCA Plot for {file_name}</a></li>'
                                #number of cells plot for the current plate
                                count_plot = os.path.join(summary_relative_path,
                                                                f'number_of_cells_bar_plot.pdf')
                                dashboard_content += f'<li><a href="{count_plot}">number of cells Plot for {file_name}</a></li>'
                                #sum file
                                sum_file = os.path.join(summary_relative_path,
                                                                f'init_statistical_summary.xlsx')
                                dashboard_content += f'<li><a href="{sum_file}">stat_summary for {file_name}</a></li>'
                                #heatmap
                                heat_plot_pdf_reg = os.path.join(summary_relative_path,
                                                                 f'heatmap.pdf')
                                dashboard_content += f'<li><a href="{heat_plot_pdf_reg}">heatmap for {file_name}</a></li>'

                                dashboard_content += "</ul>"

                        #closing the HTML content
                        dashboard_content += "</body></html>"

                        #updated dashboard content to the HTML file
                        with open(dashboard_page, 'w') as f:
                            f.write(dashboard_content)
                        print(f"Dashboard updated with  plots and saved for {item_name} at {dashboard_page}")
                    print(f"{file_name} proccessed succesfully")
    except Exception as e:
        logging.error(f"failed proccessing {plate_df}")
#####################
def gather_user_lists():
    print("Choose appropriate columns for your analysis!")
    #options based on loaded data
    print("usual unwanted columns: Number of Analyzed Fields, Height, Plane, Time, Timepoint, Cell Count, Target CO2, CO2, Target Temperature, Temperature")
    print("enter unwanted columns:")
    list1 = input("enter unwanted columns:(comma-separated): ").split(',')
    print("usual important columns: spots_chanel_3_final - Corrected Spot Intensity - Mean per Well, spots_chanel_3_final - Region Intensity - Mean per Well, spots_chanel_3_final - Spot to Region Intensity - Mean per Well, ir_chanel_3_at - Intensity ir_chanel_3_at chanel_3 Mean - Mean per Well, all_cells - chanel_3 Gabor Max 2 px w2 - Mean per Well, all_cells - chanel_3 SER Ridge 1 px - Mean per Well, all_cells - chanel_3_intensity Mean - Mean per Well, all_cells - chanel_2 Area, all_cells - chanel_1 Area, chan3_intensitymean_div_chan2area, spotschan3_intensitymean_div_chan2area")
    print("enter important_features:")
    list2 = input("enter important_features:(comma-separated): ").split(',')
    print("must index columns!!!: PC, cell_type, cell_id, compound, concentration")
    print("enter index_columns:")
    list3 = input("enter index_columns: (comma-separated): ").split(',')
    #list of strings
    list1 = [item.strip() for item in list1]
    list2 = [item.strip() for item in list2]
    list3 = [item.strip() for item in list3]
    return list1, list2, list3
#RUN MAIN
def main(main_folder, main_folder_path, final_folder_path):
    print(Fore.GREEN + "Starting MA-TMRE analysis for combined file of all plates in folder in your main directory:)")
    print(Fore.GREEN + "STEP 1 - processing text files...")
    process_folders(main_folder)
    prompt_for_review()
    print(Fore.GREEN + "STEP 2 - combining csv files...")
    process_data_combine(main_folder_path)
    prompt_for_review()
    print(Fore.GREEN + "STEP 3 - choosing column names for analysis...")
    print(Fore.RED + "==== Please make sure you choose correct column names! ====")
    list1, list2, list3 = gather_user_lists()
    process_plate_data_on_folders(final_folder_path, list1, list2, list3)
    print("Completed combined plates MA-TMRE analysis, good luck with statistics:)")
    print(Fore.RED + Style.BRIGHT + "==== IMPORTANT MESSAGE ====" +
          "PLEASE REMOVE YOUR FILES FROM THE MAIN DIRECTORY ====")

if __name__ == "__main__":
    main_folder = input("Enter folder name(containing a folder with .txt files): ")
    main_folder_path = input("Enter: the main folder name/results ")
    final_folder_path = input("Enter: the main folder name/MA-TMRE_result_combined ")
    main(main_folder, main_folder_path, final_folder_path)